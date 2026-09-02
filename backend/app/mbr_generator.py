"""
Geração automática do MBR (Monthly Business Review) em PPTX (20/08/2026,
reescrita completa) - gera o relatório com DADOS REAIS lidos diretamente das
mesmas funções de negócio que alimentam as telas do Atlas (chamadas em
processo, como funções Python normais - não via HTTP), sem abrir navegador
nenhum. Substitui a versão anterior (prints de tela via Playwright/Chromium)
por decisão explícita do usuário, que pediu um relatório no estilo do MBR do
gestor dele (resumo executivo, leitura por frente, scorecard, decisões) com a
identidade visual da Mágio Chocolates, contando a história de cada módulo com
números e não com telas printadas.

Por que chamar as funções dos routers direto (ex.: fechamento_router.
dashboard_kpis(...)) em vez de bater no HTTP: cada uma delas tem
"usuario: models.Usuario = Depends(obter_usuario_atual)" e "db: Session =
Depends(get_db)" só como VALOR PADRÃO - "Depends(...)" é apenas um marcador
que o FastAPI interpreta na hora de uma requisição de verdade. Chamando a
função Python diretamente e passando um "usuario"/"db" reais no lugar, a
injeção de dependência do FastAPI nunca entra em ação - funciona como
qualquer chamada de função comum, sem round-trip de rede e sem os problemas
de Chromium/memória da versão anterior (ver histórico em ATUALIZANDO.md).

IMPORTANTE - limitações conhecidas e assumidas nesta versão:
  - Ícones: o sistema visual da Mágio pede ícones outline (Lucide/Heroicons),
    mas renderizá-los aqui exigiria ferramental Node.js (react-icons + sharp)
    que não faz parte do backend Python do Atlas. Substituído por hierarquia
    tipográfica + badges de cor (mesma linguagem de status, sem o ícone
    literal).
  - Fonte: usa Arial em todo o documento (título e corpo) de propósito, não
    "Poppins"/"TT Chocolates" como pede a identidade visual - Arial é uma das
    poucas fontes com garantia de existir em qualquer instalação do
    PowerPoint do usuário E de ter a MESMA largura de texto tanto aqui
    (pré-visualização via LibreOffice) quanto lá (não quebra layout depois de
    entregue). Quem tiver "Poppins"/"TT Chocolates" licenciada pode trocar a
    fonte de todo o arquivo em segundos pelo PowerPoint (Página Inicial >
    Substituir > Substituir Fontes).
  - Limiares de status ("Em avanço"/"Atenção"/"Crítico"): ver _LIMIARES logo
    abaixo - são um ponto de partida razoável, não uma meta oficial definida
    pela operação. Fácil de ajustar depois se não bater com o apetite de
    risco real do time.
"""
import calendar
import html
import re
import unicodedata
from collections import defaultdict
from datetime import date, datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
from sqlalchemy.orm import Session

from . import models
from . import shelf_life as shelf_life_mod
from . import fefo as fefo_mod
from . import dashboards_externos_extrator as dash_ext
from .routers import (
    fechamento_router, baixas_operacionais_router, movimentados_router,
    cadastros_router,
)

# ---------------------------------------------------------------------------
# Identidade visual Mágio Chocolates
# ---------------------------------------------------------------------------
VERDE_AMAZONIA = RGBColor(0x4E, 0x7F, 0x84)      # cor primária
AZUL_INSTITUCIONAL = RGBColor(0x34, 0x43, 0x6C)  # cor secundária
AZUL_CLARO = RGBColor(0xB1, 0xBF, 0xE2)          # cor de apoio
OFF_WHITE = RGBColor(0xE4, 0xE0, 0xD5)           # cor neutra
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TEXTO = RGBColor(0x4A, 0x4A, 0x4A)
CINZA_CLARO = RGBColor(0xE2, 0xE4, 0xE8)
CINZA_LINHA_PAR = RGBColor(0xF4, 0xF5, 0xF7)

COR_SUCESSO = RGBColor(0x2E, 0x7D, 0x32)
COR_ATENCAO = RGBColor(0xF9, 0xA8, 0x25)
COR_ERRO = RGBColor(0xC6, 0x28, 0x28)
COR_INFO = RGBColor(0x15, 0x65, 0xC0)
COR_SEM_DADO = RGBColor(0x9E, 0x9E, 0x9E)

# Farol de Shelf-Life (22/08/2026) - 4 faixas de severidade, mesma paleta já
# validada com o usuário na simulação HTML (faixa "vencido" é a mais severa,
# mais escura que "urgente" de propósito - não é gradiente, é OUTRA cor pra
# não confundir com 0-30 dias ainda não vencido).
COR_FAROL_URGENTE = RGBColor(0xF1, 0x70, 0x4B)   # 0-30 dias
COR_FAROL_PERIGO = RGBColor(0xF2, 0xC1, 0x4E)    # 31-60 dias
COR_FAROL_ATENCAO = RGBColor(0x2C, 0xC5, 0xA8)   # 61-90 dias
COR_FAROL_VENCIDO = RGBColor(0x7B, 0x3B, 0x2E)   # já vencido

# Ver nota no docstring do módulo sobre a troca TT Chocolates/Poppins -> Arial.
FONTE_TITULO = "Arial"
FONTE_TEXTO = "Arial"

LARGURA_IN = 13.333
ALTURA_IN = 7.5
MARGEM_IN = 0.5

# (bom, atenção) - acima de "bom" é "Em avanço", entre os dois é "Atenção",
# abaixo do segundo é "Crítico". Documentado no docstring do módulo como
# ponto de partida ajustável, não meta oficial da operação.
_LIMIARES = {
    "acuracia": (95.0, 85.0),        # maior é melhor
    "fefo_quebra_pct": (5.0, 20.0),  # menor é melhor
}


def _status_maior_melhor(valor, bom, atencao):
    if valor is None:
        return ("Sem dado", COR_SEM_DADO)
    if valor >= bom:
        return ("Em avanço", COR_SUCESSO)
    if valor >= atencao:
        return ("Atenção", COR_ATENCAO)
    return ("Crítico", COR_ERRO)


def _status_menor_melhor(valor, bom, atencao):
    if valor is None:
        return ("Sem dado", COR_SEM_DADO)
    if valor <= bom:
        return ("Em avanço", COR_SUCESSO)
    if valor <= atencao:
        return ("Atenção", COR_ATENCAO)
    return ("Crítico", COR_ERRO)


def _aplica_avanco_operacional(pct, status_label, status_cor):
    """Regra de negócio pedida pelo usuário, 22/08/2026 ("Considere todo
    resultado acima de 95% como um avanço operacional nos Scorecard"):
    sobreposição final aplicada DEPOIS de qualquer outro critério que a
    linha do Scorecard já usasse pra chegar num status (gap vs. item a
    item, tendência, etc.) — se o resultado percentual da frente é >= 95%,
    o status vira "Em avanço" independente do que a lógica específica
    daquela linha tivesse concluído. Corrige, por exemplo, a linha de
    Acurácia Ponderada no Scorecard do Mês: antes desta regra, um IAP de
    98,8% podia aparecer como "Atenção" só porque a distorção item-a-item
    vs. IAP estava alta — o IAP em si já é um resultado forte e deve
    contar como avanço, a distorção é uma informação complementar, não um
    motivo pra rebaixar o status.

    Não faz nada (devolve o status recebido sem alteração) se `pct` for
    None ou menor que 95 — nesses casos a lógica específica da linha
    continua valendo como estava."""
    if pct is not None and pct >= 95.0:
        return "Em avanço", COR_SUCESSO
    return status_label, status_cor


def _status_evolucao(delta, menor_e_melhor=False, limiar=0.5):
    """Classifica uma VARIAÇÃO (mês atual vs. anterior), não um nível absoluto
    — usado pelos Scorecards de evolução/involução (Inventário por Almoxarifado,
    Mapeamento de Riscos, 20/08/2026), que precisam responder "melhorou ou
    piorou desde o mês passado", pergunta diferente de "está numa faixa boa
    hoje" (essa segunda already respondida por _status_maior_melhor/
    _status_menor_melhor). `menor_e_melhor=True` pra métricas onde cair é bom
    (taxa de quebra, taxa de furo, gasto) — inverte o sinal antes de comparar
    contra o limiar de relevância (pontos percentuais ou %, conforme a métrica;
    variações menores que o limiar são tratadas como ruído, não avanço/involução
    real)."""
    if delta is None:
        return ("Sem histórico", COR_SEM_DADO)
    direcao = -delta if menor_e_melhor else delta
    if direcao >= limiar:
        return ("Evolução", COR_SUCESSO)
    if direcao <= -limiar:
        return ("Involução", COR_ERRO)
    return ("Estável", COR_INFO)


def _mes_anterior(mes: str) -> str:
    ano, m = (int(parte) for parte in mes.split("-"))
    if m == 1:
        return f"{ano - 1}-12"
    return f"{ano}-{m - 1:02d}"


def _ultimo_dia_mes(mes: str) -> str:
    """"YYYY-MM" -> "YYYY-MM-DD" do último dia daquele mês (data-limite pra
    indicadores que trabalham em janela de dias contados a partir de uma
    data de referência)."""
    ano, m = (int(parte) for parte in mes.split("-"))
    ultimo_dia = calendar.monthrange(ano, m)[1]
    return f"{ano}-{m:02d}-{ultimo_dia:02d}"


def _truncar_serie_mensal(serie: list, mes_limite: str) -> list:
    """Várias funções de "evolução mensal" usadas pelo MBR (fechamento_router.
    dashboard_evolucao_mensal e dashboard_evolucao_ponderada_mensal,
    movimentados_router.dashboard_evolucao_mensal e
    dashboard_transferencias_evolucao_mensal, baixas_operacionais_router.
    dashboard_passivos_evolucao_mensal) não aceitam nenhum parâmetro de "mês
    de corte" - são pensadas pra uso em dashboards AO VIVO, onde mostrar até
    o mês mais recente que já tem dado no banco é o comportamento certo, e
    por isso sempre devolvem a série completa (histórico inteiro), nunca
    filtrada.

    O MBR, ao contrário, é gerado para um mês de FECHAMENTO específico
    (ex.: "2026-07") e pode ser gerado bem depois desse mês (ex.: hoje,
    "2026-08-22", já existindo algum registro parcial de agosto no banco).
    Sem esse corte, os gráficos de evolução (Painel de Inventário, Acurácia
    Ponderada, Controle de Movimentados, Mapeamento de Passivos) e os
    deltas MoM do "Recorte do Período"/Resumo Executivo liam sempre o ÚLTIMO
    ponto da série completa - ou seja, o mês mais recente com QUALQUER
    registro no banco, não o mês selecionado no filtro do relatório. Bug
    reportado pelo usuário em 22/08/2026: "mesmo selecionando o fechamento
    de julho, está trazendo dados do mês atual (agosto) em quase todos os
    módulos".

    Corta aqui, uma única vez por série, logo após a coleta em
    _coletar_dados_mbr - qualquer ponto com "mes" > mes_limite é descartado
    antes de chegar em _recorte_periodo/_analise_geral ou em qualquer slide
    (todos eles fazem d["evolucao_..."][-1] ou [-6:], então bastava a série
    já vir cortada na origem)."""
    return [ponto for ponto in serie if ponto.get("mes", "") <= mes_limite]


def _tendencia_linear(valores: list) -> dict:
    """Regressão linear simples (mínimos quadrados) sobre uma série mensal
    igualmente espaçada — devolve a inclinação (variação média por mês) e
    um selo de texto pronto pra badge de tendência (pedido do usuário,
    22/08/2026, mockups aprovados de Painel de Inventário/Acurácia
    Ponderada: "não só qual o nível hoje, pra onde isso está indo").

    Selo de texto pronto pra badge (direção + p.p./mês) — a linha
    TRAÇADA de verdade sobre o gráfico é responsabilidade de
    _pontos_tendencia_linear (mesma regressão, devolvendo a série de
    pontos da reta em vez do resumo em texto); as duas funções usam o
    MESMO cálculo de propósito, pra o número do selo e a inclinação da
    linha desenhada nunca discordarem entre si. Limiar de 0.3 p.p./mês pra
    não chamar de "tendência" uma variação que é só ruído mês a mês.

    Histórico (22/08/2026): esta função já foi a ÚNICA representação da
    tendência (um selo de texto no lugar da linha, com a justificativa de
    que o python-pptx não expõe gráfico combinado na API de alto nível).
    O usuário rejeitou essa simplificação como não-negociável ("o gráfico
    não conta a história") — _pontos_tendencia_linear + os slides que a
    chamam via _grafico_categoria_com_tendencia resolvem isso editando o
    XML do gráfico diretamente (ver _adicionar_linha_combo_eixo_unico). O
    selo de texto continua existindo — o mockup aprovado mantém os dois,
    selo no cabeçalho do gráfico E linha desenhada sobre as barras."""
    pontos = [v for v in valores if v is not None]
    n = len(pontos)
    if n < 2:
        return {"inclinacao": None, "rotulo": None, "cor": COR_SEM_DADO}
    media_x = (n - 1) / 2
    media_y = sum(pontos) / n
    numerador = sum((i - media_x) * (v - media_y) for i, v in enumerate(pontos))
    denominador = sum((i - media_x) ** 2 for i in range(n))
    inclinacao = round(numerador / denominador, 2) if denominador else 0.0
    if inclinacao > 0.3:
        rotulo = f"▲ Tendência de melhora: +{_fmt_num(inclinacao, 1)} p.p./mês"
        cor = COR_SUCESSO
    elif inclinacao < -0.3:
        rotulo = f"▼ Tendência de piora: -{_fmt_num(abs(inclinacao), 1)} p.p./mês"
        cor = COR_ERRO
    else:
        rotulo = "▬ Tendência estável no período"
        cor = COR_SEM_DADO
    return {"inclinacao": inclinacao, "rotulo": rotulo, "cor": cor}


def _pontos_tendencia_linear(valores: list):
    """Mesma regressão linear simples de _tendencia_linear (mínimos
    quadrados sobre os pontos não-nulos, ignorando posições com None ao
    compactar o eixo x — ver docstring dela), mas devolve a série de
    pontos DA RETA, um valor por posição de `valores` (None nas mesmas
    posições em que `valores` já era None, pra abrir o mesmo "buraco" na
    linha que já existe na barra), pronta pra desenhar como 2ª série de
    um gráfico combo via _adicionar_linha_combo_eixo_unico. Devolve None
    (em vez de lista) se não houver pontos suficientes (<2) pra regressão,
    mesmo critério de _tendencia_linear."""
    indices = [i for i, v in enumerate(valores) if v is not None]
    pontos = [valores[i] for i in indices]
    n = len(pontos)
    if n < 2:
        return None
    media_x = (n - 1) / 2
    media_y = sum(pontos) / n
    numerador = sum((k - media_x) * (v - media_y) for k, v in enumerate(pontos))
    denominador = sum((k - media_x) ** 2 for k in range(n))
    inclinacao = numerador / denominador if denominador else 0.0
    intercepto = media_y - inclinacao * media_x
    resultado = [None] * len(valores)
    for k, i in enumerate(indices):
        resultado[i] = intercepto + inclinacao * k
    return resultado


def _prejuizo_por_almoxarifado(itens_concentracao: list) -> dict:
    """A partir dos itens de dashboard_concentracao_valor (cada item já tem
    almoxarifado/valor/divergencia_qtd), soma o resultado financeiro NETO
    por almoxarifado — sobra (divergencia_qtd > 0) soma positivo, falta
    (divergencia_qtd < 0) soma negativo, mesma convenção de sinal usada em
    kpis_inventario['resultado_liquido']. Usado nas tabelas "Resultado por
    Almoxarifado" dos slides de Painel de Inventário/IAP/IAQ (22/08/2026,
    mockups aprovados v4/v5).

    Atenção, limitação conhecida: só cobre itens que entraram na lista de
    concentração de valor (até 50 maiores por valor — ver docstring de
    dashboard_concentracao_valor). Em recortes com muitos itens divergentes
    de valor baixo, um almoxarifado com só itens fora desse top 50 aparece
    aqui com prejuízo R$ 0, mesmo tendo divergência real — documentado no
    rodapé dos slides que usam este cálculo, não escondido."""
    por_almox = defaultdict(float)
    for item in itens_concentracao:
        qtd = item.get("divergencia_qtd") or 0
        sinal = 1 if qtd > 0 else -1
        por_almox[item.get("almoxarifado") or "—"] += sinal * (item.get("valor") or 0)
    return dict(por_almox)


# ---------------------------------------------------------------------------
# Formatação (padrão numérico brasileiro: milhar "." decimal ",")
# ---------------------------------------------------------------------------
def _fmt_num(valor, casas=0):
    if valor is None:
        return "—"
    texto = f"{valor:,.{casas}f}"
    return texto.replace(",", "§").replace(".", ",").replace("§", ".")


def _fmt_pct(valor, casas=1):
    if valor is None:
        return "—"
    return f"{_fmt_num(valor, casas)}%"


def _fmt_moeda(valor, casas=0):
    if valor is None:
        return "—"
    sinal = "-" if valor < 0 else ""
    return f"{sinal}R$ {_fmt_num(abs(valor), casas)}"


_NOMES_MESES = [
    "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
    "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro",
]


def _nome_mes(mes: str) -> str:
    ano, m = mes.split("-")
    return f"{_NOMES_MESES[int(m) - 1]} de {ano}"


# ---------------------------------------------------------------------------
# Blocos visuais reutilizáveis
# ---------------------------------------------------------------------------
def _nova_apresentacao() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(LARGURA_IN)
    prs.slide_height = Inches(ALTURA_IN)
    return prs


def _slide_em_branco(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fundo(slide, cor: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = cor


def _retangulo(slide, x, y, w, h, cor_fill=None, cor_borda=None, arredondado=True, raio=0.10):
    tipo = MSO_SHAPE.ROUNDED_RECTANGLE if arredondado else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(tipo, Inches(x), Inches(y), Inches(w), Inches(h))
    if arredondado:
        try:
            shape.adjustments[0] = raio
        except (IndexError, ValueError):
            pass
    if cor_fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor_fill
    else:
        shape.fill.background()
    if cor_borda is not None:
        shape.line.color.rgb = cor_borda
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _texto(slide, x, y, w, h, texto, tamanho=16, negrito=False, cor=CINZA_TEXTO,
           alinhamento=PP_ALIGN.LEFT, fonte=None, espacamento=1.0, ancora_meio=False, italico=False):
    caixa = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if ancora_meio:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    linhas = texto if isinstance(texto, (list, tuple)) else [texto]
    for i, linha in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = linha
        p.font.size = Pt(tamanho)
        p.font.bold = negrito
        p.font.italic = italico
        p.font.color.rgb = cor
        p.font.name = fonte or FONTE_TEXTO
        p.alignment = alinhamento
        p.line_spacing = espacamento
    return caixa


def _cabecalho(slide, titulo, mes_label, pagina, subtitulo=None):
    _texto(slide, MARGEM_IN, 0.32, 9.5, 0.22, "RELATÓRIO EXECUTIVO · MÁGIO CHOCOLATES",
           tamanho=10, negrito=True, cor=VERDE_AMAZONIA)
    _texto(slide, MARGEM_IN, 0.58, 9.7, 0.55, titulo, tamanho=27, negrito=True,
           cor=AZUL_INSTITUCIONAL, fonte=FONTE_TITULO)
    if subtitulo:
        _texto(slide, MARGEM_IN, 1.10, 9.7, 0.35, subtitulo, tamanho=13, cor=CINZA_TEXTO)
    _texto(slide, LARGURA_IN - 2.9, 0.32, 2.4, 0.3, mes_label.upper(), tamanho=11,
           negrito=True, cor=VERDE_AMAZONIA, alinhamento=PP_ALIGN.RIGHT)
    _texto(slide, LARGURA_IN - 2.9, ALTURA_IN - 0.42, 2.4, 0.3, f"{pagina:02d}", tamanho=10,
           cor=CINZA_TEXTO, alinhamento=PP_ALIGN.RIGHT)


def _cartao_kpi(slide, x, y, w, h, valor_texto, rotulo, cor_valor=AZUL_INSTITUCIONAL, contexto=None,
                 cor_contexto=None, deslocamento_rotulo=None, tamanho_valor_base=27):
    _retangulo(slide, x, y, w, h, cor_fill=BRANCO, cor_borda=CINZA_CLARO, raio=0.14)
    pad = 0.16
    largura_texto = w - 2 * pad
    # Os valores numéricos "de sempre" (R$, %, contagens) sempre couberam numa
    # linha só, mas os dashboards externos passaram a alimentar este mesmo
    # cartão com textos livres (ex.: "grupo_maior_impacto" = "Produto em
    # Processo") que quebram em 2 linhas e invadem o rótulo, cuja posição é
    # fixa - reduz a fonte do valor pra caber numa linha e, no limite, corta
    # com reticências (20/08/2026).
    #
    # 21/08/2026, pedido do usuário ("diminua um pouco os cards... o
    # indicador tem que chamar mais atenção"): cartão mais baixo (ver
    # chamadas de _linha_kpis, altura 1.25/1.30 -> 1.05/1.10) e o número
    # (o indicador) MAIOR e mais dominante em relação à moldura/rótulo -
    # fonte-base do valor sobe de 28 pra 32pt, e o rótulo encolhe um pouco
    # (10.5 -> 10pt) pra reforçar a hierarquia número > rótulo.
    #
    # 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o
    # modelo final" - cartões ainda mais baixos no modelo manual, ex.: IAP
    # 0,837in, Painel de Inventário 0,739in, medidos diretamente no arquivo
    # dela) e BUG que ela reportou nesse mesmo arquivo (2 cartões do Resumo
    # Executivo com o rótulo em 2 linhas invadindo o texto de contexto
    # abaixo, ambos com posição fixa - ver print do slide 3): cartão mais
    # baixo de novo (ver _linha_kpis) e rótulo agora GARANTIDO numa linha só
    # - reduz a fonte do rótulo (mesmo esquema de encolher-e-truncar já usado
    # no valor acima) até caber, e só trunca com reticências no piso da
    # fonte. Antes disso, a altura da caixa do rótulo era fixa (uma linha),
    # então um rótulo comprido que quebrasse em 2 linhas invadia visualmente
    # a caixa de contexto logo abaixo, cuja posição também é fixa - esse
    # bug não tinha como acontecer com o valor numérico (que já encolhe/
    # corta), só com o rótulo, que não tinha a mesma proteção.
    #
    # O valor também encolhe um pouco no tamanho-base (32 -> 27pt) pra caber
    # OK no cartão mais baixo: a caixa de texto do python-pptx não corta
    # visualmente o que não cabe nela (sem "overflow: hidden") - o que evita
    # a invasão de verdade é a altura de LINHA real do texto vs. o espaço até
    # o próximo elemento, não a altura declarada da caixa. Em 32pt/0,85in de
    # cartão o valor invadia visualmente o rótulo por baixo (visto na
    # renderização) - 27pt reabre a folga necessária.
    #
    # 21/08/2026, Fase 3 - cartões compactos sem contexto (ver
    # `deslocamento_rotulo`, Recuperação de Shelf): com o rótulo colado logo
    # abaixo do valor (sem a folga de ~0,25in que um cartão alto tem), 27pt
    # volta a invadir - medido no PRÓPRIO arquivo de referência da usuária
    # (MBR_Atlas_202607_15.pptx), o valor desses cartões específicos está em
    # 20pt, não 27pt. `tamanho_valor_base` deixa o teto de encolhimento
    # configurável por chamada em vez de mexer no valor-padrão dos outros 12
    # lugares que já usam 27pt sem problema.
    tamanho_valor = tamanho_valor_base
    texto_valor = valor_texto
    if isinstance(valor_texto, str) and valor_texto:
        while tamanho_valor > 14 and len(valor_texto) > max(1, int(largura_texto / ((tamanho_valor / 72.0) * 0.57))):
            tamanho_valor -= 1
        max_chars = max(1, int(largura_texto / ((tamanho_valor / 72.0) * 0.57)))
        if len(valor_texto) > max_chars:
            texto_valor = valor_texto[:max(1, max_chars - 1)].rstrip() + "…"
    _texto(slide, x + pad, y + 0.08, w - 2 * pad, 0.34, texto_valor, tamanho=tamanho_valor,
           negrito=True, cor=cor_valor, fonte=FONTE_TITULO)

    # Rótulo: mesmo esquema de encolher-e-truncar do valor acima, mas com um
    # fator de largura por caractere bem mais conservador (0.72, contra 0.57
    # do valor) - é MAIÚSCULO e em negrito (ver rotulo.upper() abaixo), então
    # cada caractere ocupa bem mais espaço que numa string mista qualquer; um
    # fator subestimado permitiria mais caracteres do que cabem de verdade,
    # e o texto "truncado" ainda quebraria em 2 linhas (visto na QA visual
    # com o fator antigo de 0.60 - ainda não era conservador o suficiente).
    # Garante rótulo numa linha só, então a posição fixa do contexto abaixo
    # (calculada a partir de h) nunca mais é invadida por uma 2ª linha.
    texto_rotulo = rotulo.upper()
    tamanho_rotulo = 10.0
    if texto_rotulo:
        max_chars = max(1, int(largura_texto / ((tamanho_rotulo / 72.0) * 0.72)))
        while tamanho_rotulo > 7.5 and len(texto_rotulo) > max_chars:
            tamanho_rotulo -= 0.5
            max_chars = max(1, int(largura_texto / ((tamanho_rotulo / 72.0) * 0.72)))
        if len(texto_rotulo) > max_chars:
            texto_rotulo = texto_rotulo[:max(1, max_chars - 1)].rstrip() + "…"
    # `deslocamento_rotulo` (21/08/2026, Fase 3 - cartões sem contexto do
    # slide "Recuperação de Shelf") é um deslocamento FIXO a partir do topo
    # do cartão, não relativo a h - a usuária desenhou esses cartões mais
    # baixos que os demais (cartão sem 3ª linha de contexto não precisa da
    # mesma altura), então o rótulo cola direto embaixo do valor em vez de
    # ficar ancorado no rodapé de um cartão alto. Não usar essa opção não
    # muda em nada o comportamento dos outros 12 lugares que chamam este
    # cartão (todos continuam com y_rotulo relativo a h, já validado na
    # Parte 1).
    if deslocamento_rotulo is not None and not contexto:
        y_rotulo = y + deslocamento_rotulo
    else:
        y_rotulo = y + h - (0.32 if contexto else 0.18)
    _texto(slide, x + pad, y_rotulo, w - 2 * pad, 0.16, texto_rotulo, tamanho=tamanho_rotulo, negrito=True, cor=CINZA_TEXTO)
    if contexto:
        _texto(slide, x + pad, y + h - 0.16, w - 2 * pad, 0.16, contexto, tamanho=9.5, cor=cor_contexto or CINZA_TEXTO)


def _linha_kpis(slide, y, kpis, altura=0.85, deslocamento_rotulo=None, tamanho_valor_base=27):
    n = len(kpis)
    largura_total = LARGURA_IN - 2 * MARGEM_IN
    gap = 0.22
    largura_card = (largura_total - gap * (n - 1)) / n
    x = MARGEM_IN
    for k in kpis:
        _cartao_kpi(slide, x, y, largura_card, altura, k["valor"], k["rotulo"],
                    k.get("cor", AZUL_INSTITUCIONAL), k.get("contexto"), k.get("cor_contexto"),
                    deslocamento_rotulo, tamanho_valor_base)
        x += largura_card + gap


def _caber_no_espaco(texto, largura_in, altura_in, tamanho_pt, espacamento=1.12):
    """Corta o texto (com reticências ao final) se ele não couber na área
    reservada, dado o tamanho de fonte - rede de segurança genérica contra
    estourar a caixa quando o texto de origem varia de tamanho de um mês
    pro outro (ex.: resumo_narrado muda conforme quantas categorias/baixas
    existem no recorte, e as colunas de Avanços/Atenções/Decisões mudam de
    tamanho conforme quantos pontos existem no mês). Estimativa por
    largura média de caractere, não medição real de texto (python-pptx não
    faz layout de fonte) - por isso usa uma margem de segurança (0.82) em
    vez do limite teórico exato."""
    if not texto:
        return ""
    largura_media_char = (tamanho_pt / 72.0) * 0.52
    chars_por_linha = max(10, int(largura_in / largura_media_char))
    altura_linha = (tamanho_pt / 72.0) * espacamento * 1.22
    linhas_disponiveis = max(1, int(altura_in / altura_linha))
    max_chars = int(chars_por_linha * linhas_disponiveis * 0.82)
    if len(texto) <= max_chars:
        return texto
    cortado = texto[:max_chars].rsplit(" ", 1)[0].rstrip(",.;:—-")
    return cortado + "…"


def _altura_necessaria_caixa_leitura(texto, largura_caixa, tamanho_pt, altura_rotulo=0.26, pad=0.20, espacamento=1.12):
    """Calcula a altura mínima que uma _caixa_leitura precisa pra caber `texto`
    sem acionar a rede de segurança de truncamento de _caber_no_espaco -
    replica exatamente a mesma estimativa de largura de caractere/altura de
    linha usada lá (proxy por largura média, não medição real de fonte), pra
    que a altura recomendada aqui seja consistente com o que de fato caberia
    (20/08/2026: corrige caixas com altura fixa que truncavam texto normal do
    mês porque a altura reservada tinha sido um chute, não calculada a partir
    do texto real que a caixa acaba recebendo)."""
    if not texto:
        return altura_rotulo + pad * 1.1
    largura_texto = largura_caixa - 2 * pad
    largura_media_char = (tamanho_pt / 72.0) * 0.52
    chars_por_linha = max(10, int(largura_texto / largura_media_char))
    max_chars_linha = max(1, int(chars_por_linha * 0.82))
    linhas_necessarias = max(1, -(-len(texto) // max_chars_linha))  # teto da divisão
    altura_linha = (tamanho_pt / 72.0) * espacamento * 1.22
    # +0.03in de margem (21/08/2026, CORREÇÃO DE BUG): sem isso, a altura
    # devolvida aqui é o produto EXATO linhas_necessarias * altura_linha (+
    # paddings) - quando essa altura volta pra _caber_no_espaco (chamada
    # de dentro de _caixa_leitura com a MESMA altura_linha), a divisão
    # ponto-flutuante `altura_in / altura_linha` que deveria reproduzir
    # linhas_necessarias às vezes cai um pouco abaixo por erro de
    # arredondamento (ex.: 2.9999... em vez de 3.0) - e como
    # _caber_no_espaco trunca com int() em vez de arredondar, isso conta
    # UMA LINHA A MENOS do que o pretendido e corta o texto com "…" mesmo
    # numa caixa dimensionada "exatamente" pro texto (visto na QA visual:
    # caixas de Resumo cortando a mensagem mesmo com altura_resumo já
    # calculada por esta função). A margem fixa evita cair do lado errado
    # da divisão sem mudar a fonte/densidade de forma perceptível.
    return pad * 0.6 + altura_rotulo + linhas_necessarias * altura_linha + pad * 0.5 + 0.03


def _caixa_leitura(slide, x, y, w, h, rotulo, texto, cor_fundo=OFF_WHITE, cor_rotulo=VERDE_AMAZONIA, tamanho_texto=13):
    _retangulo(slide, x, y, w, h, cor_fill=cor_fundo, raio=0.10)
    pad = 0.20
    altura_rotulo = 0.26
    _texto(slide, x + pad, y + pad * 0.6, w - 2 * pad, altura_rotulo, rotulo.upper(), tamanho=11, negrito=True, cor=cor_rotulo)
    y_texto = y + pad * 0.6 + altura_rotulo
    largura_texto = w - 2 * pad
    altura_texto = (y + h) - y_texto - pad * 0.5
    texto = _caber_no_espaco(texto, largura_texto, altura_texto, tamanho_texto)
    _texto(slide, x + pad, y_texto, largura_texto, altura_texto, texto,
           tamanho=tamanho_texto, cor=CINZA_TEXTO, espacamento=1.12)


def _linhas_estimadas(texto, largura_in, tamanho_pt, espacamento=1.1):
    """Estima quantas linhas um texto vai ocupar numa largura/fonte dadas -
    mesma lógica de _caber_no_espaco (proxy por largura média de caractere,
    já que python-pptx não faz layout de texto de verdade), usada aqui pra
    espaçar itens de lista sem empilhar um em cima do outro quando o texto
    quebra em mais linhas do que caberia com uma estimativa fixa."""
    largura_media_char = (tamanho_pt / 72.0) * 0.52
    chars_por_linha = max(10, int(largura_in / largura_media_char))
    return max(1, -(-len(texto) // chars_por_linha))  # teto da divisão


def _lista_com_marcadores(slide, x, y, w, h, itens, cor_marcador=VERDE_AMAZONIA, tamanho=12.5, espaco_linha=0.24):
    if not itens:
        itens = ["Sem observações relevantes neste recorte."]
    yy = y
    exibidos = 0
    for idx, item in enumerate(itens):
        linhas_estimadas = _linhas_estimadas(item, w - 0.22, tamanho)
        altura_item = espaco_linha * linhas_estimadas + 0.10
        # A altura `h` reservada pra esta lista é um orçamento (calculado em
        # _slide_resumo_executivo a partir da coluna com mais itens no mês),
        # não uma garantia - em meses com itens demais mesmo na fonte mínima
        # já tentada, o conteúdo real pode passar do orçamento. Antes disso
        # simplesmente estourava por baixo da caixa "Mensagem Central"
        # (desenhada depois, por cima) e o item ficava escondido sem
        # nenhum indício pro leitor; agora para de desenhar e avisa quantos
        # itens ficaram de fora (20/08/2026 - motivado pelo FEFO real virar
        # um 5º item de ATENÇÕES em meses com taxa de quebra fora da faixa).
        if h and exibidos > 0 and (yy - y + altura_item) > h:
            restantes = len(itens) - idx
            _texto(slide, x + 0.22, yy, w - 0.22, 0.24, f"+ {restantes} adicional(is) — ver Scorecard.",
                   tamanho=max(9, tamanho - 1), cor=CINZA_TEXTO, italico=True)
            return yy
        _texto(slide, x, yy, 0.18, 0.24, "•", tamanho=tamanho, negrito=True, cor=cor_marcador)
        _texto(slide, x + 0.22, yy, w - 0.22, 0.6, item, tamanho=tamanho, cor=CINZA_TEXTO, espacamento=1.1)
        yy += altura_item
        exibidos += 1
    return yy


def _badge_status(slide, x, y, w, h, texto, cor):
    shape = _retangulo(slide, x, y, w, h, cor_fill=cor, raio=0.5)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.name = FONTE_TEXTO
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    return shape


def _grafico_categoria(slide, x, y, w, h, categorias, nome_serie, valores, tipo=XL_CHART_TYPE.COLUMN_CLUSTERED,
                        cor_serie=VERDE_AMAZONIA, cores_pontos=None, formato_numero='0"%"'):
    """Gráfico nativo de UM eixo/uma série só (convenção já usada nos
    dashboards do próprio Atlas: nunca eixo duplo) - cada categoria pode ter
    cor própria via cores_pontos (ex.: status por mês/farol)."""
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series(nome_serie, valores)
    gframe = slide.shapes.add_chart(tipo, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    try:
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10)
        plot.data_labels.font.color.rgb = CINZA_TEXTO
        plot.data_labels.number_format = formato_numero
        plot.data_labels.number_format_is_linked = False
    except Exception:
        pass

    serie = plot.series[0]
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = cor_serie
    serie.format.line.fill.background()

    if cores_pontos:
        pontos = serie.points
        for i, cor in enumerate(cores_pontos):
            if cor is None or i >= len(valores):
                continue
            pontos[i].format.fill.solid()
            pontos[i].format.fill.fore_color.rgb = cor

    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(10)
        cat_ax.tick_labels.font.color.rgb = CINZA_TEXTO
        cat_ax.format.line.color.rgb = CINZA_CLARO
        val_ax = chart.value_axis
        val_ax.tick_labels.font.size = Pt(10)
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = CINZA_CLARO
        val_ax.format.line.fill.background()
        val_ax.visible = True
    except Exception:
        pass
    return chart


def _grafico_categoria_multi(slide, x, y, w, h, categorias, series, tipo=XL_CHART_TYPE.COLUMN_CLUSTERED,
                              formato_numero='#,##0', mostrar_rotulos=True):
    """Mesma convenção de UM eixo/UMA escala só (nunca dual-axis) - só que
    com VÁRIAS séries lado a lado sobre esse único eixo (ex.: Passivos vs.
    Resultado de Inventário, ambos em R$; Entradas vs. Saídas, idem). `series`
    é uma lista de (nome, valores, cor). Isso NÃO é eixo duplo - é a mesma
    régua compartilhada por todas as séries, só que com mais de uma coluna
    por categoria (20/08/2026, pros slides de Passivos/Fluxo de Inventário).

    `mostrar_rotulos=False` (22/08/2026, Fase 2) desliga o rótulo de valor -
    com MUITAS séries empilhadas e segmentos pequenos (ex.: Risco por
    Almoxarifado/Grupo do Farol de Shelf-Life, 3-4 status por coluna), o
    rótulo de cada segmento fica ilegível/sobreposto (visto na QA visual) -
    nesses casos o eixo numérico + a legenda já bastam pra leitura."""
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    for nome, valores, _cor in series:
        chart_data.add_series(nome, valores)
    gframe = slide.shapes.add_chart(tipo, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = gframe.chart
    chart.has_legend = True
    try:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = CINZA_TEXTO
    except Exception:
        pass
    chart.has_title = False

    plot = chart.plots[0]
    try:
        plot.has_data_labels = mostrar_rotulos
        if mostrar_rotulos:
            plot.data_labels.font.size = Pt(9)
            plot.data_labels.font.color.rgb = CINZA_TEXTO
            plot.data_labels.number_format = formato_numero
            plot.data_labels.number_format_is_linked = False
    except Exception:
        pass

    for serie_obj, (_nome, _valores, cor) in zip(plot.series, series):
        serie_obj.format.fill.solid()
        serie_obj.format.fill.fore_color.rgb = cor
        serie_obj.format.line.fill.background()

    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(10)
        cat_ax.tick_labels.font.color.rgb = CINZA_TEXTO
        cat_ax.format.line.color.rgb = CINZA_CLARO
        val_ax = chart.value_axis
        val_ax.tick_labels.font.size = Pt(10)
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = CINZA_CLARO
        val_ax.format.line.fill.background()
        val_ax.visible = True
    except Exception:
        pass
    return chart


# ---------------------------------------------------------------------------
# Gráficos combo (barra + linha) — edição direta do XML do gráfico
# ---------------------------------------------------------------------------
# O python-pptx não expõe gráfico combinado (dois tipos de série na mesma
# área de plotagem) na sua API de alto nível — só é possível via edição
# direta do XML do gráfico (chart._chartSpace). Pedido não-negociável do
# usuário, 22/08/2026 ("Essas premissas são inegociáveis"), rejeitando a
# simplificação anterior (selo de texto no lugar da linha de tendência) e
# pedindo a curva de Pareto/Distribuição por Magnitude iguais ao modelo
# aprovado (combo dual-axis). Três situações, três funções:
#
#   - _adicionar_linha_combo_eixo_unico: a 2ª série é a MESMA métrica/escala
#     da barra (ex.: tendência de % de acurácia sobre barra de % de
#     acurácia) — compartilha o par de eixos já criado pelo c:barChart, sem
#     violar a convenção "nunca eixo duplo" do resto do MBR (não é dual-axis,
#     é a mesma régua).
#   - _adicionar_linha_combo_eixo_duplo: a 2ª série é uma métrica de UNIDADE
#     DIFERENTE (ex.: Valor R$ x % acumulado; Nº de itens x Valor R$) — cria
#     um 2º par de eixos (valAx à direita + catAx oculto, só pra dar
#     crossAx ao valAx) e É dual-axis de propósito, excepção documentada à
#     convenção da casa porque o próprio modelo aprovado pede exatamente
#     isso nesses 2 gráficos (Concentração de Risco/Distribuição por
#     Magnitude), não nos demais.
#   - _adicionar_rotulo_total_empilhado (21/08/2026): variação da mesma
#     técnica (2ª série, eixo único, sem representação visual própria) pra
#     rótulo de TOTAL por coluna num gráfico de barras EMPILHADAS (ex.:
#     "Total de Baixas por Mês" - soma dos motivos empilhados naquele mês) -
#     não é tendência nem métrica diferente, é só onde pousar um número que
#     o próprio c:barChart não sabe calcular por coluna.
#
# As funções operam sobre um gráfico de barras já criado por
# add_chart (via _grafico_categoria/_grafico_categoria_multi ou chamada
# direta) — a nova série entra como <c:lineChart> logo depois do
# <c:barChart> existente (schema exige todos os "chart type elements"
# juntos, antes dos elementos de eixo). Os valores da série nova só entram
# no <c:numCache> (cache renderizado, o que o PowerPoint/LibreOffice
# realmente desenha) — a planilha embutida do gráfico (aberta só se o
# usuário clicar em "Editar Dados") não é atualizada; isso é aceitável
# porque o MBR não é editado depois de gerado.
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _esc_xml(valor) -> str:
    return html.escape(str(valor), quote=True)


def _hex_cor(cor: RGBColor) -> str:
    return str(cor)


def _cor_rotulo_contraste(cor_fundo: RGBColor) -> RGBColor:
    """Preto ou branco pra texto do rótulo, conforme a luminância da cor de
    fundo da barra - usado quando o rótulo passa a ficar DENTRO da coluna
    (INSIDE_END), onde a cor de fundo já não é mais branca (ver nota em
    _grafico_categoria_com_tendencia)."""
    r, g, b = cor_fundo[0], cor_fundo[1], cor_fundo[2]
    luminancia = 0.299 * r + 0.587 * g + 0.114 * b
    return RGBColor(0x2A, 0x2A, 0x2A) if luminancia >= 140 else RGBColor(0xFF, 0xFF, 0xFF)


def _adicionar_linha_combo_eixo_unico(chart, categorias, valores_linha, cor_hex,
                                       nome_serie="Tendência linear", tracejado=True):
    """Adiciona uma 2ª série como <c:lineChart> tracejado, compartilhando o
    MESMO par de eixos do <c:barChart> já existente no `chart` (mesma
    escala — ver nota do bloco acima). `valores_linha` deve ter o mesmo
    tamanho de `categorias`, com None nas posições sem ponto (mesma
    convenção de "buraco" já usada pelo resto do MBR)."""
    chart_space = chart._chartSpace
    plot_area = chart_space.find(qn("c:chart")).find(qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))
    if bar_chart is None:
        return
    ax_ids = [e.get("val") for e in bar_chart.findall(qn("c:axId"))]
    if len(ax_ids) < 2:
        return
    cat_ax_id, val_ax_id = ax_ids[0], ax_ids[1]

    n = len(categorias)
    pts_cat = "".join(f'<c:pt idx="{i}"><c:v>{_esc_xml(c)}</c:v></c:pt>' for i, c in enumerate(categorias))
    pts_val = "".join(
        f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(valores_linha) if v is not None
    )
    dash = '<a:prstDash val="dash"/>' if tracejado else ""
    line_xml = f'''<c:lineChart {nsdecls("c")}>
      <c:grouping val="standard"/>
      <c:varyColors val="0"/>
      <c:ser>
        <c:idx val="1"/>
        <c:order val="1"/>
        <c:tx><c:v>{_esc_xml(nome_serie)}</c:v></c:tx>
        <c:spPr>
          <a:ln w="17780" xmlns:a="{_A_NS}">
            <a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill>
            {dash}
          </a:ln>
        </c:spPr>
        <c:marker>
          <c:symbol val="circle"/>
          <c:size val="5"/>
          <c:spPr xmlns:a="{_A_NS}">
            <a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill>
            <a:ln><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>
          </c:spPr>
        </c:marker>
        <c:cat>
          <c:strRef>
            <c:f>Sheet1!$A$2:$A${1 + n}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{pts_cat}</c:strCache>
          </c:strRef>
        </c:cat>
        <c:val>
          <c:numRef>
            <c:f>Sheet1!$C$2:$C${1 + n}</c:f>
            <c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="{n}"/>{pts_val}</c:numCache>
          </c:numRef>
        </c:val>
        <c:smooth val="0"/>
      </c:ser>
      <c:marker val="1"/>
      <c:axId val="{cat_ax_id}"/>
      <c:axId val="{val_ax_id}"/>
    </c:lineChart>'''
    bar_chart.addnext(parse_xml(line_xml))


def _adicionar_rotulo_total_empilhado(chart, categorias, totais, cor_texto_hex="595959",
                                       formato_numero='#,##0', tamanho_pt=9.5, nome_serie="Total"):
    """Rótulo do valor TOTAL por coluna, num gráfico de barras/colunas
    EMPILHADAS (21/08/2026, pedido da usuária pro gráfico "Total de Baixas
    por Mês": "adicione o rótulo do valor total por gráfico, não por
    motivo... mas sim da somatório dos eventos por mês"). O rótulo comum
    de segmento (`plot.has_data_labels`) não serve aqui - com até 9 motivos
    empilhados numa coluna só, um rótulo por segmento fica ilegível/
    sobreposto (é exatamente por isso que esse gráfico desliga o rótulo por
    segmento - ver `mostrar_rotulos=not empilhado` em
    _slide_baixas_operacionais_externo). A técnica-padrão pra rótulo de
    TOTAL numa pilha (a mesma que Excel/PowerPoint usam quando alguém pede
    isso manualmente) é uma 2ª série de LINHA, sem linha visível e sem
    marcador, com o valor de cada ponto igual à SOMA da coluna naquela
    categoria - só o rótulo de dado dessa série aparece, ancorado acima do
    ponto (`dLblPos="t"`), pousando exatamente no topo da pilha. Como essa
    série não tem nenhuma representação visual própria, também é removida
    da legenda (`c:legendEntry`/`delete`) - senão apareceria como um item
    de legenda "Total" sem swatch nenhum pra explicar, mais confuso que
    ajuda."""
    chart_space = chart._chartSpace
    plot_area = chart_space.find(qn("c:chart")).find(qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))
    if bar_chart is None:
        return
    ax_ids = [e.get("val") for e in bar_chart.findall(qn("c:axId"))]
    if len(ax_ids) < 2:
        return
    cat_ax_id, val_ax_id = ax_ids[0], ax_ids[1]
    idx = len(bar_chart.findall(qn("c:ser")))

    n = len(categorias)
    pts_cat = "".join(f'<c:pt idx="{i}"><c:v>{_esc_xml(c)}</c:v></c:pt>' for i, c in enumerate(categorias))
    pts_val = "".join(
        f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(totais) if v is not None
    )
    line_xml = f'''<c:lineChart {nsdecls("c")}>
      <c:grouping val="standard"/>
      <c:varyColors val="0"/>
      <c:ser>
        <c:idx val="{idx}"/>
        <c:order val="{idx}"/>
        <c:tx><c:v>{_esc_xml(nome_serie)}</c:v></c:tx>
        <c:spPr>
          <a:ln xmlns:a="{_A_NS}"><a:noFill/></a:ln>
        </c:spPr>
        <c:marker><c:symbol val="none"/></c:marker>
        <c:dLbls>
          <c:numFmt formatCode="{formato_numero}" sourceLinked="0"/>
          <c:spPr xmlns:a="{_A_NS}"><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
          <c:txPr xmlns:a="{_A_NS}">
            <a:bodyPr/>
            <a:lstStyle/>
            <a:p><a:pPr><a:defRPr sz="{int(round(tamanho_pt * 100))}" b="1"><a:solidFill><a:srgbClr val="{cor_texto_hex}"/></a:solidFill></a:defRPr></a:pPr><a:endParaRPr lang="pt-BR"/></a:p>
          </c:txPr>
          <c:dLblPos val="t"/>
          <c:showLegendKey val="0"/>
          <c:showVal val="1"/>
          <c:showCatName val="0"/>
          <c:showSerName val="0"/>
          <c:showPercent val="0"/>
          <c:showBubbleSize val="0"/>
        </c:dLbls>
        <c:cat>
          <c:strRef>
            <c:f>Sheet1!$A$2:$A${1 + n}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{pts_cat}</c:strCache>
          </c:strRef>
        </c:cat>
        <c:val>
          <c:numRef>
            <c:f>Sheet1!$Z$2:$Z${1 + n}</c:f>
            <c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="{n}"/>{pts_val}</c:numCache>
          </c:numRef>
        </c:val>
        <c:smooth val="0"/>
      </c:ser>
      <c:marker val="0"/>
      <c:axId val="{cat_ax_id}"/>
      <c:axId val="{val_ax_id}"/>
    </c:lineChart>'''
    bar_chart.addnext(parse_xml(line_xml))

    try:
        legend = chart_space.find(qn("c:chart")).find(qn("c:legend"))
        if legend is not None:
            entry_xml = f'<c:legendEntry {nsdecls("c")}><c:idx val="{idx}"/><c:delete val="1"/></c:legendEntry>'
            entry_el = parse_xml(entry_xml)
            legend_pos = legend.find(qn("c:legendPos"))
            if legend_pos is not None:
                legend_pos.addnext(entry_el)
            else:
                legend.insert(0, entry_el)
    except Exception:
        pass


def _adicionar_linha_combo_eixo_duplo(chart, categorias, valores_linha, cor_hex,
                                       nome_serie, formato_numero_eixo="General", tracejado=False):
    """Adiciona uma 2ª série como <c:lineChart> em EIXO SECUNDÁRIO (dual-
    axis de propósito — ver nota do bloco acima) a um <c:barChart> já
    existente. Cria um valAx novo (à direita) e um catAx novo (oculto, só
    existe pra dar crossAx ao valAx — exigência do schema OOXML, mesmo
    padrão usado pelo próprio Excel/PowerPoint pra combo dual-axis)."""
    chart_space = chart._chartSpace
    plot_area = chart_space.find(qn("c:chart")).find(qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))
    if bar_chart is None:
        return
    ax_ids_primarios = [e.get("val") for e in bar_chart.findall(qn("c:axId"))]
    if len(ax_ids_primarios) < 2:
        return
    cat_ax1_id, val_ax1_id = ax_ids_primarios[0], ax_ids_primarios[1]

    existentes = {cat_ax1_id, val_ax1_id}
    cat_ax2_id = str(int(cat_ax1_id) + 1)
    while cat_ax2_id in existentes:
        cat_ax2_id = str(int(cat_ax2_id) + 1)
    existentes.add(cat_ax2_id)
    val_ax2_id = str(int(val_ax1_id) + 1)
    while val_ax2_id in existentes:
        val_ax2_id = str(int(val_ax2_id) + 1)

    n = len(categorias)
    pts_cat = "".join(f'<c:pt idx="{i}"><c:v>{_esc_xml(c)}</c:v></c:pt>' for i, c in enumerate(categorias))
    pts_val = "".join(
        f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(valores_linha) if v is not None
    )
    dash = '<a:prstDash val="dash"/>' if tracejado else ""
    line_xml = f'''<c:lineChart {nsdecls("c")}>
      <c:grouping val="standard"/>
      <c:varyColors val="0"/>
      <c:ser>
        <c:idx val="1"/>
        <c:order val="1"/>
        <c:tx><c:v>{_esc_xml(nome_serie)}</c:v></c:tx>
        <c:spPr>
          <a:ln w="19050" xmlns:a="{_A_NS}">
            <a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill>
            {dash}
          </a:ln>
        </c:spPr>
        <c:marker>
          <c:symbol val="circle"/>
          <c:size val="5"/>
          <c:spPr xmlns:a="{_A_NS}">
            <a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill>
            <a:ln><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>
          </c:spPr>
        </c:marker>
        <c:cat>
          <c:strRef>
            <c:f>Sheet1!$A$2:$A${1 + n}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{pts_cat}</c:strCache>
          </c:strRef>
        </c:cat>
        <c:val>
          <c:numRef>
            <c:f>Sheet1!$D$2:$D${1 + n}</c:f>
            <c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="{n}"/>{pts_val}</c:numCache>
          </c:numRef>
        </c:val>
        <c:smooth val="0"/>
      </c:ser>
      <c:marker val="1"/>
      <c:axId val="{cat_ax2_id}"/>
      <c:axId val="{val_ax2_id}"/>
    </c:lineChart>'''
    bar_chart.addnext(parse_xml(line_xml))

    val_ax2_xml = f'''<c:valAx {nsdecls("c")}>
      <c:axId val="{val_ax2_id}"/>
      <c:scaling><c:orientation val="minMax"/></c:scaling>
      <c:delete val="0"/>
      <c:axPos val="r"/>
      <c:numFmt formatCode="{_esc_xml(formato_numero_eixo)}" sourceLinked="0"/>
      <c:majorTickMark val="out"/>
      <c:minorTickMark val="none"/>
      <c:tickLblPos val="nextTo"/>
      <c:spPr xmlns:a="{_A_NS}"><a:ln><a:noFill/></a:ln></c:spPr>
      <c:txPr xmlns:a="{_A_NS}">
        <a:bodyPr/><a:lstStyle/>
        <a:p><a:pPr><a:defRPr sz="900"><a:solidFill><a:srgbClr val="4A4A4A"/></a:solidFill></a:defRPr></a:pPr><a:endParaRPr lang="pt-BR"/></a:p>
      </c:txPr>
      <c:crossAx val="{cat_ax2_id}"/>
      <c:crosses val="max"/>
    </c:valAx>'''
    cat_ax2_xml = f'''<c:catAx {nsdecls("c")}>
      <c:axId val="{cat_ax2_id}"/>
      <c:scaling><c:orientation val="minMax"/></c:scaling>
      <c:delete val="1"/>
      <c:axPos val="b"/>
      <c:majorTickMark val="out"/>
      <c:minorTickMark val="none"/>
      <c:tickLblPos val="nextTo"/>
      <c:crossAx val="{val_ax2_id}"/>
      <c:crosses val="autoZero"/>
      <c:auto val="1"/>
      <c:lblAlgn val="ctr"/>
      <c:lblOffset val="100"/>
      <c:noMultiLvlLbl val="0"/>
    </c:catAx>'''
    val_ax1_el = plot_area.find(qn("c:valAx"))
    val_ax1_el.addnext(parse_xml(val_ax2_xml))
    val_ax1_el.addnext(parse_xml(cat_ax2_xml))


def _grafico_categoria_com_tendencia(slide, x, y, w, h, categorias, nome_serie, valores,
                                      cor_serie=VERDE_AMAZONIA, formato_numero='0"%"',
                                      cor_tendencia=COR_SUCESSO, cores_pontos=None,
                                      nome_tendencia="Tendência linear"):
    """Como _grafico_categoria, mas desenha por cima uma linha de
    tendência REAL (regressão linear — ver _pontos_tendencia_linear),
    tracejada, na MESMA escala da barra — substitui, nos 3 slides de
    evolução mensal com tendência (Painel de Inventário/Item a Item, IAP,
    IAQ), o selo de texto que era a única representação da tendência
    antes disso (pedido não-negociável do usuário, 22/08/2026: "o gráfico
    não conta a história"). O selo de texto (_tendencia_linear) continua
    sendo desenhado también, no cabeçalho do gráfico — o modelo aprovado
    mantém os dois."""
    chart = _grafico_categoria(slide, x, y, w, h, categorias, nome_serie, valores,
                                cor_serie=cor_serie, cores_pontos=cores_pontos, formato_numero=formato_numero)
    pontos_linha = _pontos_tendencia_linear(valores)
    if pontos_linha:
        _adicionar_linha_combo_eixo_unico(chart, categorias, pontos_linha, cor_hex=_hex_cor(cor_tendencia),
                                           nome_serie=nome_tendencia)
        # Rótulo de valor da barra pra DENTRO da coluna (INSIDE_END) quando
        # há linha de tendência combo por cima — no topo (padrão) o rótulo
        # colide com a linha/marcador tracejado quando o valor da tendência
        # naquele mês fica próximo ou acima da barra (visto na QA visual,
        # 22/08/2026 - pv-1/pv-2/pv-5.jpg). A cor do texto passa a precisar
        # de contraste com o PREENCHIMENTO da barra (deixou de estar sobre
        # fundo branco) - por ponto quando as barras têm cores variadas
        # (cores_pontos, ex. crítico/ok), por série quando é uma cor só.
        try:
            plot = chart.plots[0]
            plot.data_labels.position = XL_LABEL_POSITION.INSIDE_END
            plot.data_labels.font.color.rgb = _cor_rotulo_contraste(cor_serie)
            if cores_pontos:
                serie = plot.series[0]
                for i, cor_pt in enumerate(cores_pontos):
                    cor_efetiva = cor_pt if cor_pt is not None else cor_serie
                    if i >= len(valores) or valores[i] is None:
                        continue
                    rotulo_pt = serie.points[i].data_label
                    rotulo_pt.position = XL_LABEL_POSITION.INSIDE_END
                    rotulo_pt.font.size = Pt(10)
                    rotulo_pt.font.color.rgb = _cor_rotulo_contraste(cor_efetiva)
        except Exception:
            pass
    return chart


def _grafico_combo_dual_eixo(slide, x, y, w, h, categorias, nome_barra, valores_barra, nome_linha, valores_linha,
                              cor_barra, cor_linha, formato_numero_barra='#,##0', formato_numero_linha='0"%"',
                              tracejado_linha=False, mostrar_rotulos_barra=False, cores_pontos_barra=None):
    """Gráfico combo com EIXO DUPLO de propósito (ver nota no bloco acima
    sobre por que esta é uma excepção documentada à convenção "nunca eixo
    duplo" do resto do MBR) — usado nos 2 gráficos em que o modelo
    aprovado (mockup v8) pede duas métricas de unidade diferente na mesma
    área de plotagem: Concentração de Risco (Valor R$ x % acumulado) e
    Distribuição por Magnitude (Nº de itens x Valor total R$)."""
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series(nome_barra, valores_barra)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    try:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = CINZA_TEXTO
    except Exception:
        pass

    plot = chart.plots[0]
    try:
        plot.has_data_labels = mostrar_rotulos_barra
        if mostrar_rotulos_barra:
            plot.data_labels.font.size = Pt(9)
            plot.data_labels.font.color.rgb = CINZA_TEXTO
            plot.data_labels.number_format = formato_numero_barra
            plot.data_labels.number_format_is_linked = False
    except Exception:
        pass

    serie = plot.series[0]
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = cor_barra
    serie.format.line.fill.background()
    if cores_pontos_barra:
        pontos = serie.points
        for i, cor in enumerate(cores_pontos_barra):
            if cor is None or i >= len(valores_barra):
                continue
            pontos[i].format.fill.solid()
            pontos[i].format.fill.fore_color.rgb = cor

    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(9)
        cat_ax.tick_labels.font.color.rgb = CINZA_TEXTO
        cat_ax.format.line.color.rgb = CINZA_CLARO
        val_ax = chart.value_axis
        val_ax.tick_labels.font.size = Pt(9)
        val_ax.tick_labels.font.color.rgb = CINZA_TEXTO
        val_ax.tick_labels.number_format = formato_numero_barra
        val_ax.tick_labels.number_format_is_linked = False
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = CINZA_CLARO
        val_ax.format.line.fill.background()
    except Exception:
        pass

    _adicionar_linha_combo_eixo_duplo(chart, categorias, valores_linha, cor_hex=_hex_cor(cor_linha),
                                       nome_serie=nome_linha, formato_numero_eixo=formato_numero_linha,
                                       tracejado=tracejado_linha)
    return chart


def _tabela(slide, x, y, w, h, cabecalhos, linhas, larguras_relativas=None,
            cor_cabecalho=AZUL_INSTITUCIONAL, tamanho_fonte=12):
    n_linhas = len(linhas) + 1
    n_col = len(cabecalhos)
    shape = slide.shapes.add_table(n_linhas, n_col, Inches(x), Inches(y), Inches(w), Inches(h))
    tabela = shape.table

    if larguras_relativas:
        total = sum(larguras_relativas)
        for i, fracao in enumerate(larguras_relativas):
            tabela.columns[i].width = Inches(w * fracao / total)

    for j, cabecalho in enumerate(cabecalhos):
        cel = tabela.cell(0, j)
        cel.fill.solid()
        cel.fill.fore_color.rgb = cor_cabecalho
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cel.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cabecalho
        p.font.size = Pt(tamanho_fonte)
        p.font.bold = True
        p.font.name = FONTE_TEXTO
        p.font.color.rgb = BRANCO

    for i, linha in enumerate(linhas, start=1):
        cor_fundo = BRANCO if i % 2 == 1 else CINZA_LINHA_PAR
        for j, valor in enumerate(linha):
            if isinstance(valor, tuple):
                texto_valor, cor_fonte, negrito = valor
            else:
                texto_valor, cor_fonte, negrito = valor, CINZA_TEXTO, False
            cel = tabela.cell(i, j)
            cel.fill.solid()
            cel.fill.fore_color.rgb = cor_fundo
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cel.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "" if texto_valor is None else str(texto_valor)
            p.font.size = Pt(tamanho_fonte - 1)
            p.font.name = FONTE_TEXTO
            p.font.color.rgb = cor_fonte
            p.font.bold = negrito
    return tabela


# ---------------------------------------------------------------------------
# Dashboards externos (Auditoria > Outros Dashboards) - leitura do html_content
# salvo em DashboardExterno e extração via dashboards_externos_extrator.py.
# Nenhum dos dois helpers levanta exceção se o slot não tiver arquivo enviado
# ou se a extração falhar (arquivo de formato inesperado) - o slide correspondente
# trata "tem_dados": False / None mostrando que o dashboard ainda não foi enviado,
# em vez de quebrar a geração do MBR inteiro por causa de um anexo (20/08/2026).
# ---------------------------------------------------------------------------
def _extrair_dashboard_externo(db: Session, chave: str, extrator, mes: str) -> dict:
    registro = db.query(models.DashboardExterno).filter_by(chave=chave).first()
    if not registro or not registro.html_content:
        return {"tem_dados": False, "enviado": False}
    try:
        resultado = extrator(registro.html_content, mes)
    except Exception:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True}
    if resultado is None:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True}
    resultado["enviado"] = True
    resultado["enviado_em"] = registro.enviado_em.strftime("%d/%m/%Y %H:%M") if registro.enviado_em else None
    return resultado


def _extrair_dashboard_externo_sem_mes(db: Session, chave: str, extrator) -> dict:
    registro = db.query(models.DashboardExterno).filter_by(chave=chave).first()
    if not registro or not registro.html_content:
        return {"tem_dados": False, "enviado": False}
    try:
        resultado = extrator(registro.html_content)
    except Exception:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True}
    if resultado is None:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True}
    resultado["tem_dados"] = True
    resultado["enviado"] = True
    resultado["enviado_em"] = registro.enviado_em.strftime("%d/%m/%Y %H:%M") if registro.enviado_em else None
    return resultado


def _normalizar_nome_indicador(nome: str) -> str:
    sem_acento = unicodedata.normalize("NFKD", nome or "").encode("ascii", "ignore").decode("ascii")
    return re.sub(r"\s+", " ", sem_acento).strip().lower()


def _extrair_dashboard_externo_por_nome(db: Session, nome_alvo: str, extrator, mes: str) -> dict:
    """Como _extrair_dashboard_externo, mas busca o DashboardExterno pelo
    nome_exibicao (não por uma chave fixa) - usado por indicadores DINÂMICOS
    que ganharam extração/slide dedicados (20/08/2026: "Dispersão de Ficha
    Técnica", pedido do usuário). A chave de um indicador dinâmico é gerada
    por slug do nome no momento em que o admin cria (ver
    dashboards_externos_router._slugificar) e pode variar se houver colisão
    de nome - buscar pelo nome (normalizado: sem acento, sem espaço duplo,
    minúsculo) é mais robusto do que fixar a chave esperada no código."""
    alvo = _normalizar_nome_indicador(nome_alvo)
    registro = next(
        (r for r in db.query(models.DashboardExterno).all() if _normalizar_nome_indicador(r.nome_exibicao) == alvo),
        None,
    )
    if not registro or not registro.html_content:
        return {"tem_dados": False, "enviado": False, "_chave_dashboard_externo": None}
    try:
        resultado = extrator(registro.html_content, mes)
    except Exception:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True, "_chave_dashboard_externo": registro.chave}
    if resultado is None:
        return {"tem_dados": False, "enviado": True, "erro_extracao": True, "_chave_dashboard_externo": registro.chave}
    resultado["enviado"] = True
    resultado["enviado_em"] = registro.enviado_em.strftime("%d/%m/%Y %H:%M") if registro.enviado_em else None
    resultado["_chave_dashboard_externo"] = registro.chave
    return resultado


def _extrair_resumo_auditoria_fefo(db: Session, mes: str) -> dict:
    """FEFO do MBR (20/08/2026): pedido do usuário pra trocar a fonte do slide de
    FEFO do dashboard "Controle de FEFO" (Auditoria > Outros Dashboards,
    dashboards_externos_extrator.extrair_fefo) pela "Auditoria FEFO importada"
    (models.AuditoriaFefo / fefo.calcular_resumo_auditoria_fefo) - a mesma base
    que já alimenta o painel "Auditoria FEFO — histórico importado" na tela FEFO.
    Motivo (mensagem do usuário): "o mesmo tem mais informações e bases de
    registro" - de fato o AuditoriaFefo guarda lote movimentado, validade e o
    lote mais antigo disponível por registro, enquanto o dashboard externo só
    tinha os totais já agregados pelo estagiário.

    Diferente de _extrair_dashboard_externo (que lê um DashboardExterno.html_content),
    aqui os dados vêm de linhas já importadas na tabela AuditoriaFefo (Excel diário
    ou dashboard HTML consolidado do André, importados na própria tela FEFO - ver
    fefo.importar_auditoria_fefo_diaria/importar_auditoria_fefo_consolidada). Por
    isso "enviado" aqui significa "existe ALGUMA linha de AuditoriaFefo já
    importada" (não depende do mês do relatório), enquanto "tem_dados" depende do
    mês (pode não ter nenhum registro DESTE mês específico)."""
    algum_registro_existe = db.query(models.AuditoriaFefo.id).first() is not None
    if not algum_registro_existe:
        return {"tem_dados": False, "enviado": False, "mes": mes}

    ano_int, mes_int = (int(parte) for parte in mes.split("-"))
    data_inicio = date(ano_int, mes_int, 1)
    ultimo_dia = calendar.monthrange(ano_int, mes_int)[1]
    data_fim = date(ano_int, mes_int, ultimo_dia)

    resumo = fefo_mod.calcular_resumo_auditoria_fefo(db, data_inicio, data_fim)
    resumo["mes"] = mes
    resumo["enviado"] = True
    if not resumo.get("total_auditaveis"):
        resumo["tem_dados"] = False
    else:
        resumo["tem_dados"] = True
    return resumo


def _coletar_scorecard_inventario_almoxarifado(db: Session, usuario: models.Usuario, mes: str) -> list:
    """Scorecard de Inventário por Almoxarifado (20/08/2026, pedido do usuário:
    "traga uma análise por almoxarifado, mostrando as evoluções e involuções.
    Com um plano de ação para cada setor, baseado no histórico de inventários
    e na conciliação de movimentados"). Combina, POR almoxarifado:
    - Acurácia item-a-item do mês e a variação vs. o mês anterior
      (fechamento_router.dashboard_evolucao_por_almox_mensal - já quebrado por
      almoxarifado E mês; nenhuma outra função do Painel de Inventário tem as
      duas dimensões ao mesmo tempo, ver claude/scorecards-inventario-riscos.md).
    - IAP (acurácia ponderada por valor) do mês
      (fechamento_router.dashboard_comparativo_por_almoxarifado - só o mês, sem
      série, por isso não entra na variação, só como leitura complementar).
    - Acurácia da reconciliação diária de Movimentados do mês e sua variação
      (movimentados_router.dashboard_evolucao_mensal(almoxarifado=X) - só
      devolve a série de UM almoxarifado por chamada, por isso o loop abaixo).
    Não inventa "plano de ação" livre - a leitura e o próximo passo de cada
    linha são montados por regra a partir desses mesmos números (ver
    _linha_scorecard_almoxarifado).

    "movimentados_aplicavel" (21/08/2026, correção pedida pelo usuário):
    reaproveita Almoxarifado.participa_contagem_diaria - o MESMO cadastro que
    já exclui Box, Box_2, Ativação, Degustação e Loja da Cobertura de
    Conferência (divergencias_router.py) - porque esses almoxarifados
    simplesmente não fazem conciliação diária de Movimentados por decisão
    operacional, não por lacuna de dado. Sem essa marcação, a regra do "pior
    dos dois sinais" tratava a ausência estrutural de Movimentados como
    "Sem histórico" e arrastava pra baixo até almoxarifados com evolução real
    de acurácia de fechamento."""
    almoxarifados = cadastros_router.listar_almoxarifados_cadastro(incluir_inativos=False, usuario=usuario, db=db)
    comparativo_mes = {
        item["almoxarifado"]: item
        for item in fechamento_router.dashboard_comparativo_por_almoxarifado(mes=mes, usuario=usuario, db=db)
    }
    evolucao_fech_por_almox = defaultdict(list)
    for linha in fechamento_router.dashboard_evolucao_por_almox_mensal(usuario=usuario, db=db):
        evolucao_fech_por_almox[linha["almoxarifado"]].append(linha)

    resultado = []
    for almox in almoxarifados:
        codigo = almox.codigo
        serie_fech = sorted(
            (l for l in evolucao_fech_por_almox.get(codigo, []) if l["mes"] <= mes), key=lambda l: l["mes"]
        )
        acuracia_pct = serie_fech[-1]["acuracia_pct"] if serie_fech else None
        delta_acuracia_pp = None
        if len(serie_fech) >= 2 and serie_fech[-1]["acuracia_pct"] is not None and serie_fech[-2]["acuracia_pct"] is not None:
            delta_acuracia_pp = round(serie_fech[-1]["acuracia_pct"] - serie_fech[-2]["acuracia_pct"], 2)

        serie_mov = sorted(
            (l for l in movimentados_router.dashboard_evolucao_mensal(almoxarifado=codigo, usuario=usuario, db=db) if l["mes"] <= mes),
            key=lambda l: l["mes"],
        )
        movimentados_pct = serie_mov[-1]["pct_acuracia"] if serie_mov else None
        delta_movimentados_pp = None
        if len(serie_mov) >= 2 and serie_mov[-1]["pct_acuracia"] is not None and serie_mov[-2]["pct_acuracia"] is not None:
            delta_movimentados_pp = round(serie_mov[-1]["pct_acuracia"] - serie_mov[-2]["pct_acuracia"], 2)

        comp = comparativo_mes.get(codigo)
        resultado.append({
            "almoxarifado": codigo,
            "nome_exibicao": almox.nome_exibicao or codigo,
            "acuracia_pct": acuracia_pct,
            "delta_acuracia_pp": delta_acuracia_pp,
            "iap_pct": comp.get("iap_pct") if comp else None,
            "movimentados_pct": movimentados_pct,
            "delta_movimentados_pp": delta_movimentados_pp,
            "movimentados_aplicavel": bool(almox.participa_contagem_diaria),
        })
    return resultado


def _linha_scorecard_almoxarifado(item: dict) -> dict:
    """Classifica e escreve a leitura/próximo passo de UMA linha do Scorecard
    de Inventário por Almoxarifado, por regra (sem heurística nova de negócio -
    reaproveita os mesmos limiares de acurácia (_LIMIARES) já usados no resto
    do MBR). Prioriza o pior dos dois sinais (fechamento vs. movimentados) pra
    decidir o status da linha - um almoxarifado só está "em avanço" se os dois
    estiverem, no mínimo, estáveis.

    Quando Movimentados não é aplicável (item["movimentados_aplicavel"] ==
    False - ver docstring de _coletar_scorecard_inventario_almoxarifado), o
    sinal de Movimentados é ignorado por completo: o status e o próximo passo
    vêm só do fechamento. Sem isso, um almoxarifado que nunca faz conciliação
    de Movimentados (decisão operacional, não lacuna de dado) ficaria travado
    em "Sem histórico" pra sempre, mesmo com evolução real de acurácia."""
    movimentados_aplicavel = item.get("movimentados_aplicavel", True)
    delta_fech = item["delta_acuracia_pp"]
    delta_mov = item["delta_movimentados_pp"] if movimentados_aplicavel else None
    label_fech, cor_fech = _status_evolucao(delta_fech)

    if movimentados_aplicavel:
        label_mov, cor_mov = _status_evolucao(delta_mov)
        ordem = {"Involução": 0, "Sem histórico": 1, "Estável": 2, "Evolução": 3}
        if ordem[label_fech] <= ordem[label_mov]:
            status_label, status_cor = label_fech, cor_fech
        else:
            status_label, status_cor = label_mov, cor_mov
    else:
        label_mov = None
        status_label, status_cor = label_fech, cor_fech

    partes_leitura = [f"Acurácia {_fmt_pct(item['acuracia_pct'])}"]
    if delta_fech is not None:
        sinal = "+" if delta_fech >= 0 else ""
        partes_leitura.append(f"({sinal}{_fmt_pct(delta_fech)} vs. mês anterior)")
    if movimentados_aplicavel:
        partes_leitura.append(f"· Movimentados {_fmt_pct(item['movimentados_pct'])}")
        if delta_mov is not None:
            sinal = "+" if delta_mov >= 0 else ""
            partes_leitura.append(f"({sinal}{_fmt_pct(delta_mov)})")
    else:
        partes_leitura.append("· Movimentados: não aplicável a este almoxarifado")
    if item.get("iap_pct") is not None:
        partes_leitura.append(f"· IAP {_fmt_pct(item['iap_pct'])}")
    leitura = " ".join(partes_leitura)

    acuracia_critica = item["acuracia_pct"] is not None and item["acuracia_pct"] < _LIMIARES["acuracia"][1]
    if label_fech == "Involução" and label_mov == "Involução":
        proximo_passo = "Reconferência prioritária: acurácia e conciliação de movimentados pioraram juntas neste almoxarifado."
    elif label_fech == "Involução":
        proximo_passo = "Investigar causa raiz das divergências recorrentes de fechamento neste almoxarifado."
    elif label_mov == "Involução":
        proximo_passo = "Reforçar a conciliação diária de movimentados — reconciliação piorou vs. o mês anterior."
    elif acuracia_critica:
        proximo_passo = "Intensificar a cadência de conferência até a acurácia voltar pra faixa saudável."
    elif label_fech == "Sem histórico" and (label_mov == "Sem histórico" or label_mov is None):
        proximo_passo = "Sem histórico suficiente ainda — acompanhar a partir do próximo fechamento."
    else:
        proximo_passo = "Manter a cadência atual de fechamento e conciliação — sem sinal de piora no mês."

    return {
        "frente": item["nome_exibicao"], "status_label": status_label, "status_cor": status_cor,
        "leitura": leitura, "proximo_passo": proximo_passo,
    }


# 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o modelo
# final"): o Scorecard de Mapeamento de Riscos foi removido do relatório (ela
# editou o MBR manualmente pra tirar esse tópico). As duas funções que só
# existiam pra montar aquele slide - _coletar_scorecard_mapeamento_riscos
# (reaproveitava resumo_shelf_life/mapeamento_risco_obsolescencia/
# dispersao_ficha_tecnica_externo/testes_industriais_externo/fefo_externo já
# coletados, e ainda buscava o MESMO indicador do mês anterior via
# _extrair_dashboard_externo(_por_nome)/_extrair_resumo_auditoria_fefo só pra
# calcular a variação) e _linha_risco_com_evolucao (montava cada linha) -
# foram removidas junto, senão ficariam fazendo 3 consultas extra ao mês
# anterior a cada MBR gerado sem nenhum slide pra usar o resultado. Os dicts
# que elas reaproveitavam (resumo_shelf_life, mapeamento_risco_obsolescencia
# etc.) continuam sendo coletados normalmente - outros slides/KPIs ainda leem
# deles (ver comentário de decisão antes de _slide_controle_movimentados).


# ---------------------------------------------------------------------------
# Diário de Bordo / Rotina Master (21/08/2026, pedido do usuário: "Adicione
# ao slide 'Atlas' a visão da rotina master, meu diário de bordo [...] crie
# um indicador paralelo mostrando curva de evolução associada a constância
# e disciplina de manter as tarefas em dia"). Esse indicador não mede
# estoque - mede a disciplina operacional de manter o diário de bordo em
# dia, como leitura complementar de "impacto do Atlas" (a mesma constância
# que sustenta a qualidade dos dados usados no resto deste relatório).
#
# Vem de um app separado do Atlas (Mágio Rotinas / "Rotina Master",
# https://rotinabusiness.lovable.app/), sem integração automática ainda -
# por isso, diferente de todo outro indicador deste MBR, os números aqui
# NÃO vêm de uma consulta ao banco do Atlas. Vieram de navegar direto no
# Dashboard de Performance daquele app, com o filtro de período (De/Até) já
# no mês de fechamento (1º ao último dia) - ver _DIARIO_BORDO_POR_MES.
# Fins de semana são excluídos da leitura de constância porque não têm
# rotina devida naquele app (cumprimento 0% num sábado/domingo é ausência
# de tarefa, não uma falha real).
#
# Só existe dado pro(s) mês(es) já coletado(s) manualmente - pra qualquer
# outro mês, _coletar_indicador_diario_bordo devolve "tem_dados": False em
# vez de reaproveitar por engano o número de um mês errado. Se este
# indicador for pra continuar no MBR mês a mês, precisa de uma coleta manual
# nova (ou uma integração automática, que não existe hoje) a cada geração.
#
# 02/09/2026 (pedido do usuário: "Agora o HTML do diario de bordo passa a
# alimentar o indicador de performance e cumprimento de rotina"): esse dict
# manual virou FALLBACK. _coletar_indicador_diario_bordo agora tenta
# primeiro a extração automática do export HTML do Dashboard de Performance
# (enviado como indicador dinâmico em Auditoria > Outros Dashboards - ver
# dashboards_externos_extrator.extrair_diario_bordo) e só cai neste dict
# quando não há extração automática disponível pra aquele mês. A extração
# automática só traz os KPIs de topo (cumprimento geral, rotinas cumpridas/
# devidas, no prazo/atraso) - NÃO reconstrói sequência/lapsos por dia útil
# nem quebra semanal (o gráfico diário do export é um SVG Recharts com só
# 16 dos 31 dias rotulados no eixo - risco real de atribuir valor ao dia
# errado, ver docstring de extrair_diario_bordo). _slide_diario_bordo mostra
# uma versão mais enxuta do slide quando esses campos não vêm no dado.
_NOMES_DIARIO_BORDO_CANDIDATOS = (
    "Constância e Disciplina — Diário de Bordo",
    "Diário de Bordo",
    "Dashboard de Performance",
)

_DIARIO_BORDO_POR_MES = {
    "2026-07": {
        "cumprimento_geral_pct": 97.0,
        "rotinas_cumpridas": 294,
        "rotinas_devidas": 302,
        "pct_no_prazo": 97.0,
        "pct_em_atraso": 3.0,
        "media_dias_uteis_pct": 84.7,
        "maior_sequencia_dias_uteis_100": 8,
        "lapsos_dias_uteis": ["09/07", "10/07", "30/07"],
        "semanas": [
            {"rotulo": "01–03/07", "cumprimento_pct": 92.0},
            {"rotulo": "06–10/07", "cumprimento_pct": 58.6},
            {"rotulo": "13–17/07", "cumprimento_pct": 100.0},
            {"rotulo": "20–24/07", "cumprimento_pct": 97.4},
            {"rotulo": "27–31/07", "cumprimento_pct": 78.6},
        ],
        "coletado_em": "21/08/2026",
    },
}


def _coletar_indicador_diario_bordo(db: Session, mes: str) -> dict:
    for nome in _NOMES_DIARIO_BORDO_CANDIDATOS:
        automatico = _extrair_dashboard_externo_por_nome(db, nome, dash_ext.extrair_diario_bordo, mes)
        if automatico.get("tem_dados"):
            return automatico

    dado = _DIARIO_BORDO_POR_MES.get(mes)
    if not dado:
        return {"tem_dados": False, "mes": mes}
    return {**dado, "tem_dados": True, "mes": mes}


# Chaves dos 5 slots nativos (ver dashboards_externos_router.SLOTS) - cada um já
# tem slide dedicado com extração específica acima; qualquer outra chave de
# DashboardExterno é um indicador dinâmico (18/08/2026, ver _coletar_dashboards_extras).
_CHAVES_DASHBOARDS_NATIVOS_MBR = {
    "controle_fefo", "testes_industriais", "farol_shelf_life",
    "recuperacao_shelf", "baixas_operacionais",
}


def _coletar_dashboards_extras(db: Session, chaves_excluir: set = None) -> list:
    """Indicadores dinâmicos (Outros Dashboards > Adicionar Indicador,
    18/08/2026) - qualquer DashboardExterno com conteúdo enviado que não é um
    dos 5 slots nativos entra aqui automaticamente, com extração genérica (ver
    dashboards_externos_extrator.extrair_generico). "dado": None significa que
    a extração não achou tabela nem metadado confiável nesse arquivo - o slide
    correspondente mostra um aviso em vez de inventar conteúdo.

    `chaves_excluir` (20/08/2026): indicadores dinâmicos que já ganharam
    extração/slide DEDICADOS (ex: "Dispersão de Ficha Técnica" - ver
    _extrair_dashboard_externo_por_nome) não devem cair de novo aqui com a
    extração genérica - senão apareceriam duplicados no MBR."""
    chaves_bloqueadas = _CHAVES_DASHBOARDS_NATIVOS_MBR | (chaves_excluir or set())
    registros = (
        db.query(models.DashboardExterno)
        .filter(~models.DashboardExterno.chave.in_(chaves_bloqueadas))
        .all()
    )
    extras = []
    for registro in registros:
        if not registro.html_content:
            continue
        try:
            resultado = dash_ext.extrair_generico(registro.html_content)
        except Exception:
            resultado = None
        extras.append({
            "chave": registro.chave,
            "nome_exibicao": registro.nome_exibicao,
            "enviado_em": registro.enviado_em.strftime("%d/%m/%Y %H:%M") if registro.enviado_em else None,
            "dado": resultado,
        })
    extras.sort(key=lambda item: item["nome_exibicao"].lower())
    return extras


# ---------------------------------------------------------------------------
# Coleta de dados (chama as funções de negócio do Atlas diretamente)
# ---------------------------------------------------------------------------
def _coletar_dados_mbr(db: Session, usuario: models.Usuario, mes: str) -> dict:
    ano_int, mes_int = (int(parte) for parte in mes.split("-"))

    dados = {
        "kpis_inventario": fechamento_router.dashboard_kpis(almoxarifado=None, mes=mes, usuario=usuario, db=db),
        # 22/08/2026 (bug reportado pelo usuário: relatório de julho trazendo dados de
        # agosto "em quase todos os módulos"): dashboard_evolucao_mensal não tem
        # parâmetro de corte - sempre devolve a série completa até o mês mais recente
        # que já tiver QUALQUER registro no banco, pensada pra dashboard ao vivo. Corta
        # pra mes_limite=mes aqui, uma única vez, pra _recorte_periodo/_analise_geral e
        # todo slide que faz [-1]/[-6:] nessa série pararem no mês do relatório, não no
        # mês corrente real (ver docstring de _truncar_serie_mensal).
        "evolucao_inventario": _truncar_serie_mensal(
            fechamento_router.dashboard_evolucao_mensal(almoxarifado=None, usuario=usuario, db=db), mes
        ),
        "comparativo_acuracia": fechamento_router.dashboard_comparativo_acuracia(almoxarifado=None, mes=mes, usuario=usuario, db=db),
        "evolucao_ponderada": _truncar_serie_mensal(
            fechamento_router.dashboard_evolucao_ponderada_mensal(almoxarifado=None, usuario=usuario, db=db), mes
        ),
        # Pareto (concentração de valor) e distribuição por magnitude - já recortados pelo
        # mês do relatório (não o histórico inteiro), pra "mais exemplos do período" serem
        # de fato exemplos DESTE mês (20/08/2026).
        # top_n bumped de 10 pra 20 (22/08/2026, mockup aprovado "Concentração de
        # Risco" v8): o Pareto agora mostra 20 SKUs individuais + uma barra
        # agregada pro resto da cauda, e o Top 10 detalhado ganhou colunas de
        # Almoxarifado/Qtd. Sistêmica/Qtd. Conferência/Dif. que já vêm neste
        # mesmo endpoint — sem chamada nova ao backend. A lista de itens em si
        # já vinha com até 50 (curva[:50] no router), então isso só muda o
        # cálculo de top_n_pct_do_valor, não o que fica disponível pro slide.
        "concentracao_valor": fechamento_router.dashboard_concentracao_valor(almoxarifado=None, mes=mes, top_n=20, usuario=usuario, db=db),
        "distribuicao_magnitude": fechamento_router.dashboard_distribuicao_magnitude(almoxarifado=None, mes=mes, usuario=usuario, db=db),
        # Os 3 modelos (item a item/IAQ/IAP) por almoxarifado, só do mês do
        # relatório (22/08/2026, mockups aprovados v4/v5) - alimenta a tabela
        # "Resultado por Almoxarifado" dos slides de Painel de Inventário e
        # Acurácia Ponderada (IAP/IAQ). O "Prejuízo" em R$ dessas tabelas é
        # derivado de concentracao_valor acima (ver _prejuizo_por_almoxarifado),
        # não vem deste endpoint (que só devolve os 3 percentuais).
        "comparativo_por_almoxarifado": fechamento_router.dashboard_comparativo_por_almoxarifado(mes=mes, usuario=usuario, db=db),
        # Top 10 por faixa de magnitude (22/08/2026, mockup aprovado v8, slide
        # novo "Detalhamento por Faixa de Magnitude") - uma chamada por faixa
        # (só 4), mesma função já usada pelo duplo-clique na tela de Acurácia
        # Ponderada, agora reaproveitada pra imprimir no MBR.
        "magnitude_por_faixa_itens": [
            fechamento_router.dashboard_itens_por_magnitude(faixa_idx=i, almoxarifado=None, mes=mes, usuario=usuario, db=db)
            for i in range(4)
        ],
        "resumo_passivos": baixas_operacionais_router.resumo_executivo(
            ano=ano_int, mes=mes_int, data_inicio=None, data_fim=None, almoxarifado=None, motivo=None, usuario=usuario, db=db
        ),
        # Evolução mensal REAL de Passivos, já cruzada com o Fluxo de Inventário
        # (entradas/saídas/resultado de TODOS os inventários) mês a mês - histórico
        # completo, sem filtro de ano/mês (o slide usa os últimos meses, mesmo padrão
        # de evolucao_inventario/evolucao_ponderada acima) (20/08/2026). Truncado pra
        # mes_limite=mes pelo mesmo motivo dos dois acima (22/08/2026).
        "evolucao_passivos_fluxo": _truncar_serie_mensal(
            baixas_operacionais_router.dashboard_passivos_evolucao_mensal(
                ano=None, mes=None, data_inicio=None, data_fim=None, almoxarifado=None, motivo=None, usuario=usuario, db=db
            ),
            mes,
        ),
        # 22/08/2026: estava chamado com ano=None/mes=None (todo o histórico agregado);
        # a própria docstring da função diz que ela é "a mesma quebra Passivos x
        # Resultado de Inventário do Resumo Executivo, só que por ALMOXARIFADO" - ou
        # seja, deveria usar o mesmo recorte de mês que resumo_passivos acima (ano_int/
        # mes_int), não o acumulado de todos os meses.
        "resultado_por_almoxarifado": baixas_operacionais_router.dashboard_resultado_por_almoxarifado(
            ano=ano_int, mes=mes_int, data_inicio=None, data_fim=None, almoxarifado=None, motivo=None, usuario=usuario, db=db
        ),
        "resumo_shelf_life": shelf_life_mod.calcular_resumo_shelf_life(db, incluir_itens=True, limite_itens=5),
        # Controle de Movimentados (reconciliação diária sistema x físico, origem ==
        # "movimentacao") - indicador PRÓPRIO, separado do Fechamento de Inventário
        # periódico. Ver docstring de movimentados_router.py (20/08/2026).
        "resumo_movimentados": movimentados_router.dashboard_resumo(mes=mes, almoxarifado=None, usuario=usuario, db=db),
        # Quebra por almoxarifado do mês (22/08/2026, mockup aprovado v9) -
        # alimenta a tabela nova "Resultado por Almoxarifado (Movimentados)"
        # no slide Controle de Movimentados.
        "movimentados_por_almoxarifado": movimentados_router.dashboard_por_almoxarifado(mes=mes, usuario=usuario, db=db),
        # Ambas as "evolução mensal" abaixo leem de snapshots persistidos que a própria
        # docstring do router documenta como "o mês corrente sempre reflete o estado
        # mais recente" - ou seja, histórico completo, sem corte. Truncadas pra
        # mes_limite=mes pelo mesmo motivo de evolucao_inventario/evolucao_ponderada
        # acima (22/08/2026). resumo_transferencias fica de fora de propósito - é
        # "totais atuais, sem recorte de período" por design (Transferencia não tem
        # conceito de mês corrente pra filtro, ver docstring do router), não é o bug.
        "evolucao_movimentados": _truncar_serie_mensal(
            movimentados_router.dashboard_evolucao_mensal(almoxarifado=None, usuario=usuario, db=db), mes
        ),
        "resumo_transferencias": movimentados_router.dashboard_transferencias_resumo(usuario=usuario, db=db),
        "evolucao_transferencias": _truncar_serie_mensal(
            movimentados_router.dashboard_transferencias_evolucao_mensal(usuario=usuario, db=db), mes
        ),
        # Dashboards externos (Auditoria > Outros Dashboards) - substituem/complementam
        # números calculados pelo Atlas por dados reais dos arquivos .html que a equipe
        # já mantém em paralelo, pedido do usuário (20/08/2026, ver dashboards_externos_extrator.py).
        # FEFO (atualizado 20/08/2026, pedido do usuário: "use o arquivo HTML pra alimentar
        # a construção do MBR no módulo FEFO, o mesmo tem mais informações e bases de
        # registro" -> confirmado "Auditoria FEFO importada"): a fonte NÃO é mais o dashboard
        # "Controle de FEFO" de Outros Dashboards, e sim a tabela AuditoriaFefo (já usada pelo
        # painel "Auditoria FEFO — histórico importado" da tela FEFO) - ver
        # _extrair_resumo_auditoria_fefo acima.
        "fefo_externo": _extrair_resumo_auditoria_fefo(db, mes),
        "testes_industriais_externo": _extrair_dashboard_externo(db, "testes_industriais", dash_ext.extrair_testes_industriais, mes),
        # Os três abaixo não têm uma dimensão de mês limpa no arquivo de origem (ver
        # docstring de dashboards_externos_extrator.py) - entram no MBR como retrato
        # datado (número real, rotulado com a data/período real do arquivo), não
        # filtrados pelo mês do relatório.
        "farol_shelf_externo": _extrair_dashboard_externo_sem_mes(db, "farol_shelf_life", dash_ext.extrair_farol_shelf),
        "recuperacao_shelf_externo": _extrair_dashboard_externo_sem_mes(db, "recuperacao_shelf", dash_ext.extrair_recuperacao_shelf),
        "baixas_operacionais_externo": _extrair_dashboard_externo_sem_mes(db, "baixas_operacionais", dash_ext.extrair_baixas_operacionais_externo),
        # Dispersão de Ficha Técnica (20/08/2026, pedido do usuário: "Adicione
        # Dispersão de Ficha técnica na apresentação do MBR") - indicador
        # DINÂMICO (cadastrado em Outros Dashboards > Adicionar Indicador),
        # mas com extração/slide dedicados (ver docstring de
        # _extrair_dashboard_externo_por_nome e dashboards_externos_extrator.
        # extrair_dispersao_ficha_tecnica) - por isso é buscado pelo nome de
        # exibição, não por uma chave fixa como os 5 slots nativos.
        "dispersao_ficha_tecnica_externo": _extrair_dashboard_externo_por_nome(
            db, "Dispersão de Ficha Técnica", dash_ext.extrair_dispersao_ficha_tecnica, mes
        ),
    }

    # Mapeamento de Risco - Obsolescência (18/08/2026, pedido do usuário) - ver
    # shelf_life.calcular_mapeamento_risco_obsolescencia.
    dados["mapeamento_risco_obsolescencia"] = shelf_life_mod.calcular_mapeamento_risco_obsolescencia(db)

    # Indicadores dinâmicos (18/08/2026, pedido do usuário: "adicione a opção de
    # adicionar mais indicadores e adicionar automaticamente na construção do
    # MBR") - qualquer DashboardExterno enviado que não é um dos 5 slots nativos
    # acima entra aqui, com extração genérica (ver _coletar_dashboards_extras).
    # Exclui a chave de "Dispersão de Ficha Técnica" (se encontrada acima) pra
    # não duplicar - ela já tem slide dedicado.
    chave_dispersao_ficha = dados["dispersao_ficha_tecnica_externo"].pop("_chave_dashboard_externo", None)
    dados["dashboards_extras"] = _coletar_dashboards_extras(
        db, chaves_excluir={chave_dispersao_ficha} if chave_dispersao_ficha else None
    )

    # Scorecard por almoxarifado (20/08/2026, pedido do usuário) - montado
    # depois do dict principal por padrão, embora não reaproveite nada dele
    # (o Scorecard de Mapeamento de Riscos, que reaproveitava indicadores
    # já coletados acima, foi removido em 22/08/2026 - ver comentário de
    # decisão antes de _slide_controle_movimentados).
    dados["scorecard_inventario_almoxarifado"] = _coletar_scorecard_inventario_almoxarifado(db, usuario, mes)
    dados["diario_bordo"] = _coletar_indicador_diario_bordo(db, mes)

    return dados


# ---------------------------------------------------------------------------
# Narrativa executiva automática (baseada nos números reais coletados acima)
# ---------------------------------------------------------------------------
def _recorte_periodo(d: dict) -> dict:
    """"Recorte do período" (20/08/2026, pedido do usuário): compara o mês do
    relatório com o anterior nas frentes que TÊM série histórica real (mesma
    unidade - pontos percentuais - pra comparação fazer sentido) e aponta o
    maior avanço e a maior involução MoM. Passivos (R$) fica de fora dessa
    comparação de "maior/menor" (unidade diferente, R$ não é comparável a pp
    de acurácia) e ganha uma linha própria, à parte.

    Só usa séries que o Atlas de fato calcula mês a mês - não inventa
    comparação pra frente que não tem histórico (Shelf Life hoje é uma foto,
    não uma série).

    "Acurácia do Inventário" (item-a-item) SAIU da lista de comparáveis em
    22/08/2026 (pedido do usuário: IAP substitui o item-a-item como
    indicador de acurácia no Resumo Executivo) - com o item-a-item fora,
    "Acurácia Ponderada (IAP)" e "Controle de Movimentados" são as duas
    frentes que disputam "maior avanço/involução do mês" aqui."""
    comparaveis = []

    evol_pond = d["evolucao_ponderada"]
    if len(evol_pond) >= 2 and evol_pond[-1].get("variacao_iap_pp") is not None:
        comparaveis.append({"frente": "Acurácia Ponderada (IAP)", "delta_pp": evol_pond[-1]["variacao_iap_pp"]})

    evol_mov = d["evolucao_movimentados"]
    if len(evol_mov) >= 2 and evol_mov[-1].get("pct_acuracia") is not None and evol_mov[-2].get("pct_acuracia") is not None:
        comparaveis.append({"frente": "Controle de Movimentados", "delta_pp": round(evol_mov[-1]["pct_acuracia"] - evol_mov[-2]["pct_acuracia"], 2)})

    maior_avanco = max(comparaveis, key=lambda c: c["delta_pp"]) if comparaveis else None
    if maior_avanco and maior_avanco["delta_pp"] <= 0:
        maior_avanco = None
    maior_involucao = min(comparaveis, key=lambda c: c["delta_pp"]) if comparaveis else None
    if maior_involucao and maior_involucao["delta_pp"] >= 0:
        maior_involucao = None
    # evita apontar a MESMA frente como avanço e involução (só acontece com 1 item comparável)
    if maior_avanco and maior_involucao and maior_avanco["frente"] == maior_involucao["frente"]:
        maior_involucao = None

    variacao_passivos_valor = None
    evol_passivos = d["evolucao_passivos_fluxo"]
    if len(evol_passivos) >= 2:
        variacao_passivos_valor = round(evol_passivos[-1]["valor"] - evol_passivos[-2]["valor"], 2)

    return {
        "comparaveis": comparaveis,
        "maior_avanco": maior_avanco,
        "maior_involucao": maior_involucao,
        "variacao_passivos_valor": variacao_passivos_valor,
    }


def _analise_geral(d: dict):
    """Monta as 3 colunas do resumo executivo (Avanços / Atenções /
    Decisões) e a mensagem central, tudo derivado dos limiares reais - nada
    de texto fixo/genérico independente dos números do mês."""
    avancos, atencoes, decisoes = [], [], []

    kpis_inv = d["kpis_inventario"]
    comp = d["comparativo_acuracia"]
    evol_pond = d["evolucao_ponderada"]
    passivos = d["resumo_passivos"]["passivos"]
    resultado_inv = d["resumo_passivos"]["resultado_inventario"]
    baixas_pacote = (d.get("baixas_operacionais_externo") or {}).get("resumo") or {}
    shelf = d["resumo_shelf_life"]
    movimentados = d["resumo_movimentados"]
    evol_mov = d["evolucao_movimentados"]

    # Acurácia Ponderada / IAP como indicador oficial no Resumo Executivo
    # (22/08/2026, pedido do usuário: "substitua a lógica de inventário Item
    # a Item e trocar pela análise IAP. Mostrando curva de aprendizado e
    # melhoria contínua do processo" - só neste slide, ver decisão registrada
    # na conversa: Painel de Inventário/Scorecard continuam com item-a-item).
    # Usa evolucao_ponderada (série real, já truncada no mês do relatório)
    # pra falar de TENDÊNCIA no período, não só do nível do mês.
    if evol_pond and evol_pond[-1].get("iap_pct") is not None:
        iap_atual = evol_pond[-1]["iap_pct"]
        primeiro = evol_pond[0]
        if len(evol_pond) >= 2 and primeiro.get("iap_pct") is not None:
            delta_periodo = round(iap_atual - primeiro["iap_pct"], 2)
            meses_com_delta = [e for e in evol_pond[1:] if e.get("variacao_iap_pp") is not None]
            melhora_meses = sum(1 for e in meses_com_delta if e["variacao_iap_pp"] > 0)
            piora_meses = sum(1 for e in meses_com_delta if e["variacao_iap_pp"] < 0)
            if delta_periodo > 0 and melhora_meses >= piora_meses:
                texto_iap = (
                    f"IAP (Acurácia por Valor) sobe de {_fmt_pct(primeiro['iap_pct'])} em {_nome_mes(primeiro['mes'])} para "
                    f"{_fmt_pct(iap_atual)} em {_nome_mes(evol_pond[-1]['mes'])} — +{_fmt_pct(delta_periodo)} no período"
                )
                if melhora_meses > piora_meses:
                    texto_iap += (
                        f", melhora em {melhora_meses} de {len(meses_com_delta)} mês(es) monitorado(s). "
                        "Curva de aprendizado consistente, não um pico isolado."
                    )
                else:
                    texto_iap += "."
                avancos.append(texto_iap)
            else:
                atencoes.append(
                    f"IAP (Acurácia por Valor) em {_fmt_pct(iap_atual)}, variação de {_fmt_pct(delta_periodo)} desde "
                    f"{_nome_mes(primeiro['mes'])} — sem curva de aprendizado consistente ainda ({piora_meses} de "
                    f"{len(meses_com_delta)} mês(es) com queda)."
                )
        else:
            label_iap, _ = _status_maior_melhor(iap_atual, *_LIMIARES["acuracia"])
            (avancos if label_iap == "Em avanço" else atencoes).append(
                f"IAP (Acurácia por Valor) em {_fmt_pct(iap_atual)} neste mês (ainda sem histórico anterior suficiente pra falar de tendência)."
            )

    if comp.get("item_a_item_pct") is not None:
        atencoes.append(
            f"Item a item em {_fmt_pct(comp['item_a_item_pct'])} funciona como termômetro de disciplina de apontamento — "
            "não é mais o indicador oficial de acurácia (ver Acurácia Ponderada)."
        )

    gap = comp.get("gap_item_vs_iap_pp")
    if gap is not None and abs(gap) >= 3:
        atencoes.append(
            f"A leitura ponderada por valor (IAP) muda em {_fmt_pct(abs(gap))} a foto do item-a-item — "
            "sinal de que o risco financeiro está concentrado em poucos SKUs, não de que o inventário piorou."
        )

    # Baixas por Pacote / Baixas Operacionais (22/08/2026, pedido do usuário:
    # "Troque o indicador de Mapeamento de passivos por Baixas por pacote
    # (Baixas operacionais)") - fonte é o dashboard externo (mesmo dado do
    # KPI do Resumo Executivo e do slide "Baixas Operacionais (Controle
    # Paralelo)"), não o Mapeamento de Passivos nativo (que continua existindo
    # à parte, em sua própria seção/scorecard).
    if baixas_pacote.get("prejuizo_total"):
        atencoes.append(
            f"{_fmt_moeda(baixas_pacote['prejuizo_total'])} em baixas por pacote no período "
            f"({_fmt_pct(baixas_pacote.get('pct_concentrado'))} concentrado em {baixas_pacote.get('motivo_concentrado', '—')})."
        )
    elif passivos["valor"]:
        if passivos["valor"] > 0:
            atencoes.append(f"{_fmt_moeda(passivos['valor'])} em passivos aprovados mapeados no período ({_fmt_num(passivos['quantidade'])} baixas).")
        else:
            avancos.append("Nenhum passivo residual aprovado sem mapeamento no período.")
    else:
        avancos.append("Nenhum passivo residual aprovado sem mapeamento no período.")

    if resultado_inv["resultado_valor"] is not None:
        if resultado_inv["resultado_valor"] >= 0:
            avancos.append(f"Resultado de inventário acumulado positivo em {_fmt_moeda(resultado_inv['resultado_valor'])}.")
        else:
            atencoes.append(f"Resultado de inventário acumulado negativo em {_fmt_moeda(resultado_inv['resultado_valor'])}.")

    if shelf["valor_total"]:
        decisoes.append(f"Priorizar consumo/transferência dos {_fmt_moeda(shelf['valor_total'])} em lotes vencidos ou a vencer em até 90 dias ({_fmt_num(shelf['total_lotes_em_risco'])} lotes).")
    else:
        avancos.append("Sem valor relevante de estoque em risco de validade neste recorte.")

    # Controle de Movimentados (item pedido explicitamente pelo usuário, 20/08/2026):
    # destaca como AÇÃO com impacto mensurável na acurácia desde a implantação - primeiro
    # mês com dado no snapshot mensal (ResumoMovimentacaoMensal) é lido como "implantação".
    if evol_mov and evol_mov[0].get("pct_acuracia") is not None and movimentados.get("pct_acuracia") is not None:
        primeiro = evol_mov[0]
        delta_implantacao = round(movimentados["pct_acuracia"] - primeiro["pct_acuracia"], 2)
        if len(evol_mov) >= 2 and delta_implantacao > 0:
            avancos.append(
                f"Controle de Movimentados: acurácia da reconciliação diária subiu de {_fmt_pct(primeiro['pct_acuracia'])} "
                f"(em {_nome_mes(primeiro['mes'])}, primeiro mês monitorado) para {_fmt_pct(movimentados['pct_acuracia'])} agora "
                f"(+{_fmt_pct(delta_implantacao)}) — ação com impacto direto e mensurável na acurácia."
            )
        elif movimentados.get("pct_acuracia") is not None:
            atencoes.append(
                f"Controle de Movimentados: acurácia da reconciliação diária em {_fmt_pct(movimentados['pct_acuracia'])} "
                f"desde a implantação (em {_nome_mes(primeiro['mes'])}) — ainda sem ganho consistente frente ao ponto de partida."
            )
    elif movimentados.get("itens_analisados"):
        avancos.append(f"Controle de Movimentados ativo: {_fmt_num(movimentados['itens_analisados'])} item(ns) reconciliado(s) neste mês.")

    # FEFO (atualizado 20/08/2026, pedido do usuário: trocar a fonte pra "Auditoria FEFO
    # importada", ver _extrair_resumo_auditoria_fefo) - histórico importado na própria tela
    # FEFO (Excel diário ou dashboard HTML consolidado do André), com mais detalhe por
    # registro (lote movimentado, validade, lote mais antigo disponível) do que o dashboard
    # "Controle de FEFO" usado antes - por isso entra como fato normal em avanços/atenções,
    # com o mesmo limiar (_LIMIARES["fefo_quebra_pct"]) já usado no restante do relatório.
    fefo = d["fefo_externo"]
    if fefo.get("tem_dados"):
        label_fefo, cor_fefo = _status_menor_melhor(fefo["taxa_quebra_pct"], *_LIMIARES["fefo_quebra_pct"])
        texto_fefo = (
            f"FEFO: taxa de quebra em {_fmt_pct(fefo['taxa_quebra_pct'])} neste mês "
            f"({_fmt_num(fefo['total_quebras'])} de {_fmt_num(fefo['total_auditaveis'])} movimentos auditáveis, "
            f"dados da Auditoria FEFO importada)."
        )
        if label_fefo == "Em avanço":
            avancos.append(texto_fefo)
        elif label_fefo == "Atenção":
            atencoes.append(texto_fefo)
        else:
            atencoes.append(texto_fefo + " — acima do limiar aceitável.")
    elif not fefo.get("enviado"):
        decisoes.append(
            "Auditoria FEFO importada ainda não tem nenhum histórico importado na tela FEFO — sem esse arquivo, o "
            "FEFO não entra neste relatório com dado real."
        )
    else:
        decisoes.append(f"Auditoria FEFO importada: nenhum movimento auditável importado para {_nome_mes(fefo.get('mes', ''))}.")

    if not decisoes:
        decisoes.append("Manter a cadência atual de fechamento e monitoramento — sem decisão crítica pendente neste recorte.")

    recorte = _recorte_periodo(d)

    if len(atencoes) == 0:
        mensagem = (
            "A operação está sob controle em todas as frentes mapeadas neste mês — inventário auditado, passivos "
            "mapeados, sem decisão crítica pendente. O ganho aqui é de visibilidade: o mesmo controle que antes "
            "vivia em planilhas isoladas agora é automático, a cada fechamento."
        )
    else:
        mensagem = (
            f"O mapeamento da operação mostra controle real do estoque, não percepção: {len(avancos)} frente(s) em "
            f"avanço e {len(atencoes)} de atenção neste mês, com as ações já priorizadas ao lado — inventário, "
            "passivos, validade e movimentação lidos juntos, pra agir antes que a atenção vire perda."
        )

    return {"avancos": avancos, "atencoes": atencoes, "decisoes": decisoes, "mensagem": mensagem, "recorte_periodo": recorte}


def _linha_scorecard(frente, status_label, status_cor, leitura, proximo_passo):
    return {"frente": frente, "status_label": status_label, "status_cor": status_cor, "leitura": leitura, "proximo_passo": proximo_passo}


def _montar_scorecard(d: dict):
    kpis_inv = d["kpis_inventario"]
    passivos = d["resumo_passivos"]["passivos"]
    shelf = d["resumo_shelf_life"]
    movimentados = d["resumo_movimentados"]

    label_acuracia, cor_acuracia = _status_maior_melhor(kpis_inv["acuracia_geral_pct"], *_LIMIARES["acuracia"])
    label_acuracia, cor_acuracia = _aplica_avanco_operacional(kpis_inv["acuracia_geral_pct"], label_acuracia, cor_acuracia)
    linha_inventario = _linha_scorecard(
        "Painel de Inventário", label_acuracia, cor_acuracia,
        f"Acurácia geral {_fmt_pct(kpis_inv['acuracia_geral_pct'])} sobre {_fmt_num(kpis_inv['total_itens'])} itens avaliados.",
        "Investigar SKUs recorrentes." if kpis_inv["total_divergentes"] else "Manter cadência de fechamento atual.",
    )

    # Baixas Operacionais (Pacote) no lugar de Mapeamento de Passivos (22/08/2026,
    # pedido do usuário: "troque o indicador de Mapeamento de passivos por
    # Baixas por pacote (Baixas operacionais)... atualize todos os Scorecards" -
    # mesma fonte do KPI e do slide dedicado "Baixas Operacionais (Controle
    # Paralelo)"). Cai pro Mapeamento de Passivos nativo se o dashboard externo
    # ainda não foi enviado, pra não quebrar o scorecard.
    baixas_pacote = (d.get("baixas_operacionais_externo") or {}).get("resumo") or {}
    if baixas_pacote.get("prejuizo_total") is not None:
        label_baixas = "Atenção" if baixas_pacote["prejuizo_total"] else "Em avanço"
        cor_baixas = COR_ATENCAO if baixas_pacote["prejuizo_total"] else COR_SUCESSO
        linha_passivos = _linha_scorecard(
            "Baixas Operacionais (Pacote)", label_baixas, cor_baixas,
            f"{_fmt_moeda(baixas_pacote['prejuizo_total'])} no período — "
            f"{_fmt_pct(baixas_pacote.get('pct_concentrado'))} concentrado em {baixas_pacote.get('motivo_concentrado', '—')}.",
            "Investigar causa raiz da concentração antes do próximo fechamento." if baixas_pacote["prejuizo_total"] else "Sem baixa relevante — manter monitoramento.",
        )
    else:
        label_passivos = "Atenção" if passivos["valor"] else "Em avanço"
        cor_passivos = COR_ATENCAO if passivos["valor"] else COR_SUCESSO
        linha_passivos = _linha_scorecard(
            "Mapeamento de Passivos", label_passivos, cor_passivos,
            f"{_fmt_moeda(passivos['valor'])} em passivos aprovados mapeados no período.",
            "Classificar pendências como ajuste de processo ou perda real." if passivos["valor"] else "Sem pendência — manter monitoramento.",
        )

    label_shelf = "Atenção" if shelf["total_lotes_em_risco"] else "Em avanço"
    cor_shelf = COR_ATENCAO if shelf["total_lotes_em_risco"] else COR_SUCESSO
    linha_shelf = _linha_scorecard(
        "Shelf Life", label_shelf, cor_shelf,
        f"{_fmt_num(shelf['total_lotes_em_risco'])} lote(s) em risco de validade, {_fmt_moeda(shelf['valor_total'])} em valor exposto.",
        "Priorizar consumo/transferência dos lotes vencidos ou a vencer em 30 dias." if shelf["total_lotes_em_risco"] else "Sem lote em risco — manter monitoramento.",
    )

    # Controle de Movimentados (20/08/2026) - substitui a antiga linha "Movimentados & FEFO":
    # a taxa de quebra de FEFO saiu do scorecard (ver _analise_geral) porque hoje ela compara
    # a transferência com o estoque de lote ATUAL, não com uma leitura de disponibilidade
    # tirada no momento exato da transferência - não é um número fechado o bastante pra virar
    # status de scorecard. Essa linha agora reflete a reconciliação diária de verdade.
    # "Em avanço" por tendência consistente desde o primeiro mês monitorado
    # (22/08/2026, pedido do usuário: "atualize... e destaque avanços
    # alcançados" - mesma regra de _analise_geral: sem essa checagem, o
    # limiar absoluto (_LIMIARES["acuracia"]) classificava Movimentados como
    # "Atenção" mesmo com alta real e ininterrupta em todos os meses
    # monitorados, o que contradizia o texto de Avanços do próprio Resumo
    # Executivo dizendo que a frente avançou).
    evol_mov_sc = d["evolucao_movimentados"]
    tendencia_positiva_mov = (
        len(evol_mov_sc) >= 2
        and evol_mov_sc[0].get("pct_acuracia") is not None
        and movimentados.get("pct_acuracia") is not None
        and movimentados["pct_acuracia"] > evol_mov_sc[0]["pct_acuracia"]
    )
    if tendencia_positiva_mov:
        label_mov, cor_mov = "Em avanço", COR_SUCESSO
        leitura_mov = (
            f"Acurácia da reconciliação diária em {_fmt_pct(movimentados.get('pct_acuracia'))} — alta consistente desde "
            f"{_nome_mes(evol_mov_sc[0]['mes'])} (+{_fmt_pct(movimentados['pct_acuracia'] - evol_mov_sc[0]['pct_acuracia'])})."
        )
    else:
        label_mov, cor_mov = _status_maior_melhor(movimentados.get("pct_acuracia"), *_LIMIARES["acuracia"])
        leitura_mov = f"Acurácia da reconciliação diária em {_fmt_pct(movimentados.get('pct_acuracia'))} sobre {_fmt_num(movimentados.get('itens_analisados'))} item(ns) analisado(s) no mês."
    label_mov, cor_mov = _aplica_avanco_operacional(movimentados.get("pct_acuracia"), label_mov, cor_mov)
    linha_movimentados = _linha_scorecard(
        "Controle de Movimentados", label_mov, cor_mov, leitura_mov,
        "Investigar itens com divergência não resolvida." if movimentados.get("itens_com_divergencia") else "Manter cadência de conferência diária.",
    )

    # IAP pondera cada divergência pelo valor financeiro do SKU - um IAP alto
    # (>= 95%, regra de negócio pedida pelo usuário 22/08/2026) já é, por si
    # só, um avanço operacional; o gap vs. a leitura item-a-item é uma
    # informação complementar (mostra o quanto a ponderação por valor
    # "esconde" da dispersão item a item), não um motivo pra rebaixar esse
    # status - ver _aplica_avanco_operacional.
    comp = d["comparativo_acuracia"]
    gap = comp.get("gap_item_vs_iap_pp")
    label_gap = "Atenção" if (gap is not None and abs(gap) >= 3) else "Em avanço"
    cor_gap = COR_ATENCAO if label_gap == "Atenção" else COR_SUCESSO
    label_gap, cor_gap = _aplica_avanco_operacional(comp.get("iap_pct"), label_gap, cor_gap)
    linha_ponderada = _linha_scorecard(
        "Acurácia Ponderada", label_gap, cor_gap,
        f"IAP (ponderado por valor) em {_fmt_pct(comp.get('iap_pct'))}, {_fmt_pct(abs(gap)) if gap is not None else '—'} de distância do item-a-item.",
        "Priorizar SKUs de maior impacto financeiro na correção." if label_gap == "Atenção" else "Distorção sob controle — manter leitura mensal.",
    )

    # Mapeamento de Risco - Obsolescência (18/08/2026, pedido do usuário) - cruza
    # Shelf Life (validade) com Movimentados (giro) - ver
    # shelf_life.calcular_mapeamento_risco_obsolescencia.
    risco_obs = d["mapeamento_risco_obsolescencia"]
    if risco_obs.get("quantidade_criticos"):
        label_risco_obs, cor_risco_obs = "Crítico", COR_ERRO
    elif risco_obs.get("quantidade_itens"):
        label_risco_obs, cor_risco_obs = "Atenção", COR_ATENCAO
    else:
        label_risco_obs, cor_risco_obs = "Em avanço", COR_SUCESSO
    linha_risco_obs = _linha_scorecard(
        "Mapeamento de Risco (Obsolescência)", label_risco_obs, cor_risco_obs,
        (f"{_fmt_num(risco_obs.get('quantidade_itens'))} item(ns) em risco de obsolescência, "
         f"{_fmt_moeda(risco_obs.get('valor_total_risco'))} em valor exposto.")
        if risco_obs.get("quantidade_itens") else "Nenhum item combina validade próxima com giro insuficiente neste recorte.",
        "Priorizar ação comercial/promocional nos itens críticos (giro zero) antes que vençam." if risco_obs.get("quantidade_criticos")
        else ("Monitorar itens em atenção — giro insuficiente pra escoar no ritmo atual." if risco_obs.get("quantidade_itens")
              else "Sem risco de obsolescência — manter monitoramento."),
    )

    return [linha_inventario, linha_passivos, linha_ponderada, linha_shelf, linha_risco_obs, linha_movimentados]


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def _slide_capa(prs: Presentation, mes_label: str):
    slide = _slide_em_branco(prs)
    _fundo(slide, AZUL_INSTITUCIONAL)
    _texto(slide, 1.0, 2.55, 11.3, 0.3, "MÁGIO CHOCOLATES", tamanho=14, negrito=True, cor=AZUL_CLARO)
    _texto(slide, 1.0, 2.95, 11.3, 1.1, "MBR Executivo - Controle", tamanho=40, negrito=True,
           cor=BRANCO, fonte=FONTE_TITULO)
    _texto(slide, 1.0, 3.80, 11.3, 0.5, "Inteligência e controle operacional · Atlas", tamanho=18, cor=AZUL_CLARO)
    agora = datetime.now().strftime("%d/%m/%Y %H:%M")
    _texto(slide, 1.0, 6.55, 11.3, 0.4, f"{mes_label}  ·  Gerado automaticamente pelo Atlas em {agora}",
           tamanho=13, cor=OFF_WHITE)
    return slide


def _slide_abertura_secao(prs: Presentation, mes_label: str, pagina: int, numero_secao: int, titulo_secao: str,
                           descricao_secao: str, itens_secao: list):
    """Capa de seção (pedido do usuário, 18/08/2026: "traga uma visão
    detalhada modulando os grupos de relatório" nas 7 categorias definidas
    com ele) - fundo escuro (mesma cor da capa), pra marcar visualmente a
    virada de um grupo temático pro outro, sem usar accent stripe/linha
    decorativa (ver diretrizes de design da skill de pptx - evitadas de
    propósito)."""
    slide = _slide_em_branco(prs)
    _fundo(slide, AZUL_INSTITUCIONAL)
    # "DE 4" (era "DE 5" até 22/08/2026 - a fusão de Seção 3 nativa + Seção 4
    # "Outros" numa seção só, pedido não-negociável da usuária, reduziu de
    # 5 pra 4 seções - ver comentário de decisão dentro de montar_pptx_mbr).
    _texto(slide, MARGEM_IN, 0.32, 4.0, 0.3, f"SEÇÃO {numero_secao} DE 4", tamanho=11, negrito=True, cor=AZUL_CLARO)
    _texto(slide, LARGURA_IN - 2.9, 0.32, 2.4, 0.3, mes_label.upper(), tamanho=11, negrito=True,
           cor=AZUL_CLARO, alinhamento=PP_ALIGN.RIGHT)
    _texto(slide, MARGEM_IN, 1.15, LARGURA_IN - 2 * MARGEM_IN, 1.0, titulo_secao, tamanho=34, negrito=True,
           cor=BRANCO, fonte=FONTE_TITULO)
    _texto(slide, MARGEM_IN, 2.15, LARGURA_IN - 2 * MARGEM_IN, 0.6, descricao_secao, tamanho=15, cor=AZUL_CLARO)

    y = 3.15
    _texto(slide, MARGEM_IN, y, 4.0, 0.28, "NESTA SEÇÃO", tamanho=11, negrito=True, cor=AZUL_CLARO)
    y += 0.4
    for item in itens_secao:
        _texto(slide, MARGEM_IN, y, 0.24, 0.3, "•", tamanho=14, negrito=True, cor=AZUL_CLARO)
        _texto(slide, MARGEM_IN + 0.28, y, LARGURA_IN - 2 * MARGEM_IN - 0.28, 0.3, item, tamanho=13, cor=BRANCO)
        y += 0.4

    _texto(slide, LARGURA_IN - 2.9, ALTURA_IN - 0.42, 2.4, 0.3, f"{pagina:02d}", tamanho=10,
           cor=AZUL_CLARO, alinhamento=PP_ALIGN.RIGHT)
    return slide


def _slide_resumo_executivo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Resumo Executivo", mes_label, pagina, "Controle operacional do estoque em números — o que mudou, o que exige decisão")

    kpis_inv = d["kpis_inventario"]
    passivos = d["resumo_passivos"]["passivos"]
    shelf = d["resumo_shelf_life"]
    movimentados = d["resumo_movimentados"]
    comp = d["comparativo_acuracia"]
    evol_pond = d["evolucao_ponderada"]
    farol_shelf = d.get("farol_shelf_externo") or {}
    baixas_pacote = (d.get("baixas_operacionais_externo") or {}).get("resumo") or {}
    label_mov, cor_mov = _status_maior_melhor(movimentados.get("pct_acuracia"), *_LIMIARES["acuracia"])

    # KPI 1: IAP (Acurácia Ponderada) no lugar do item-a-item (22/08/2026,
    # pedido do usuário: "substitua a lógica de inventário Item a Item e
    # trocar pela análise IAP" - SÓ neste slide, ver decisão registrada na
    # conversa; Painel de Inventário e Scorecard continuam com item-a-item).
    label_iap, cor_iap = _status_maior_melhor(comp.get("iap_pct"), *_LIMIARES["acuracia"])
    variacao_iap_mes = evol_pond[-1].get("variacao_iap_pp") if evol_pond else None
    if variacao_iap_mes is not None:
        contexto_iap = f"{'▲' if variacao_iap_mes > 0 else '▼' if variacao_iap_mes < 0 else '▬'} {_fmt_pct(abs(variacao_iap_mes))} vs. mês anterior"
        cor_contexto_iap = COR_SUCESSO if variacao_iap_mes > 0 else (COR_ERRO if variacao_iap_mes < 0 else CINZA_TEXTO)
    else:
        contexto_iap, cor_contexto_iap = label_iap, cor_iap

    # KPI 2: Baixas por Pacote (Baixas Operacionais) no lugar de Passivos
    # Mapeados (22/08/2026, pedido do usuário) - fonte é o dashboard externo
    # "Dashboard Baixas Operacionais" (mesmo dado do KPI e do slide dedicado),
    # que em 22/08/2026 passou a ser a ÚNICA fonte pra esse assunto no MBR -
    # o slide nativo "Mapeamento de Passivos" foi removido (ver comentário de
    # decisão antes de _slide_controle_movimentados). O fallback
    # abaixo (`passivos["valor"]`, de d["resumo_passivos"]) continua existindo
    # só pra este KPI específico não quebrar se o dashboard externo ainda não
    # tiver sido enviado - não é mais um slide dedicado nativo no relatório.
    #
    # Rótulo encurtado pra "Baixas por Pacote" (22/08/2026, bug reportado pela
    # usuária no modelo manual: com o texto completo "Baixas por Pacote
    # (Baixas Operacionais)", 40 caracteres, o rótulo quebrava em 2 linhas
    # dentro do cartão de KPI e invadia a caixa de contexto logo abaixo, cuja
    # posição é fixa - mesma classe de bug já corrigida em títulos de slide
    # (Faixas, Baixas Operacionais), agora também no cartão de KPI. A fonte
    # ("Dashboard Baixas Operacionais") já está clara pelo resto do relatório;
    # não precisa repetir aqui dentro do espaço apertado do cartão.
    if baixas_pacote.get("prejuizo_total") is not None:
        kpi_baixas = {
            "valor": _fmt_moeda(baixas_pacote["prejuizo_total"]), "rotulo": "Baixas por Pacote",
            "cor": COR_ERRO, "contexto": f"{_fmt_pct(baixas_pacote.get('pct_concentrado'))} em {baixas_pacote.get('motivo_concentrado', '—')}",
        }
    else:
        # dashboard externo não enviado ainda - cai pro cálculo nativo de passivos (d["resumo_passivos"]), sem quebrar o slide.
        kpi_baixas = {
            "valor": _fmt_moeda(passivos["valor"]), "rotulo": "Passivos Mapeados",
            "cor": AZUL_INSTITUCIONAL, "contexto": f"{_fmt_num(passivos['quantidade'])} baixas aprovadas",
        }

    # KPI 3: Valor em Risco de Validade a partir do Farol de Shelf externo
    # (22/08/2026, pedido do usuário: "traga também o valor em risco baseado
    # no relatório de Farol de Shelf") - antes vinha de resumo_shelf_life
    # (cálculo nativo do Atlas); o Farol de Shelf externo (mesmo valor do
    # slide dedicado "Farol de Shelf-Life") é, desde 22/08/2026, a ÚNICA
    # fonte pra esse assunto no MBR - o slide nativo "Shelf Life" foi
    # removido (ver mesmo comentário de decisão citado acima).
    #
    # Rótulo encurtado pra "Valor em Risco de Validade" (22/08/2026, mesmo bug
    # de rótulo em 2 linhas reportado pela usuária, ver comentário do KPI 2
    # acima) - com o sufixo "(Farol de Shelf)" o texto passava de 40
    # caracteres e quebrava em 2 linhas, invadindo o contexto abaixo. Sem o
    # sufixo, o rótulo fica igual ao do caminho nativo (fallback abaixo) -
    # tudo bem, é o mesmo conceito de indicador nos dois casos, só a fonte
    # do dado muda.
    if farol_shelf.get("perda_potencial_total") is not None:
        qtd_farol = farol_shelf.get("qtd_lotes") or {}
        total_lotes_farol = sum(v for v in qtd_farol.values() if isinstance(v, int))
        kpi_risco_validade = {
            "valor": _fmt_moeda(farol_shelf["perda_potencial_total"]), "rotulo": "Valor em Risco de Validade",
            "cor": COR_ATENCAO, "contexto": f"{_fmt_num(total_lotes_farol)} lotes · {_fmt_num(qtd_farol.get('vencidos'))} já vencidos",
        }
    else:
        # dashboard externo não enviado ainda - cai pro cálculo nativo, sem quebrar o slide.
        kpi_risco_validade = {
            "valor": _fmt_moeda(shelf["valor_total"]), "rotulo": "Valor em Risco de Validade",
            "cor": AZUL_INSTITUCIONAL, "contexto": f"{_fmt_num(shelf['total_lotes_em_risco'])} lotes",
        }

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_pct(comp.get("iap_pct")), "rotulo": "IAP — Acurácia por Valor", "cor": cor_iap, "contexto": contexto_iap, "cor_contexto": cor_contexto_iap},
        kpi_baixas,
        kpi_risco_validade,
        {"valor": _fmt_pct(movimentados.get("pct_acuracia")), "rotulo": "Controle de Movimentados", "cor": AZUL_INSTITUCIONAL, "contexto": label_mov, "cor_contexto": cor_mov},
    ], altura=0.85)

    analise = _analise_geral(d)

    # "Recorte do Período" (20/08/2026, pedido do usuário): uma linha de destaque logo
    # abaixo dos KPIs com o maior avanço e a maior involução MoM entre as frentes que têm
    # série histórica comparável (mesma unidade - pp) - ver _recorte_periodo.
    recorte = analise["recorte_periodo"]
    y_recorte = 2.95
    altura_recorte = 0.42
    partes_recorte = []
    if recorte["maior_avanco"]:
        partes_recorte.append(("▲ Maior avanço do mês: ", CINZA_TEXTO, False))
        partes_recorte.append((f"{recorte['maior_avanco']['frente']} (+{_fmt_pct(recorte['maior_avanco']['delta_pp'])})", COR_SUCESSO, True))
    if recorte["maior_involucao"]:
        if partes_recorte:
            partes_recorte.append(("     ", CINZA_TEXTO, False))
        partes_recorte.append(("▼ Maior involução do mês: ", CINZA_TEXTO, False))
        partes_recorte.append((f"{recorte['maior_involucao']['frente']} ({_fmt_pct(recorte['maior_involucao']['delta_pp'])})", COR_ERRO, True))
    if recorte["variacao_passivos_valor"] is not None and recorte["variacao_passivos_valor"] != 0:
        if partes_recorte:
            partes_recorte.append(("     ", CINZA_TEXTO, False))
        sinal = "+" if recorte["variacao_passivos_valor"] > 0 else ""
        cor_passivo_var = COR_ERRO if recorte["variacao_passivos_valor"] > 0 else COR_SUCESSO
        # "aprovados (fechamento nativo)" pra não confundir com o KPI "Baixas por Pacote"
        # acima, que é outra fonte (dashboard externo) - 22/08/2026.
        partes_recorte.append(("Passivos aprovados (fechamento nativo): ", CINZA_TEXTO, False))
        partes_recorte.append((f"{sinal}{_fmt_moeda(recorte['variacao_passivos_valor'])} vs. mês anterior", cor_passivo_var, True))

    if partes_recorte:
        _retangulo(slide, MARGEM_IN, y_recorte, LARGURA_IN - 2 * MARGEM_IN, altura_recorte, cor_fill=OFF_WHITE, raio=0.08)
        tf_recorte = slide.shapes.add_textbox(Inches(MARGEM_IN + 0.18), Inches(y_recorte), Inches(LARGURA_IN - 2 * MARGEM_IN - 0.36), Inches(altura_recorte)).text_frame
        tf_recorte.word_wrap = True
        tf_recorte.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_recorte = tf_recorte.paragraphs[0]
        for texto_parte, cor_parte, negrito_parte in partes_recorte:
            run = p_recorte.add_run()
            run.text = texto_parte
            run.font.size = Pt(11.5)
            run.font.name = FONTE_TEXTO
            run.font.bold = negrito_parte
            run.font.color.rgb = cor_parte
        y_colunas = y_recorte + altura_recorte + 0.18
    else:
        y_colunas = 3.10

    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 2 * 0.25) / 3
    largura_texto_col = largura_col - 0.22

    # Quantos itens cada coluna acaba tendo varia todo mês (mais ou menos
    # avanços/atenções/decisões conforme os números do recorte) - por isso o
    # tamanho de fonte da lista e a posição da caixa "Mensagem Central" são
    # calculados a partir da altura que o conteúdo realmente vai ocupar
    # (_linhas_estimadas), não de um layout fixo que só foi testado com um
    # mês de exemplo. Tenta fontes menores até caber no espaço reservado
    # pras 3 colunas antes da caixa de mensagem; no pior caso (muitos itens
    # longos), fica no menor tamanho e aceita ficar mais apertado.
    y_zona_segura_fim = ALTURA_IN - 0.55  # reserva o rodapé (número de página)
    # Altura reservada calculada a partir do texto real da mensagem (20/08/2026: antes era
    # um valor fixo de 1.05 que não escalava com o tamanho da mensagem - truncava com
    # reticências em meses normais, não só no pior caso).
    altura_reservada_mensagem = max(1.05, _altura_necessaria_caixa_leitura(analise["mensagem"], LARGURA_IN - 2 * MARGEM_IN, 12))
    gap_antes_mensagem = 0.15
    orcamento_colunas = y_zona_segura_fim - y_colunas - 0.36 - gap_antes_mensagem - altura_reservada_mensagem

    tamanho_lista = 11.5
    espaco_linha = 0.24
    for tentativa in (11.5, 10.5, 9.5):
        espaco_tentativa = 0.24 * (tentativa / 11.5)
        maior_altura = max(
            sum(espaco_tentativa * _linhas_estimadas(item, largura_texto_col, tentativa) + 0.10 for item in analise[chave])
            for chave in ("avancos", "atencoes", "decisoes")
        )
        tamanho_lista, espaco_linha = tentativa, espaco_tentativa
        if maior_altura <= orcamento_colunas:
            break

    altura_lista_usada = max(
        sum(espaco_linha * _linhas_estimadas(item, largura_texto_col, tamanho_lista) + 0.10 for item in analise[chave])
        for chave in ("avancos", "atencoes", "decisoes")
    )
    altura_lista_usada = min(altura_lista_usada, orcamento_colunas)

    _texto(slide, MARGEM_IN, y_colunas, largura_col, 0.28, "AVANÇOS", tamanho=13, negrito=True, cor=COR_SUCESSO)
    _lista_com_marcadores(slide, MARGEM_IN, y_colunas + 0.36, largura_col, altura_lista_usada, analise["avancos"],
                           cor_marcador=COR_SUCESSO, tamanho=tamanho_lista, espaco_linha=espaco_linha)

    x2 = MARGEM_IN + largura_col + 0.25
    _texto(slide, x2, y_colunas, largura_col, 0.28, "ATENÇÕES", tamanho=13, negrito=True, cor=COR_ATENCAO)
    _lista_com_marcadores(slide, x2, y_colunas + 0.36, largura_col, altura_lista_usada, analise["atencoes"],
                           cor_marcador=COR_ATENCAO, tamanho=tamanho_lista, espaco_linha=espaco_linha)

    x3 = x2 + largura_col + 0.25
    _texto(slide, x3, y_colunas, largura_col, 0.28, "DECISÕES", tamanho=13, negrito=True, cor=AZUL_INSTITUCIONAL)
    _lista_com_marcadores(slide, x3, y_colunas + 0.36, largura_col, altura_lista_usada, analise["decisoes"],
                           cor_marcador=AZUL_INSTITUCIONAL, tamanho=tamanho_lista, espaco_linha=espaco_linha)

    y_mensagem = y_colunas + 0.36 + altura_lista_usada + gap_antes_mensagem
    altura_mensagem = max(altura_reservada_mensagem, y_zona_segura_fim - y_mensagem)
    _caixa_leitura(slide, MARGEM_IN, y_mensagem, LARGURA_IN - 2 * MARGEM_IN, altura_mensagem, "Mensagem Central",
                   analise["mensagem"], cor_fundo=OFF_WHITE, cor_rotulo=VERDE_AMAZONIA, tamanho_texto=12)
    return slide


def _slide_scorecard(prs: Presentation, mes_label: str, pagina: int, d: dict):
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Scorecard do Mês", mes_label, pagina, "Leitura executiva e próximo passo por frente da operação")

    linhas_dados = _montar_scorecard(d)
    linhas_tabela = []
    for item in linhas_dados:
        linhas_tabela.append([
            (item["frente"], AZUL_INSTITUCIONAL, True),
            (item["status_label"], item["status_cor"], True),
            (item["leitura"], CINZA_TEXTO, False),
            (item["proximo_passo"], CINZA_TEXTO, False),
        ])

    _tabela(slide, MARGEM_IN, 1.65, LARGURA_IN - 2 * MARGEM_IN, 4.9,
            ["Frente", "Status", "Leitura Executiva", "Próximo Passo"], linhas_tabela,
            larguras_relativas=[2.0, 1.1, 3.4, 3.0], tamanho_fonte=13)
    return slide


def _slide_painel_inventario(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Inventário Item a Item (22/08/2026, reescrito pra bater 1:1 com o
    modelo aprovado v4 — mockup e prints de referência da usuária): KPIs +
    evolução mensal com a LINHA DE TENDÊNCIA REAL desenhada sobre o
    gráfico (ver _grafico_categoria_com_tendencia, não mais só o selo de
    texto) + as 3 tabelas (Resultado por Almoxarifado/Top10 Faltas/Top10
    Sobras) ABAIXO do indicador + caixa de Resumo no rodapé — tudo em UM
    slide só, como no modelo aprovado (antes dividido em 2 slides, com
    "SKUs Recorrentes"/"Cobertura de Conferência" no lugar das tabelas -
    a usuária rejeitou isso como "totalmente diferente do modelo
    aprovado... tabelas abaixo do indicador contando a história do
    recorte", 22/08/2026, "premissas inegociáveis"). SKUs Recorrentes e
    Cobertura de Conferência passaram a viver como conteúdo INFORMATIVO num
    slide companheiro à parte (_slide_painel_inventario_detalhe) - esse
    slide companheiro foi removido depois, ainda em 22/08/2026, a pedido da
    usuária ("fiz uma versão manual [...] esse é o modelo final"; ver
    comentário de decisão antes de _slide_acuracia_ponderada_iap). Não há
    mais nenhum slide no MBR mostrando SKUs Recorrentes/Cobertura de
    Conferência.

    Pedido explícito da mesma mensagem: dar ênfase ao IAP/IAQ (ponderados)
    e deixar esta leitura item a item "apenas para conhecimento" - por
    isso este slide vem DEPOIS dos 4 slides de Acurácia Ponderada na
    ordem do relatório (ver montar_pptx_mbr), mesmo sendo o indicador
    histórico do Atlas; o CONTEÚDO do slide em si segue o modelo aprovado
    sem alteração — a ênfase se resolve por ORDEM no relatório, não
    reduzindo o que este slide mostra."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Inventário Item a Item", mes_label, pagina,
               "Acurácia geral, evolução mensal com tendência, resultado por almoxarifado e ranking financeiro — "
               "leitura item a item, apenas para conhecimento (a leitura ponderada IAP/IAQ é a referência do mês)")

    kpis = d["kpis_inventario"]
    label, cor = _status_maior_melhor(kpis["acuracia_geral_pct"], *_LIMIARES["acuracia"])
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_pct(kpis["acuracia_geral_pct"]), "rotulo": "Acurácia Geral", "contexto": label, "cor_contexto": cor},
        {"valor": _fmt_num(kpis["total_itens"]), "rotulo": "Itens Avaliados"},
        {"valor": _fmt_num(kpis["total_divergentes"]), "rotulo": "Itens Divergentes"},
        {"valor": _fmt_moeda(kpis["resultado_liquido"]), "rotulo": "Resultado Líquido", "cor": COR_SUCESSO if (kpis["resultado_liquido"] or 0) >= 0 else COR_ERRO},
    ], altura=0.85)

    largura_cheia = LARGURA_IN - 2 * MARGEM_IN
    evolucao = d["evolucao_inventario"][-6:]
    valores = [item.get("acuracia_pct") for item in evolucao]
    tendencia = _tendencia_linear(valores)
    if len(evolucao) >= 2:
        categorias = [_nome_mes(item["mes"])[:3] + "/" + item["mes"][2:4] for item in evolucao]
        _texto(slide, MARGEM_IN, 2.78, largura_cheia - 3.2, 0.24, "EVOLUÇÃO MENSAL (MOM) — ACURÁCIA E TENDÊNCIA",
               tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        if tendencia["rotulo"]:
            _texto(slide, MARGEM_IN + largura_cheia - 3.2, 2.78, 3.2, 0.24, tendencia["rotulo"], tamanho=9,
                   negrito=True, cor=tendencia["cor"], alinhamento=PP_ALIGN.RIGHT)
        # Cores por barra iguais ao modelo aprovado (.barra.critico/.barra.ok):
        # vermelho abaixo de 50%, laranja a partir de 50% — a própria barra já
        # avisa o nível crítico, sem depender só do cartão de KPI no topo.
        cores_pontos = [COR_ERRO if (v or 0) < 50 else COR_ATENCAO for v in valores]
        cor_tendencia = tendencia["cor"] if tendencia["cor"] != COR_SEM_DADO else COR_SUCESSO
        _grafico_categoria_com_tendencia(slide, MARGEM_IN, 3.04, largura_cheia, 1.00, categorias,
                                          "Acurácia geral", valores, cor_serie=COR_ATENCAO, cores_pontos=cores_pontos,
                                          cor_tendencia=cor_tendencia)
        y_tabelas = 3.04 + 1.00 + 0.10
    else:
        _caixa_leitura(slide, MARGEM_IN, 2.78, largura_cheia, 1.30, "Evolução",
                       "Sem histórico suficiente de fechamentos para montar a série mensal ainda.")
        y_tabelas = 2.78 + 1.30 + 0.10

    itens_conc = d["concentracao_valor"].get("itens", [])
    y_fim_tabelas, pior, faltas, _sobras = _tabelas_almoxarifado_e_top10(
        slide, y_tabelas, itens_conc, d["comparativo_por_almoxarifado"], "item_a_item_pct", "Acur.",
        ordenar_por="valor", rotulo_qtd_ou_valor="Valor",
    )

    y_resumo = y_fim_tabelas + 0.08
    y_zona_segura_fim = ALTURA_IN - 0.42
    if pior and faltas:
        if tendencia["rotulo"] and tendencia.get("inclinacao") is not None and abs(tendencia["inclinacao"]) > 0.3:
            direcao_tendencia = f", em trajetória de {'melhora' if tendencia['inclinacao'] > 0 else 'queda'} no período (tendência de {_fmt_num(abs(tendencia['inclinacao']), 1)} p.p./mês)"
        else:
            direcao_tendencia = ""
        texto_resumo = (
            f"O mês fecha com acurácia geral de {_fmt_pct(kpis['acuracia_geral_pct'])} ({label.lower()}){direcao_tendencia}. "
            f"{pior['almoxarifado']} tem a menor acurácia item a item do mês ({_fmt_pct(pior.get('item_a_item_pct'))}) — maior "
            f"prioridade de reconferência. O maior exemplo de falta é {faltas[0]['sku']} "
            f"({(faltas[0].get('descricao') or '—')[:30]}, -{_fmt_moeda(faltas[0]['valor'])})."
        )
    else:
        texto_resumo = "Sem dado suficiente neste recorte para fechar a leitura do ciclo."
    altura_resumo = min(max(0.65, _altura_necessaria_caixa_leitura(texto_resumo, largura_cheia, 10)),
                         max(0.5, y_zona_segura_fim - y_resumo - 0.24))
    _caixa_leitura(slide, MARGEM_IN, y_resumo, largura_cheia, altura_resumo, "Resumo do Inventário Item a Item",
                   texto_resumo, cor_fundo=OFF_WHITE, tamanho_texto=10)

    y_rodape = y_resumo + altura_resumo + 0.04
    if y_rodape < y_zona_segura_fim:
        _texto(slide, MARGEM_IN, y_rodape, largura_cheia, 0.20,
               "\"Resultado\" por almoxarifado e os dois rankings cobrem só os itens da lista de concentração de valor "
               "(até 50 maiores por valor do mês).",
               tamanho=8, italico=True, cor=CINZA_TEXTO)
    return slide


def _tabelas_almoxarifado_e_top10(slide, y_topo, itens_conc, comp_almox, campo_pct, rotulo_pct,
                                   ordenar_por="valor", rotulo_qtd_ou_valor="Valor", altura_tabela=1.35,
                                   max_linhas=5):
    """3 tabelas lado a lado — Resultado por Almoxarifado + Top Faltas + Top
    Sobras — reaproveitada pelos 3 slides que o pedem com o mesmo layout
    (Painel de Inventário, IAP, IAQ; 22/08/2026, mockups aprovados v4/v5).
    `campo_pct` escolhe qual dos 3 modelos mostrar na coluna de acurácia
    (item_a_item_pct/iaq_pct/iap_pct); `ordenar_por` escolhe se o ranking é
    por valor financeiro (R$, usado em Painel de Inventário e IAP) ou por
    quantidade (unidades, usado em IAQ — reordena os MESMOS itens de
    concentracao_valor pela magnitude de divergencia_qtd, sem chamada nova
    ao backend).

    Limitação avisada no rodapé de cada slide que chama isto: `itens_conc`
    é a lista de concentração de valor (até 50 maiores POR VALOR do mês) —
    quando ordenado por quantidade (IAQ), um item de alto volume mas baixo
    valor unitário pode estar fora dessa lista e não aparecer no ranking.

    21/08/2026 — CORREÇÃO DE BUG: `altura_tabela` era uma constante fixa de
    2.75in (dimensionada pra até 8 linhas + cabeçalho) independente do
    espaço vertical de fato disponível no slide chamador — nos 3 slides que
    usam isto (Painel/IAP/IAQ), a soma de KPIs + gráfico de evolução + este
    bloco já passava da altura do slide (7.5in) ANTES de reservar espaço pra
    caixa de Resumo/rodapé, o que jogava a caixa de Resumo inteira pra fora
    da área visível (nem aparecia renderizada) — visto na QA visual dos 3
    slides. `altura_tabela`/`max_linhas` agora são parâmetros com default
    reduzido (6 linhas em vez de 8, 1.55in em vez de 2.75in) calibrados pelo
    orçamento vertical real de cada slide chamador (ver `_slide_acuracia_
    ponderada_iap/iaq` e `_slide_painel_inventario`)."""
    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 2 * 0.28) / 3
    x1, x2, x3 = MARGEM_IN, MARGEM_IN + largura_col + 0.28, MARGEM_IN + 2 * (largura_col + 0.28)
    altura_titulo = 0.22
    prejuizo_almox = _prejuizo_por_almoxarifado(itens_conc)

    _texto(slide, x1, y_topo, largura_col, altura_titulo, f"RESULTADO POR ALMOXARIFADO ({rotulo_pct.upper()})",
           tamanho=9.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    comp_ordenado = sorted(comp_almox, key=lambda x: x.get(campo_pct) if x.get(campo_pct) is not None else 999)
    if comp_ordenado:
        linhas_almox = []
        for item in comp_ordenado[:max_linhas]:
            almox = item["almoxarifado"]
            prejuizo = prejuizo_almox.get(almox, 0.0)
            linhas_almox.append([
                almox, _fmt_pct(item.get(campo_pct)),
                (f"-{_fmt_moeda(abs(prejuizo))}" if prejuizo < 0 else _fmt_moeda(prejuizo),
                 COR_ERRO if prejuizo < 0 else COR_SUCESSO, True),
            ])
        _tabela(slide, x1, y_topo + altura_titulo, largura_col, altura_tabela,
                ["Almox.", rotulo_pct[:8], "Result."], linhas_almox,
                larguras_relativas=[1.5, 0.8, 1.1], tamanho_fonte=9)
    else:
        _caixa_leitura(slide, x1, y_topo + altura_titulo, largura_col, altura_tabela, "Almoxarifado",
                       "Sem dado de fechamento por almoxarifado neste recorte.")

    chave_ordem = (lambda i: i["valor"]) if ordenar_por == "valor" else (lambda i: abs(i.get("divergencia_qtd") or 0))
    faltas = sorted([i for i in itens_conc if (i.get("divergencia_qtd") or 0) < 0], key=chave_ordem, reverse=True)[:max_linhas]
    sobras = sorted([i for i in itens_conc if (i.get("divergencia_qtd") or 0) > 0], key=chave_ordem, reverse=True)[:max_linhas]

    def _rotulo_item(item):
        return _fmt_moeda(item["valor"]) if ordenar_por == "valor" else f"{_fmt_num(abs(item.get('divergencia_qtd') or 0))} un."

    _texto(slide, x2, y_topo, largura_col, altura_titulo, f"TOP FALTAS DO MÊS ({rotulo_qtd_ou_valor.upper()})",
           tamanho=9.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    if faltas:
        linhas_faltas = [[it["sku"], (it.get("descricao") or "—")[:15], (f"-{_rotulo_item(it)}", COR_ERRO, True)] for it in faltas]
        _tabela(slide, x2, y_topo + altura_titulo, largura_col, altura_tabela,
                ["SKU", "Descrição", rotulo_qtd_ou_valor[:8]], linhas_faltas,
                larguras_relativas=[1.1, 1.6, 1.0], tamanho_fonte=9)
    else:
        _caixa_leitura(slide, x2, y_topo + altura_titulo, largura_col, altura_tabela, "Faltas",
                       "Sem falta calculável neste recorte.")

    _texto(slide, x3, y_topo, largura_col, altura_titulo, f"TOP SOBRAS DO MÊS ({rotulo_qtd_ou_valor.upper()})",
           tamanho=9.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    if sobras:
        linhas_sobras = [[it["sku"], (it.get("descricao") or "—")[:15], (f"+{_rotulo_item(it)}", COR_SUCESSO, True)] for it in sobras]
        _tabela(slide, x3, y_topo + altura_titulo, largura_col, altura_tabela,
                ["SKU", "Descrição", rotulo_qtd_ou_valor[:8]], linhas_sobras,
                larguras_relativas=[1.1, 1.6, 1.0], tamanho_fonte=9)
    else:
        _caixa_leitura(slide, x3, y_topo + altura_titulo, largura_col, altura_tabela, "Sobras",
                       "Sem sobra calculável neste recorte.")

    return y_topo + altura_titulo + altura_tabela, (comp_ordenado[0] if comp_ordenado else None), faltas, sobras


# 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o modelo
# final"): o slide "Painel de Inventário — Detalhamento Financeiro" (função
# _slide_painel_inventario_detalhe, companheiro informativo do Painel de
# Inventário com SKUs Recorrentes + Cobertura de Conferência) foi removido -
# ela editou o MBR manualmente pra tirar esse tópico. Os dados que só esse
# slide usava (dados["top_recorrentes"], dados["cobertura_conferencia"])
# foram removidos junto de _coletar_dados_mbr, senão ficariam sendo buscados
# no banco a cada MBR gerado sem nenhum slide pra usar o resultado.


def _slide_acuracia_ponderada_iap(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Acurácia Ponderada — por Valor (IAP), slide próprio (22/08/2026,
    mockup aprovado v5 — separa o que antes era um único slide IAQ+IAP em
    dois, cada um com sua própria evolução/tendência, resultado por
    almoxarifado e ranking financeiro em R$).

    22/08/2026, pedido não-negociável da usuária: a linha de tendência
    agora é desenhada de verdade sobre o gráfico (ver
    _grafico_categoria_com_tendencia), não só o selo de texto; e o slide
    ganhou uma caixa de Resumo no rodapé, depois das 3 tabelas ("mal
    enquadrado sem a mensagem no final, abaixo das tabelas")."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Acurácia Ponderada — por Valor (IAP)", mes_label, pagina,
               "Leitura ponderada pelo valor financeiro do inventário, fechamento por almoxarifado e ranking financeiro")

    comp = d["comparativo_acuracia"]
    evolucao = d["evolucao_ponderada"][-6:]
    ultimo = evolucao[-1] if evolucao else {}
    gap = comp.get("gap_item_vs_iap_pp")
    delta_mom = ultimo.get("variacao_iap_pp")
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_pct(comp.get("iap_pct")), "rotulo": "IAP (por Valor)", "cor": VERDE_AMAZONIA},
        {"valor": (f"R$ {_fmt_num(ultimo.get('valor_mod'), 2)}" if ultimo.get("valor_mod") is not None else "—"),
         "rotulo": "Valor do Inventário em Jogo", "contexto": "Valor Mod — sobra vs. falta"},
        {"valor": _fmt_pct(abs(gap)) if gap is not None else "—", "rotulo": "Distorção (Item vs. IAP)",
         "cor": COR_ATENCAO if (gap is not None and abs(gap) >= 3) else COR_SUCESSO},
        {"valor": (f"+{_fmt_pct(delta_mom)}" if delta_mom is not None and delta_mom >= 0 else _fmt_pct(delta_mom)),
         "rotulo": "Variação do IAP (MoM)", "cor": COR_SUCESSO if (delta_mom or 0) >= 0 else COR_ERRO,
         "contexto": "Melhora" if (delta_mom or 0) >= 0 else ("Piora" if delta_mom is not None else None),
         "cor_contexto": COR_SUCESSO if (delta_mom or 0) >= 0 else COR_ERRO},
    ], altura=0.85)

    largura_cheia = LARGURA_IN - 2 * MARGEM_IN
    if len(evolucao) >= 2:
        categorias = [_nome_mes(item["mes"])[:3] + "/" + item["mes"][2:4] for item in evolucao]
        valores_iap = [item.get("iap_pct") or 0 for item in evolucao]
        tendencia = _tendencia_linear([item.get("iap_pct") for item in evolucao])
        _texto(slide, MARGEM_IN, 2.78, largura_cheia - 3.2, 0.24, "EVOLUÇÃO MENSAL (MOM) — IAP", tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        if tendencia["rotulo"]:
            _texto(slide, MARGEM_IN + largura_cheia - 3.2, 2.78, 3.2, 0.24, tendencia["rotulo"], tamanho=9,
                   negrito=True, cor=tendencia["cor"], alinhamento=PP_ALIGN.RIGHT)
        cor_tendencia = tendencia["cor"] if tendencia["cor"] != COR_SEM_DADO else COR_SUCESSO
        _grafico_categoria_com_tendencia(slide, MARGEM_IN, 3.04, largura_cheia, 1.00, categorias, "IAP (por valor)",
                                          valores_iap, cor_serie=AZUL_INSTITUCIONAL, cor_tendencia=cor_tendencia)
        y_tabelas = 3.04 + 1.00 + 0.10
    else:
        tendencia = {"inclinacao": None, "rotulo": None, "cor": COR_SEM_DADO}
        _caixa_leitura(slide, MARGEM_IN, 2.78, largura_cheia, 1.30, "Evolução do IAP",
                       "Sem histórico suficiente de fechamentos ainda.")
        y_tabelas = 2.78 + 1.30 + 0.10

    itens_conc = d["concentracao_valor"].get("itens", [])
    y_fim_tabelas, pior, faltas, sobras = _tabelas_almoxarifado_e_top10(
        slide, y_tabelas, itens_conc, d["comparativo_por_almoxarifado"], "iap_pct", "IAP",
        ordenar_por="valor", rotulo_qtd_ou_valor="Valor",
    )

    # Caixa de Resumo (22/08/2026, mockup aprovado v5) — mesmo padrão visual
    # "resumo-ciclo" já usado no Inventário Item a Item, agora também aqui,
    # depois das 3 tabelas.
    y_resumo = y_fim_tabelas + 0.08
    y_zona_segura_fim = ALTURA_IN - 0.42
    partes_resumo = [f"O IAP fecha o mês em {_fmt_pct(comp.get('iap_pct'))}"]
    if delta_mom is not None:
        partes_resumo.append(f", com {'melhora' if delta_mom >= 0 else 'piora'} de {_fmt_pct(abs(delta_mom))} frente ao mês anterior")
    if tendencia.get("inclinacao") is not None and abs(tendencia["inclinacao"]) > 0.3:
        partes_resumo.append(f" e tendência linear de {_fmt_num(abs(tendencia['inclinacao']), 1)} p.p./mês no período")
    if ultimo.get("valor_mod") is not None:
        partes_resumo.append(f" — há {_fmt_moeda(abs(ultimo['valor_mod']))} em jogo entre sobra e falta (Valor Mod)")
    partes_resumo.append(
        ". O IAP pondera por valor financeiro — poucos itens de alto valor bastam pra manter o indicador alto mesmo "
        "com mais SKUs divergentes item a item."
    )
    if pior:
        partes_resumo.append(f" {pior['almoxarifado']} concentra o maior prejuízo ponderado por valor no período.")
    texto_resumo = "".join(partes_resumo)
    altura_resumo = min(max(0.65, _altura_necessaria_caixa_leitura(texto_resumo, largura_cheia, 10)),
                         max(0.5, y_zona_segura_fim - y_resumo - 0.24))
    _caixa_leitura(slide, MARGEM_IN, y_resumo, largura_cheia, altura_resumo, "Resumo da Acurácia Ponderada por Valor (IAP)",
                   texto_resumo, cor_fundo=OFF_WHITE, tamanho_texto=10)

    y_rodape = y_resumo + altura_resumo + 0.04
    if y_rodape < y_zona_segura_fim:
        _texto(slide, MARGEM_IN, y_rodape, largura_cheia, 0.20,
               "IAP pondera cada divergência pelo valor financeiro do SKU. Ranking e resultado por almoxarifado "
               "cobrem só os itens da lista de concentração de valor (até 50 maiores por valor do mês).",
               tamanho=8, italico=True, cor=CINZA_TEXTO)
    return slide


def _slide_acuracia_ponderada_iaq(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Acurácia Ponderada — por Quantidade (IAQ), slide próprio (22/08/2026,
    mockup aprovado v5 — espelha o slide de IAP, mas ponderando/ordenando
    por unidades, não por R$).

    22/08/2026, pedido não-negociável da usuária: linha de tendência real
    desenhada sobre o gráfico (ver _grafico_categoria_com_tendencia) +
    caixa de Resumo no rodapé, depois das 3 tabelas — mesmo ajuste feito
    no slide de IAP."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Acurácia Ponderada — por Quantidade (IAQ)", mes_label, pagina,
               "Leitura ponderada pela quantidade de itens do inventário, fechamento por almoxarifado e ranking por unidades")

    comp = d["comparativo_acuracia"]
    evolucao = d["evolucao_ponderada"][-6:]
    ultimo = evolucao[-1] if evolucao else {}
    gap_iaq = None
    if comp.get("item_a_item_pct") is not None and comp.get("iaq_pct") is not None:
        gap_iaq = round(comp["item_a_item_pct"] - comp["iaq_pct"], 2)
    delta_mom = ultimo.get("variacao_iaq_pp")
    # Nota: o mockup aprovado v5 tinha uma 2ª métrica "Quantidade em Jogo" (Qtd.
    # Mod, equivalente por unidades ao Valor Mod do IAP) — mas o próprio mockup
    # já marcava esse número como "ilustrativo", e não existe hoje um cálculo
    # real equivalente no backend (dashboard_evolucao_ponderada_mensal só tem
    # valor_mod, em R$, não uma versão em unidades). Em vez de inventar o
    # número, o 2º cartão aqui usa "Itens Divergentes no Mês" (kpis_inventario,
    # já real e coerente com uma leitura por quantidade).
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_pct(comp.get("iaq_pct")), "rotulo": "IAQ (por Quantidade)"},
        {"valor": _fmt_num(d["kpis_inventario"].get("total_divergentes")), "rotulo": "Itens Divergentes no Mês"},
        {"valor": _fmt_pct(abs(gap_iaq)) if gap_iaq is not None else "—", "rotulo": "Distorção (Item vs. IAQ)",
         "cor": COR_ATENCAO if (gap_iaq is not None and abs(gap_iaq) >= 3) else COR_SUCESSO},
        {"valor": (f"+{_fmt_pct(delta_mom)}" if delta_mom is not None and delta_mom >= 0 else _fmt_pct(delta_mom)),
         "rotulo": "Variação do IAQ (MoM)", "cor": COR_SUCESSO if (delta_mom or 0) >= 0 else COR_ERRO,
         "contexto": "Melhora" if (delta_mom or 0) >= 0 else ("Piora" if delta_mom is not None else None),
         "cor_contexto": COR_SUCESSO if (delta_mom or 0) >= 0 else COR_ERRO},
    ], altura=0.85)

    largura_cheia = LARGURA_IN - 2 * MARGEM_IN
    if len(evolucao) >= 2:
        categorias = [_nome_mes(item["mes"])[:3] + "/" + item["mes"][2:4] for item in evolucao]
        valores_iaq = [item.get("iaq_pct") or 0 for item in evolucao]
        tendencia = _tendencia_linear([item.get("iaq_pct") for item in evolucao])
        _texto(slide, MARGEM_IN, 2.78, largura_cheia - 3.2, 0.24, "EVOLUÇÃO MENSAL (MOM) — IAQ", tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        if tendencia["rotulo"]:
            _texto(slide, MARGEM_IN + largura_cheia - 3.2, 2.78, 3.2, 0.24, tendencia["rotulo"], tamanho=9,
                   negrito=True, cor=tendencia["cor"], alinhamento=PP_ALIGN.RIGHT)
        cor_tendencia = tendencia["cor"] if tendencia["cor"] != COR_SEM_DADO else COR_SUCESSO
        _grafico_categoria_com_tendencia(slide, MARGEM_IN, 3.04, largura_cheia, 1.00, categorias, "IAQ (quantidade)",
                                          valores_iaq, cor_serie=VERDE_AMAZONIA, cor_tendencia=cor_tendencia)
        y_tabelas = 3.04 + 1.00 + 0.10
    else:
        tendencia = {"inclinacao": None, "rotulo": None, "cor": COR_SEM_DADO}
        _caixa_leitura(slide, MARGEM_IN, 2.78, largura_cheia, 1.30, "Evolução do IAQ",
                       "Sem histórico suficiente de fechamentos ainda.")
        y_tabelas = 2.78 + 1.30 + 0.10

    itens_conc = d["concentracao_valor"].get("itens", [])
    y_fim_tabelas, pior, faltas, sobras = _tabelas_almoxarifado_e_top10(
        slide, y_tabelas, itens_conc, d["comparativo_por_almoxarifado"], "iaq_pct", "IAQ",
        ordenar_por="quantidade", rotulo_qtd_ou_valor="Qtd.",
    )

    # Caixa de Resumo (22/08/2026, mockup aprovado v5) — mesmo padrão do slide de IAP.
    y_resumo = y_fim_tabelas + 0.08
    y_zona_segura_fim = ALTURA_IN - 0.42
    partes_resumo = [f"O IAQ fecha o mês em {_fmt_pct(comp.get('iaq_pct'))}"]
    if delta_mom is not None:
        partes_resumo.append(f", com {'melhora' if delta_mom >= 0 else 'piora'} de {_fmt_pct(abs(delta_mom))} frente ao mês anterior")
    if tendencia.get("inclinacao") is not None and abs(tendencia["inclinacao"]) > 0.3:
        conector = "mesmo com" if (delta_mom is not None and delta_mom < 0 and tendencia["inclinacao"] > 0) else "e"
        partes_resumo.append(f" — {conector} tendência linear de {_fmt_num(abs(tendencia['inclinacao']), 1)} p.p./mês no período")
    partes_resumo.append(
        ". Diferente do IAP, o IAQ pesa mais divergências em itens de giro alto e baixo valor unitário."
    )
    if pior:
        partes_resumo.append(f" {pior['almoxarifado']} concentra a maior perda em unidades no período.")
    texto_resumo = "".join(partes_resumo)
    altura_resumo = min(max(0.65, _altura_necessaria_caixa_leitura(texto_resumo, largura_cheia, 10)),
                         max(0.5, y_zona_segura_fim - y_resumo - 0.24))
    _caixa_leitura(slide, MARGEM_IN, y_resumo, largura_cheia, altura_resumo, "Resumo da Acurácia Ponderada por Quantidade (IAQ)",
                   texto_resumo, cor_fundo=OFF_WHITE, tamanho_texto=10)

    y_rodape = y_resumo + altura_resumo + 0.04
    if y_rodape < y_zona_segura_fim:
        _texto(slide, MARGEM_IN, y_rodape, largura_cheia, 0.20,
               "IAQ pondera cada divergência pela quantidade de itens do SKU (não pelo valor). Ranking por unidades "
               "calculado sobre os 50 maiores itens POR VALOR do mês — pode omitir um item de alto volume e baixo "
               "valor unitário fora dessa lista.",
               tamanho=8, italico=True, cor=CINZA_TEXTO)
    return slide


def _slide_acuracia_ponderada_detalhe(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Segundo slide de Acurácia Ponderada (20/08/2026, pedido do usuário;
    ampliado em 22/08/2026, mockup aprovado v8): curva de Pareto agora com
    20 SKUs individuais + uma categoria agregada pro resto da cauda ("de
    ponta a ponta" até 100% — o ponto agregado usa 100% de propósito, já
    que por definição inclui todo o restante dos itens divergentes, sem
    precisar do valor exato de cada um deles), e uma tabela Top 10
    detalhada (Almoxarifado/Qtd. Sistêmica/Qtd. Conferência/Dif./Valor/%
    Acum. — todos campos que já vêm de dashboard_concentracao_valor, sem
    chamada nova ao backend) no lugar da versão anterior de 3 linhas.

    22/08/2026, pedido não-negociável da usuária ("curva de pareto com
    top 10 itens, depois análise pela distribuição da magnitude"): a
    curva agora é um gráfico COMBO dual-axis em largura total — barras de
    Valor (R$) por SKU (+ barra agregada da cauda, com o valor REAL
    somado do restante, não ilustrativo) e a linha de % acumulado sobre
    eixo secundário — igual ao modelo aprovado, no lugar do gráfico de
    linha isolado que só mostrava o % acumulado sem as barras de valor.
    "Distribuição por Magnitude" saiu deste slide (virou o gráfico combo
    do próximo slide, Detalhamento por Faixa — ver _slide_acuracia_
    ponderada_faixas) para abrir espaço pro Pareto em largura total,
    igual ao modelo aprovado. Eixo duplo aqui é excepção documentada à
    convenção "nunca eixo duplo" do resto do MBR — ver nota no bloco de
    _grafico_combo_dual_eixo."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Acurácia Ponderada — Concentração de Risco", mes_label, pagina,
               "Curva de Pareto — onde o valor em risco está concentrado neste mês")

    largura_cheia = LARGURA_IN - 2 * MARGEM_IN
    pareto = d["concentracao_valor"]
    itens_pareto = pareto.get("itens", [])[:20]
    total_divergentes = pareto.get("total_itens_divergentes") or 0
    valor_total = pareto.get("valor_total")
    if itens_pareto:
        categorias = [item["sku"] for item in itens_pareto]
        valores_barra = [item["valor"] for item in itens_pareto]
        valores_pct_acum = [item["pct_valor_acumulado"] for item in itens_pareto]
        resto = total_divergentes - len(itens_pareto)
        if resto > 0:
            soma_top = sum(valores_barra)
            valor_resto = max((valor_total or 0) - soma_top, 0) if valor_total is not None else 0
            categorias.append(f"+{_fmt_num(resto)} itens")
            valores_barra.append(valor_resto)
            valores_pct_acum.append(100.0)
        _texto(slide, MARGEM_IN, 1.65, largura_cheia, 0.24,
               f"CONCENTRAÇÃO DE VALOR (CURVA DE PARETO) — TOP {len(itens_pareto)} SKUS + CAUDA",
               tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_combo_dual_eixo(slide, MARGEM_IN, 1.93, largura_cheia, 1.85, categorias,
                                  "Valor (R$)", valores_barra, "% acumulado do valor", valores_pct_acum,
                                  cor_barra=VERDE_AMAZONIA, cor_linha=COR_ATENCAO,
                                  formato_numero_barra='#,##0', formato_numero_linha='0"%"')

        top_n_pct = pareto.get("top_n_pct_do_valor")
        top_n = pareto.get("top_n")
        linha_pareto_texto = (
            f"Os {top_n} maiores SKUs divergentes concentram {_fmt_pct(top_n_pct)} do valor em risco do mês "
            f"({_fmt_moeda(pareto.get('valor_total'))} no total, {_fmt_num(total_divergentes)} itens divergentes)."
            if top_n_pct is not None else "Concentração de valor não disponível neste recorte."
        )
        _texto(slide, MARGEM_IN, 3.83, largura_cheia, 0.35, linha_pareto_texto, tamanho=10, cor=CINZA_TEXTO)
    else:
        _caixa_leitura(slide, MARGEM_IN, 1.65, largura_cheia, 2.4, "Curva de Pareto",
                       "Sem item divergente com custo cadastrado neste recorte para montar a curva de concentração de valor.")

    y_tabela_titulo = 4.30
    _texto(slide, MARGEM_IN, y_tabela_titulo, largura_cheia, 0.24, "TOP 10 MAIORES EXEMPLOS DO MÊS",
           tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    top10 = pareto.get("itens", [])[:10]
    if top10:
        linhas_top10 = [
            [
                it["sku"], (it.get("descricao") or "—")[:22], it.get("almoxarifado") or "—",
                _fmt_num(it.get("qtd_sistema")), _fmt_num(it.get("qtd_contagem")),
                (_fmt_num(it.get("divergencia_qtd")), COR_ERRO if (it.get("divergencia_qtd") or 0) < 0 else COR_SUCESSO, False),
                _fmt_moeda(it["valor"]), _fmt_pct(it["pct_valor_acumulado"]),
            ]
            for it in top10
        ]
        y_zona_segura_fim = ALTURA_IN - 0.42
        _tabela(slide, MARGEM_IN, y_tabela_titulo + 0.26, largura_cheia, y_zona_segura_fim - (y_tabela_titulo + 0.26),
                ["SKU", "Descrição", "Almox.", "Qtd. Sist.", "Qtd. Conf.", "Dif.", "Valor", "% Acum."], linhas_top10,
                larguras_relativas=[1.05, 2.1, 1.3, 0.95, 0.95, 0.8, 1.1, 0.85], tamanho_fonte=9.5)
    else:
        _texto(slide, MARGEM_IN, y_tabela_titulo + 0.30, largura_cheia, 0.4,
               "Sem exemplo disponível neste recorte.", tamanho=10.5, cor=CINZA_TEXTO)
    return slide


def _slide_acuracia_ponderada_faixas(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Detalhamento por Faixa de Magnitude (22/08/2026, mockup aprovado v8) —
    Top 10 por valor dentro de cada faixa (0-5un/5-20un/20-100un/mais de
    100 un.), lado a lado, via dashboard_itens_por_magnitude (função já
    existente e já usada pelo duplo-clique na tela — uma chamada por
    faixa, ver _coletar_dados_mbr/magnitude_por_faixa_itens).

    22/08/2026, pedido não-negociável da usuária: ganhou o gráfico combo
    "Distribuição por Magnitude da Divergência" (barras = Nº de itens,
    linha = Valor total R$, eixo duplo — mesma excepção documentada da
    Concentração de Risco) ACIMA das 4 colunas — antes vivia como painel
    lateral no slide de Concentração de Risco, o modelo aprovado o move
    pra este slide, em largura total; tabelas por faixa foram de Top 5
    pra Top 10; e ganhou a caixa "Resumo do Detalhamento por Faixa" depois
    das 4 colunas, igual ao modelo aprovado."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    # Título encurtado (22/08/2026): a versão completa ("... de Magnitude")
    # passa de ~45 caracteres e quebra em 2 linhas dentro da caixa de título
    # do cabeçalho (27pt/9.7in), invadindo a caixa de subtítulo de posição
    # fixa - "de Magnitude" já fica implícito no subtítulo e nos rótulos das
    # 4 colunas (0 a 5 un./5 a 20 un./etc.). Também alinha com o nome já
    # usado na lista de itens da capa de seção (_secao, ver montar_pptx_mbr).
    _cabecalho(slide, "Acurácia Ponderada — Detalhamento por Faixa", mes_label, pagina,
               "Top 10 por valor dentro de cada faixa — onde o risco financeiro está dentro de cada grupo de tamanho")

    largura_total = LARGURA_IN - 2 * MARGEM_IN
    y_zona_segura_fim = ALTURA_IN - 0.42
    magnitude = d["distribuicao_magnitude"]
    faixas_resumo_lista = magnitude.get("faixas", [])
    faixas_resumo = {f["faixa"]: f for f in faixas_resumo_lista}

    # Orçamento vertical explícito (21/08/2026, CORREÇÃO DE BUG): a versão
    # anterior deste slide reservava um "espaço pro resumo" (1.05+0.22) só
    # no cálculo de altura_tabelas_faixa, mas os offsets fixos usados
    # DEPOIS pra posicionar a tabela (+0.50) e o resumo (+0.10) não entravam
    # nessa conta — o valor realmente disponível pro resumo, medido na QA
    # visual, ficava perto de 0.43in (não os ~1.05in pretendidos), cortando
    # o texto do resumo com "…" bem antes do fim da frase. Reduz também o
    # gráfico combo/legenda de topo (1.10 -> 0.80in) pra abrir espaço.
    y_topo = 1.62
    if faixas_resumo_lista and magnitude.get("total_divergentes"):
        categorias_mag = [f["faixa"] for f in faixas_resumo_lista]
        valores_qtd = [f.get("quantidade_itens") or 0 for f in faixas_resumo_lista]
        valores_valor = [f.get("valor_total") or 0 for f in faixas_resumo_lista]
        _texto(slide, MARGEM_IN, y_topo, largura_total, 0.20, "DISTRIBUIÇÃO POR MAGNITUDE DA DIVERGÊNCIA",
               tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_combo_dual_eixo(slide, MARGEM_IN, y_topo + 0.22, largura_total, 0.90, categorias_mag,
                                  "Nº de itens", valores_qtd, "Valor total (R$)", valores_valor,
                                  cor_barra=AZUL_INSTITUCIONAL, cor_linha=COR_ATENCAO,
                                  formato_numero_barra="0", formato_numero_linha='"R$" #,##0',
                                  mostrar_rotulos_barra=True)
        pct_pequenas = magnitude.get("pct_divergencias_pequenas")
        maior_faixa_valor = max(faixas_resumo_lista, key=lambda f: f.get("valor_total") or 0, default=None)
        if pct_pequenas is not None and maior_faixa_valor:
            texto_magnitude = (
                f"{_fmt_pct(pct_pequenas)} das divergências do mês são pequenas (0 a 5 unidades), mas a faixa "
                f"\"{maior_faixa_valor['faixa']}\" — {_fmt_num(maior_faixa_valor.get('quantidade_itens'))} item(ns) — concentra o "
                f"maior valor em risco ({_fmt_moeda(maior_faixa_valor.get('valor_total'))}). É esse contraste que a métrica item a "
                "item não mostra. Abaixo, o top 10 por valor dentro de cada faixa."
            )
        else:
            texto_magnitude = "Sem divergência neste recorte para distribuir por magnitude."
        y_caption = y_topo + 0.22 + 0.90 + 0.03
        _texto(slide, MARGEM_IN, y_caption, largura_total, 0.22, texto_magnitude, tamanho=9.5, cor=CINZA_TEXTO)
        y_faixas = y_caption + 0.22 + 0.04
    else:
        _caixa_leitura(slide, MARGEM_IN, y_topo, largura_total, 1.50, "Distribuição por magnitude",
                       "Sem divergência neste recorte para distribuir por magnitude.")
        y_faixas = y_topo + 1.62

    faixas_itens = d.get("magnitude_por_faixa_itens", [])
    cores_faixa = [COR_SUCESSO, COR_INFO, COR_ATENCAO, COR_ERRO]

    n = len(faixas_itens) or 1
    gap = 0.18
    largura_col = (largura_total - gap * (n - 1)) / n
    # Cabeçalho de cada coluna de faixa (título + subtítulo) antes da
    # tabela, e reserva pro resumo/rodapé DEPOIS da tabela — os mesmos
    # offsets fixos usados abaixo pra posicionar tabela/resumo, agora
    # somados aqui também, pra altura_tabelas_faixa refletir o espaço que
    # de fato resta (não um orçamento que ignora esses offsets).
    altura_cabecalho_faixa = 0.42
    gap_antes_resumo = 0.06
    reserva_rodape = 0.22
    # 0.86in reservado pro resumo: a tabela de 11 linhas (cabeçalho + Top 10)
    # precisa de ~2.50in reais pra não estourar - medido empiricamente (não
    # só estimado): 2.34in e até 2.40in ainda cortavam a última linha na
    # renderização real (LibreOffice/PowerPoint impõem uma altura mínima por
    # linha que a estimativa por caractere não captura com exatidão) - só a
    # partir de ~2.50in a tabela efetivamente coube. O texto do resumo (a
    # versão mais longa, com os 2 parágrafos de contraste) precisa de
    # ~0.85in a 9pt — 0.86in dá uma margem mínima de segurança.
    reserva_resumo = 0.86
    altura_tabelas_faixa = max(
        2.50,
        y_zona_segura_fim - y_faixas - altura_cabecalho_faixa - gap_antes_resumo - reserva_resumo - reserva_rodape,
    )

    algum_item = False
    pior_faixa_para_resumo = None
    x = MARGEM_IN
    for idx, bloco in enumerate(faixas_itens):
        faixa_nome = bloco.get("faixa", "—")
        resumo_faixa = faixas_resumo.get(faixa_nome, {})
        cor_faixa = cores_faixa[idx % len(cores_faixa)]
        _texto(slide, x, y_faixas, largura_col, 0.24, faixa_nome.upper(), tamanho=10, negrito=True, cor=cor_faixa)
        subtitulo_faixa = (
            f"{_fmt_num(resumo_faixa.get('quantidade_itens'))} itens · {_fmt_moeda(resumo_faixa.get('valor_total'))} no total"
            if resumo_faixa else ""
        )
        _texto(slide, x, y_faixas + 0.24, largura_col, 0.22, subtitulo_faixa, tamanho=9, cor=CINZA_TEXTO)
        itens_faixa = sorted(bloco.get("itens", []), key=lambda i: i.get("valor") or 0, reverse=True)[:10]
        if itens_faixa:
            algum_item = True
            linhas = [[(it.get("descricao") or it.get("sku") or "—")[:22], _fmt_moeda(it.get("valor"))] for it in itens_faixa]
            _tabela(slide, x, y_faixas + altura_cabecalho_faixa, largura_col, altura_tabelas_faixa,
                    ["Descrição", "Valor"], linhas, larguras_relativas=[2.0, 1.0], tamanho_fonte=9)
        else:
            _caixa_leitura(slide, x, y_faixas + altura_cabecalho_faixa, largura_col, altura_tabelas_faixa, "Sem itens",
                           "Nenhum item divergente nesta faixa neste recorte.", tamanho_texto=9.5)
        x += largura_col + gap

    if not faixas_itens or not algum_item:
        _caixa_leitura(slide, MARGEM_IN, y_faixas, largura_total, y_zona_segura_fim - y_faixas,
                       "Detalhamento por Faixa", "Sem divergência neste recorte para detalhar por faixa de magnitude.")
        return slide

    # Resumo do Detalhamento por Faixa (22/08/2026, mockup aprovado v8) — compara a
    # faixa com mais itens (geralmente "0 a 5 un.") contra a faixa de maior valor em
    # risco (geralmente "mais de 100 un.") — o mesmo contraste que o gráfico combo
    # acima mostra visualmente, agora em texto.
    y_resumo = y_faixas + altura_cabecalho_faixa + altura_tabelas_faixa + gap_antes_resumo
    largura_cheia = largura_total
    faixa_mais_itens = max(faixas_resumo_lista, key=lambda f: f.get("quantidade_itens") or 0, default=None)
    faixa_mais_valor = max(faixas_resumo_lista, key=lambda f: f.get("valor_total") or 0, default=None)
    if faixa_mais_itens and faixa_mais_valor:
        if faixa_mais_itens["faixa"] == faixa_mais_valor["faixa"]:
            texto_resumo = (
                f"A faixa \"{faixa_mais_valor['faixa']}\" concentra tanto o maior número de itens divergentes "
                f"({_fmt_num(faixa_mais_valor.get('quantidade_itens'))}) quanto o maior valor em risco "
                f"({_fmt_moeda(faixa_mais_valor.get('valor_total'))}) neste recorte — sem o contraste de volume x valor "
                "que costuma aparecer entre faixas pequenas e grandes."
            )
        else:
            texto_resumo = (
                f"A faixa \"{faixa_mais_valor['faixa']}\" tem só {_fmt_num(faixa_mais_valor.get('quantidade_itens'))} item(ns), mas "
                f"concentra o maior valor em risco ({_fmt_moeda(faixa_mais_valor.get('valor_total'))}) — poucos SKUs de giro alto "
                f"puxando o financeiro sozinhos. Já \"{faixa_mais_itens['faixa']}\" tem {_fmt_num(faixa_mais_itens.get('quantidade_itens'))} "
                f"itens somando {_fmt_moeda(faixa_mais_itens.get('valor_total'))} — volume grande, mas diluído por item. É esse "
                "contraste que a leitura item a item não mostra."
            )
    else:
        texto_resumo = "Sem dado suficiente neste recorte para fechar a leitura por faixa."
    altura_resumo = min(max(0.65, _altura_necessaria_caixa_leitura(texto_resumo, largura_cheia, 9)),
                         max(0.5, y_zona_segura_fim - y_resumo - reserva_rodape))
    _caixa_leitura(slide, MARGEM_IN, y_resumo, largura_cheia, altura_resumo, "Resumo do Detalhamento por Faixa",
                   texto_resumo, cor_fundo=OFF_WHITE, tamanho_texto=9)

    y_rodape = y_resumo + altura_resumo + 0.04
    if y_rodape < y_zona_segura_fim:
        _texto(slide, MARGEM_IN, y_rodape, largura_cheia, 0.20,
               "Top 10 por valor calculado dentro de cada faixa (dashboard_itens_por_magnitude).",
               tamanho=8, italico=True, cor=CINZA_TEXTO)
    return slide


def _leitura_passivos(resumo: dict) -> str:
    """Versão curta (2-3 frases), pensada pra caber num slide, do mesmo
    dado por trás de _montar_resumo_narrado (baixas_operacionais_router) -
    aquele texto foi escrito pro pop-up de duplo-clique da tela (parágrafo
    longo, detalhando cada categoria uma por uma) e estoura qualquer caixa
    de slide de forma feia se usado direto aqui. Em vez de cortar esse
    texto no meio de uma frase, monta uma leitura própria e objetiva a
    partir dos MESMOS números - foco no insight que mais importa (o quanto
    do valor aprovado como baixa ainda não foi confirmado pelo inventário
    físico), por pedido explícito do usuário de "focar em resultado,
    números, menos teoria"."""
    passivos = resumo["passivos"]
    resultado_inv = resumo["resultado_inventario"]
    por_categoria = passivos.get("por_categoria", {})
    aguardando = por_categoria.get("aguardando_divergencia", {})

    partes = []
    if passivos["valor"]:
        if aguardando.get("valor"):
            partes.append(
                f"{_fmt_moeda(passivos['valor'])} em passivos aprovados no período, dos quais "
                f"{_fmt_moeda(aguardando['valor'])} ainda aguardam cruzamento com uma divergência."
            )
        else:
            partes.append(f"{_fmt_moeda(passivos['valor'])} em passivos aprovados no período, já mapeados por categoria.")
    else:
        partes.append("Nenhum passivo aprovado sem mapeamento neste recorte.")

    valor_resultado = resultado_inv.get("resultado_valor")
    if valor_resultado is not None:
        if valor_resultado >= 0:
            partes.append(f"O inventário acumulado no período é positivo em {_fmt_moeda(valor_resultado)}.")
        else:
            diferenca = abs(passivos["valor"] - abs(valor_resultado))
            partes.append(
                f"A perda física medida pelo inventário ({_fmt_moeda(abs(valor_resultado))}) é menor que o valor "
                f"aprovado como baixa — diferença de {_fmt_moeda(diferenca)}, possivelmente de baixas de um período "
                "ainda não coberto pelo inventário físico."
            )
    return " ".join(partes)


# 22/08/2026, pedido não-negociável da usuária ("substitua os indicadores
# nativos - Mapeamento de Passivos passa a ser os arquivos externos com
# modelos já aprovados"): os 3 slides nativos que viviam aqui -
# _slide_mapeamento_passivos, _slide_passivos_evolucao e _slide_shelf_life -
# foram REMOVIDOS do relatório (não só reordenados). O Dashboard Baixas
# Operacionais e o Farol de Shelf-Life, ambos dashboards externos com modelo
# já aprovado (ver _slide_baixas_operacionais_externo/_evolucao e
# _slide_farol_shelf_externo/_risco_almoxarifado, abaixo), passam a ser a
# ÚNICA fonte pra esses dois assuntos.
#
# Aviso explícito ciente e confirmado pela usuária antes desta remoção: os
# 3 slides nativos removidos usavam dado 100% exato, direto do banco do
# Atlas (fechamento do mês exato; categorização de motivo própria do
# Atlas), enquanto os dashboards externos que os substituem são retrato
# datado de uma janela diferente (~60 dias corridos no caso de Baixas
# Operacionais; momento da exportação no caso do Farol de Shelf-Life) e
# usam categorização de motivo própria da equipe (Stock Savvy), não a do
# Atlas - os números dos dois PODEM não bater entre si por desenho. A
# usuária optou por manter só a fonte externa mesmo assim.
#
# Os cálculos que alimentavam esses 3 slides (d["resumo_passivos"],
# d["evolucao_passivos_fluxo"], d["resultado_por_almoxarifado"],
# d["resumo_shelf_life"]) continuam rodando em _coletar_dados_mbr sem
# alteração - NÃO foram removidos, porque outros slides do relatório (Resumo
# Executivo e Scorecard do Mês) continuam lendo dessas mesmas chaves.
#
# 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o
# modelo final"): os dois slides que viviam aqui - "Mapeamento de Risco —
# Obsolescência" (_slide_mapeamento_risco_obsolescencia) e "Scorecard de
# Mapeamento de Riscos" (_slide_scorecard_mapeamento_riscos, capítulo-síntese
# que fechava a Seção 3) - foram REMOVIDOS do relatório; ela editou o MBR
# manualmente pra tirar esses dois tópicos. d["mapeamento_risco_obsolescencia"]
# continua sendo calculado em _coletar_dados_mbr (Scorecard do Mês e Resumo
# Executivo ainda leem dessa chave, ver _linha_scorecard mais acima e KPI 2/3
# do Resumo Executivo mais abaixo) - só o slide dedicado saiu. Já
# d["scorecard_mapeamento_riscos"] e as duas funções que só existiam pra
# montar esse slide (_coletar_scorecard_mapeamento_riscos e
# _linha_risco_com_evolucao) foram removidas também - ver comentário de
# decisão mais acima, logo antes da seção "Diário de Bordo / Rotina Master".


def _slide_controle_movimentados(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Controle de Movimentados (20/08/2026, pedido do usuário) - slide
    próprio, separado de FEFO: mede a reconciliação diária sistema x físico
    (livro-caixa bruto) desde a implantação, e destaca explicitamente o
    ganho de acurácia como uma AÇÃO com impacto mensurável - não só mais um
    número solto."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Controle de Movimentados", mes_label, pagina,
               "Reconciliação diária (sistema x físico) — impacto na acurácia desde a implantação")

    resumo = d["resumo_movimentados"]
    evolucao = d["evolucao_movimentados"]
    label_mov, cor_mov = _status_maior_melhor(resumo.get("pct_acuracia"), *_LIMIARES["acuracia"])

    delta_implantacao = None
    primeiro = evolucao[0] if evolucao else None
    if primeiro and primeiro.get("pct_acuracia") is not None and resumo.get("pct_acuracia") is not None:
        delta_implantacao = round(resumo["pct_acuracia"] - primeiro["pct_acuracia"], 2)

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(resumo.get("itens_analisados")), "rotulo": "Itens Analisados no Mês"},
        {"valor": _fmt_pct(resumo.get("pct_acuracia")), "rotulo": "Acurácia da Reconciliação", "contexto": label_mov, "cor_contexto": cor_mov},
        {"valor": _fmt_num(resumo.get("itens_com_divergencia")), "rotulo": "Itens com Divergência",
         "cor": COR_ATENCAO if resumo.get("itens_com_divergencia") else COR_SUCESSO},
        {"valor": (f"+{_fmt_pct(delta_implantacao)}" if delta_implantacao is not None and delta_implantacao >= 0 else _fmt_pct(delta_implantacao)),
         "rotulo": "Ganho desde a Implantação", "cor": COR_SUCESSO if (delta_implantacao or 0) >= 0 else COR_ERRO},
    ], altura=0.85)

    # Altura da seção de topo encolhida (era 3.35/3.65) pra abrir espaço pra
    # tabela "Resultado por Almoxarifado" abaixo (22/08/2026, mockup aprovado v9).
    altura_secao_topo = 2.10
    largura_esquerda = 7.5
    if len(evolucao) >= 2:
        categorias = [_nome_mes(item["mes"])[:3] + "/" + item["mes"][2:4] for item in evolucao]
        valores = [item.get("pct_acuracia") or 0 for item in evolucao]
        _texto(slide, MARGEM_IN, 3.05, largura_esquerda, 0.28, "ACURÁCIA DA RECONCILIAÇÃO — DESDE A IMPLANTAÇÃO", tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_categoria(slide, MARGEM_IN, 3.35, largura_esquerda, altura_secao_topo, categorias, "Acurácia", valores, cor_serie=VERDE_AMAZONIA)
    else:
        _caixa_leitura(slide, MARGEM_IN, 3.05, largura_esquerda, altura_secao_topo + 0.30, "Evolução",
                       "Ainda só há um mês de histórico monitorado - a série de evolução aparece a partir do segundo mês de dados.")

    x_direita = MARGEM_IN + largura_esquerda + 0.35
    largura_direita = LARGURA_IN - MARGEM_IN - x_direita
    if primeiro and delta_implantacao is not None:
        if delta_implantacao > 0:
            texto_impacto = (
                f"Desde a implantação do Controle de Movimentados, em {_nome_mes(primeiro['mes'])}, a acurácia da reconciliação "
                f"diária subiu de {_fmt_pct(primeiro['pct_acuracia'])} para {_fmt_pct(resumo['pct_acuracia'])} — um ganho de "
                f"{_fmt_pct(delta_implantacao)}. Esse é o efeito direto de passar a reconciliar o livro-caixa bruto todo dia, em "
                "vez de só no fechamento periódico."
            )
        else:
            texto_impacto = (
                f"Desde a implantação, em {_nome_mes(primeiro['mes'])}, a acurácia da reconciliação diária está em "
                f"{_fmt_pct(resumo['pct_acuracia'])}, ainda sem ganho consistente frente ao ponto de partida "
                f"({_fmt_pct(primeiro['pct_acuracia'])}) — vale olhar os itens com divergência recorrente antes do próximo fechamento."
            )
    else:
        texto_impacto = "Ainda não há dado suficiente pra comparar com o ponto de partida da implantação."
    _caixa_leitura(slide, x_direita, 3.05, largura_direita, altura_secao_topo + 0.30, "Impacto da Implantação", texto_impacto,
                   cor_fundo=OFF_WHITE, tamanho_texto=12)

    # Resultado por Almoxarifado (22/08/2026, mockup aprovado v9) - quebra do
    # mês já real (movimentados_router.dashboard_por_almoxarifado, já existia,
    # só não estava plugada no MBR ainda).
    largura_cheia = LARGURA_IN - 2 * MARGEM_IN
    y_tabela = 3.35 + altura_secao_topo + 0.16
    _texto(slide, MARGEM_IN, y_tabela, largura_cheia, 0.22, "RESULTADO POR ALMOXARIFADO (MOVIMENTADOS)",
           tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    por_almox = d.get("movimentados_por_almoxarifado", [])
    y_zona_segura_fim = ALTURA_IN - 0.42
    if por_almox:
        ordenado = sorted(por_almox, key=lambda x: x.get("pct_acuracia") if x.get("pct_acuracia") is not None else 999)
        linhas_almox = []
        for item in ordenado[:7]:
            divergentes = (item.get("itens_analisados") or 0) - (item.get("itens_sem_divergencia") or 0)
            pct = item.get("pct_acuracia")
            _, cor_pct = _status_maior_melhor(pct, *_LIMIARES["acuracia"])
            linhas_almox.append([
                item["almoxarifado"], _fmt_num(item.get("itens_analisados")), _fmt_num(divergentes),
                (_fmt_pct(pct), cor_pct, True),
            ])
        _tabela(slide, MARGEM_IN, y_tabela + 0.24, largura_cheia, max(0.5, y_zona_segura_fim - (y_tabela + 0.24)),
                ["Almoxarifado", "Itens Analisados", "Itens c/ Divergência", "Acurácia"], linhas_almox,
                larguras_relativas=[2.0, 1.3, 1.3, 1.0], tamanho_fonte=9.5)
    else:
        _texto(slide, MARGEM_IN, y_tabela + 0.28, largura_cheia, 0.3,
               "Sem quebra por almoxarifado neste recorte.", tamanho=10, cor=CINZA_TEXTO)
    return slide


def _slide_scorecard_inventario_almoxarifado(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Scorecard de Inventário por Almoxarifado (20/08/2026, pedido do
    usuário: "traga uma análise por almoxarifado, mostrando as evoluções e
    involuções. Com um plano de ação para cada setor, baseado no histórico
    de inventários e na conciliação de movimentados"). Mesmo formato de
    tabela do Scorecard do Mês (Seção 1) - Status aqui é sobre VARIAÇÃO vs.
    o mês anterior (evolução/involução), não sobre nível absoluto (ver
    _status_evolucao) - pergunta diferente da do Painel de Inventário, que
    já mostra o nível atual por almoxarifado no slide anterior."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Scorecard de Inventário — por Almoxarifado", mes_label, pagina,
               "Evolução vs. o mês anterior e plano de ação por almoxarifado, cruzando fechamento e Controle de Movimentados")

    itens = d["scorecard_inventario_almoxarifado"]
    if not itens:
        _caixa_leitura(slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, "Sem almoxarifado cadastrado",
                       "Nenhum almoxarifado ativo encontrado no Cadastro — sem base para montar o scorecard por setor.",
                       cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13)
        return slide

    linhas_scorecard = [_linha_scorecard_almoxarifado(item) for item in itens]
    ordem_status = {"Involução": 0, "Sem histórico": 1, "Estável": 2, "Evolução": 3}
    linhas_scorecard.sort(key=lambda l: ordem_status.get(l["status_label"], 1))

    linhas_tabela = [
        [(l["frente"], AZUL_INSTITUCIONAL, True), (l["status_label"], l["status_cor"], True),
         (l["leitura"], CINZA_TEXTO, False), (l["proximo_passo"], CINZA_TEXTO, False)]
        for l in linhas_scorecard
    ]
    n = len(linhas_tabela)
    tamanho_fonte = 11.5 if n <= 8 else (10.5 if n <= 11 else 9.5)
    _tabela(slide, MARGEM_IN, 1.65, LARGURA_IN - 2 * MARGEM_IN, min(5.35, ALTURA_IN - 1.95),
            ["Almoxarifado", "Status", "Leitura Executiva", "Próximo Passo"], linhas_tabela,
            larguras_relativas=[1.6, 1.0, 3.6, 3.3], tamanho_fonte=tamanho_fonte)
    return slide


def _slide_externo_indisponivel(slide, dado: dict, nome_dashboard: str) -> bool:
    """Preenche o slide com uma mensagem de estado (não enviado / erro de leitura)
    quando um dashboard externo ainda não tem conteúdo utilizável - usado pelos
    slides de FEFO, Testes Industriais, Farol de Shelf, Recuperação de Shelf e
    Baixas Operacionais externo (20/08/2026). Retorna True se preencheu (o
    chamador deve `return slide` em seguida), False se há dado normal a exibir."""
    if not dado.get("enviado"):
        titulo = f"{nome_dashboard} ainda não enviado"
        texto = (
            f"Suba o arquivo .html autocontido do {nome_dashboard} em Auditoria > Outros Dashboards "
            "pra que ele apareça neste relatório."
        )
        cor = COR_ATENCAO
    elif dado.get("erro_extracao"):
        titulo = f"Não foi possível ler o {nome_dashboard}"
        texto = (
            "O arquivo enviado não pôde ser processado (formato inesperado) — reenvie o .html "
            "em Auditoria > Outros Dashboards."
        )
        cor = COR_ERRO
    else:
        return False
    _caixa_leitura(slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, titulo, texto,
                   cor_fundo=OFF_WHITE, cor_rotulo=cor, tamanho_texto=13)
    return True


def _slide_fefo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """FEFO (atualizado 20/08/2026, pedido do usuário: "use o arquivo HTML pra
    alimentar a construção do MBR no módulo FEFO, o mesmo tem mais informações e
    bases de registro" -> confirmado "Auditoria FEFO importada"). Fonte trocada do
    dashboard "Controle de FEFO" (Auditoria > Outros Dashboards) pra AuditoriaFefo
    - o histórico importado direto na tela FEFO (painel "Auditoria FEFO — histórico
    importado"), com lote movimentado/validade/lote mais antigo disponível por
    registro (ver fefo.calcular_resumo_auditoria_fefo e
    _extrair_resumo_auditoria_fefo acima). NÃO reaproveita _slide_externo_indisponivel
    porque essa fonte não vem de um DashboardExterno - o caminho de importação é a
    tela FEFO, não Auditoria > Outros Dashboards."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "FEFO", mes_label, pagina, "Dados reais da Auditoria FEFO importada — filtrado pelo mês deste relatório")

    fefo = d["fefo_externo"]

    if not fefo.get("enviado"):
        _caixa_leitura(
            slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, "Auditoria FEFO ainda não importada",
            "Nenhum histórico foi importado ainda no painel \"Auditoria FEFO — histórico importado\" da tela FEFO "
            "(Excel diário ou dashboard HTML consolidado) — sem esse histórico, o FEFO não entra neste relatório "
            "com dado real.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13,
        )
        return slide

    if not fefo.get("tem_dados"):
        _caixa_leitura(
            slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, f"Sem movimentos auditáveis em {mes_label}",
            "A Auditoria FEFO importada não tem movimentos auditáveis registrados para este mês — confira se o "
            "histórico do período está importado no painel \"Auditoria FEFO — histórico importado\" da tela FEFO.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13,
        )
        return slide

    label_fefo, cor_fefo = _status_menor_melhor(fefo["taxa_quebra_pct"], *_LIMIARES["fefo_quebra_pct"])
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(fefo["total_auditaveis"]), "rotulo": "Movimentos Auditáveis no Mês", "cor": COR_INFO},
        {"valor": _fmt_num(fefo["total_quebras"]), "rotulo": "Quebras de FEFO", "cor": COR_ERRO if fefo["total_quebras"] else COR_SUCESSO},
        {"valor": _fmt_pct(fefo["taxa_quebra_pct"]), "rotulo": "Taxa de Quebra", "cor": cor_fefo, "contexto": label_fefo, "cor_contexto": cor_fefo},
        {"valor": _fmt_num(fefo.get("total_sem_correspondencia")), "rotulo": "Sem Correspondência no Mês"},
    ], altura=0.85)

    top = fefo.get("top_produtos_com_quebra") or []
    if top:
        categorias = [t["produto"][:26] for t in reversed(top)]
        valores = [t["quebras"] for t in reversed(top)]
        _texto(slide, MARGEM_IN, 3.05, 6.3, 0.26, "PRODUTOS COM MAIS QUEBRAS NO MÊS", tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_categoria(slide, MARGEM_IN, 3.35, 6.3, 2.55, categorias, "Quebras", valores,
                            tipo=XL_CHART_TYPE.BAR_CLUSTERED, cor_serie=COR_ERRO, formato_numero='0')
    else:
        _caixa_leitura(slide, MARGEM_IN, 3.05, 6.3, 2.55, "Produtos com mais quebras",
                        "Nenhuma quebra registrada neste mês.")

    x_direita = MARGEM_IN + 6.3 + 0.35
    largura_direita = LARGURA_IN - MARGEM_IN - x_direita
    _texto(slide, x_direita, 3.05, largura_direita, 0.26, "QUEBRAS POR DESTINO", tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
    por_destino = fefo.get("top_destinos_com_quebra") or []
    if por_destino:
        linhas = [[pd["destino"][:30], _fmt_num(pd["quebras"])] for pd in por_destino]
        _tabela(slide, x_direita, 3.35, largura_direita, 2.55, ["Destino", "Quebras"], linhas,
                larguras_relativas=[2.6, 1.0], tamanho_fonte=10.5)
    else:
        _caixa_leitura(slide, x_direita, 3.35, largura_direita, 2.25, "Por destino", "Sem dados de destino neste mês.")

    fontes = fefo.get("fontes_no_periodo") or []
    fontes_label = " + ".join(
        "auditoria diária" if f == "auditoria_diaria" else "dashboard consolidado" if f == "dashboard_consolidado" else f
        for f in fontes
    ) or "—"
    _texto(
        slide, MARGEM_IN, 6.15, LARGURA_IN - 2 * MARGEM_IN, 0.5,
        f"Fonte: Auditoria FEFO importada (tela FEFO), origem no mês: {fontes_label} — "
        f"{_fmt_num(fefo.get('total_sem_correspondencia'))} movimento(s) sem correspondência neste mês não entram na taxa de quebra.",
        tamanho=10, cor=CINZA_TEXTO,
    )
    return slide


def _slide_testes_industriais(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Testes Industriais (20/08/2026, pedido do usuário: "adicionar todos os
    anexos na construção do MBR") - dashboard de Controle de Testes Industriais
    tem data por registro (campo "mes"), então filtra exato pelo mês deste
    relatório, igual ao FEFO (ver dashboards_externos_extrator.extrair_testes_industriais)."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Testes Industriais", mes_label, pagina,
               "Dados reais do dashboard de Controle de Testes Industriais — filtrado pelo mês deste relatório")

    dado = d["testes_industriais_externo"]
    if _slide_externo_indisponivel(slide, dado, "Controle de Testes Industriais"):
        return slide

    if not dado.get("tem_dados"):
        _caixa_leitura(
            slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, f"Sem testes registrados em {mes_label}",
            "O dashboard de Controle de Testes Industriais enviado não tem itens consumidos para este mês — confira "
            "se o arquivo está atualizado em Auditoria > Outros Dashboards.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13,
        )
        return slide

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(dado["total_itens"]), "rotulo": "Itens Consumidos no Mês"},
        {"valor": _fmt_moeda(dado["gasto_total"]), "rotulo": "Gasto Total no Mês", "cor": COR_ATENCAO},
        {"valor": _fmt_num(dado["ops"]), "rotulo": "OPs Testadas"},
        {"valor": _fmt_moeda(dado["custo_medio_op"]), "rotulo": "Custo Médio por OP"},
    ], altura=0.85)

    # 02/09/2026 (decisão do usuário sobre a nova lógica do indicador): itens
    # sem custo cadastrado são excluídos do Gasto Total/Custo Médio por OP
    # (não zerar/distorcer esses números), mas avisa quantos ficaram de fora
    # em vez de simplesmente sumir com eles sem explicação.
    itens_sem_custo = dado.get("itens_sem_custo") or 0
    if itens_sem_custo:
        _texto(
            slide, MARGEM_IN, 2.55, LARGURA_IN - 2 * MARGEM_IN, 0.26,
            f"Obs.: {_fmt_num(itens_sem_custo)} item(ns) consumido(s) no mês sem custo cadastrado foram excluídos "
            "deste total — confira o cadastro de custo desses materiais.",
            tamanho=10, cor=COR_ATENCAO, negrito=True,
        )

    top = dado.get("top_materias_primas") or []
    if top:
        categorias = [t["nome"][:26] for t in reversed(top)]
        valores = [t["custo"] for t in reversed(top)]
        _texto(slide, MARGEM_IN, 3.05, LARGURA_IN - 2 * MARGEM_IN, 0.26, "MATÉRIAS-PRIMAS COM MAIOR CUSTO NO MÊS (R$)",
               tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_categoria(slide, MARGEM_IN, 3.35, LARGURA_IN - 2 * MARGEM_IN, 3.0, categorias, "Custo", valores,
                            tipo=XL_CHART_TYPE.BAR_CLUSTERED, cor_serie=COR_ATENCAO, formato_numero='#,##0')
    else:
        _caixa_leitura(slide, MARGEM_IN, 3.05, LARGURA_IN - 2 * MARGEM_IN, 3.0, "Matérias-primas",
                        "Nenhuma matéria-prima consumida neste mês.")

    _texto(
        slide, MARGEM_IN, 6.55, LARGURA_IN - 2 * MARGEM_IN, 0.4,
        f"Fonte: dashboard de Controle de Testes Industriais (Auditoria > Outros Dashboards), enviado em {dado.get('enviado_em') or '—'}.",
        tamanho=10, cor=CINZA_TEXTO,
    )
    return slide


def _slide_dispersao_ficha_tecnica(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Dispersão de Ficha Técnica (20/08/2026, pedido do usuário: "Adicione
    Dispersão de Ficha técnica na apresentação do MBR") - dashboard
    "Dispersão de Lote — Produção" (Ficha Técnica × Consumo real por Ordem
    de Produção), enviado como indicador dinâmico em Auditoria > Outros
    Dashboards. Tem data por registro (campo "mes"), então filtra exato
    pelo mês deste relatório, igual a FEFO/Testes Industriais (ver
    dashboards_externos_extrator.extrair_dispersao_ficha_tecnica). Sem
    limiar de status definido pelo usuário pra taxa de furo (ao contrário
    de FEFO/_LIMIARES) - mostra o valor real sem rótulo "Em avanço/
    Atenção/Crítico" pra não inventar uma meta que não foi combinada."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Dispersão de Ficha Técnica", mes_label, pagina,
               "Ficha Técnica × Consumo real por Ordem de Produção — filtrado pelo mês deste relatório")

    dado = d["dispersao_ficha_tecnica_externo"]
    if _slide_externo_indisponivel(slide, dado, "Dispersão de Ficha Técnica"):
        return slide

    if not dado.get("tem_dados"):
        _caixa_leitura(
            slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, f"Sem OPs analisadas em {mes_label}",
            "O dashboard de Dispersão de Ficha Técnica enviado não tem Ordens de Produção registradas para este mês "
            "— confira se o arquivo está atualizado em Auditoria > Outros Dashboards.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13,
        )
        return slide

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(dado["ops_analisadas"]), "rotulo": "OPs Analisadas no Mês"},
        {"valor": _fmt_pct(dado["taxa_furo_pct"]), "rotulo": "Taxa de Furo",
         "cor": COR_ATENCAO if (dado["taxa_furo_pct"] or 0) > 0 else COR_SUCESSO},
        {"valor": _fmt_moeda(dado["impacto_liquido"]), "rotulo": "Impacto Líquido (Perda − Economia)",
         "cor": COR_ERRO if (dado["impacto_liquido"] or 0) > 0 else COR_SUCESSO},
        {"valor": _fmt_num(dado["ops_criticas"]), "rotulo": "OPs Críticas",
         "cor": COR_ERRO if dado["ops_criticas"] else COR_SUCESSO},
    ], altura=0.85)

    # Tendência Financeira mês a mês (22/08/2026, pedido do usuário: "adicione
    # rótulo de dados no indicador de Dispersão de Ficha Técnica") - Perda,
    # Economia e Impacto Líquido, últimos 6 meses até o mês do relatório
    # (extrator agrega isso com exatidão a partir do JSON embutido - ver
    # docstring de extrair_dispersao_ficha_tecnica). Rótulos de valor já
    # vêm de fábrica em _grafico_categoria_multi (has_data_labels=True).
    evolucao = (dado.get("evolucao_mensal") or [])[-6:]
    if len(evolucao) >= 2:
        categorias_evol = [_nome_mes(item["mes"])[:3] + "/" + item["mes"][2:4] for item in evolucao]
        _texto(slide, MARGEM_IN, 2.90, LARGURA_IN - 2 * MARGEM_IN, 0.24,
               "TENDÊNCIA FINANCEIRA — PERDA × ECONOMIA × IMPACTO LÍQUIDO (R$/MÊS)",
               tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_categoria_multi(
            slide, MARGEM_IN, 3.18, LARGURA_IN - 2 * MARGEM_IN, 1.42, categorias_evol,
            [("Perda", [i["perda"] for i in evolucao], COR_ERRO),
             ("Economia", [i["economia"] for i in evolucao], COR_SUCESSO),
             ("Impacto Líquido", [i["impacto_liquido"] for i in evolucao], AZUL_INSTITUCIONAL)],
            formato_numero='#,##0',
        )
        y_secao_seguinte = 4.78
    else:
        y_secao_seguinte = 3.05

    top_perda = dado.get("top_materiais_perda") or []
    altura_secao = 6.35 - y_secao_seguinte
    if top_perda:
        amostra = top_perda[:6]
        categorias = [t["descricao"][:26] for t in reversed(amostra)]
        valores = [t["impacto"] for t in reversed(amostra)]
        _texto(slide, MARGEM_IN, y_secao_seguinte, 6.3, 0.24, "MATERIAIS COM MAIOR PERDA NO MÊS (R$)", tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        _grafico_categoria(slide, MARGEM_IN, y_secao_seguinte + 0.28, 6.3, altura_secao - 0.28, categorias, "Perda", valores,
                            tipo=XL_CHART_TYPE.BAR_CLUSTERED, cor_serie=COR_ERRO, formato_numero='#,##0')
    else:
        _caixa_leitura(slide, MARGEM_IN, y_secao_seguinte, 6.3, altura_secao, "Materiais com maior perda",
                        "Nenhum material com perda líquida neste mês.")

    x_direita = MARGEM_IN + 6.3 + 0.35
    largura_direita = LARGURA_IN - MARGEM_IN - x_direita
    _texto(slide, x_direita, y_secao_seguinte, largura_direita, 0.24, "MATERIAIS COM MAIOR ECONOMIA", tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
    top_economia = dado.get("top_materiais_economia") or []
    if top_economia:
        linhas = [[te["descricao"][:24], _fmt_num(te["ops"]), _fmt_moeda(-te["impacto"])] for te in top_economia[:6]]
        _tabela(slide, x_direita, y_secao_seguinte + 0.28, largura_direita, altura_secao - 0.28, ["Material", "OPs", "Economia"], linhas,
                larguras_relativas=[2.4, 0.8, 1.4], tamanho_fonte=10)
    else:
        _caixa_leitura(slide, x_direita, y_secao_seguinte + 0.28, largura_direita, altura_secao - 0.28, "Economia",
                        "Sem economia líquida registrada neste mês.")

    _texto(
        slide, MARGEM_IN, 6.45, LARGURA_IN - 2 * MARGEM_IN, 0.5,
        f"Fonte: dashboard Dispersão de Ficha Técnica (Auditoria > Outros Dashboards), enviado em {dado.get('enviado_em') or '—'} — "
        f"Materiais crônicos (≥ {dado.get('limiar_freq_ops', 5)} OPs): {_fmt_num(dado.get('materiais_cronicos'))} · "
        f"Concentração Top 20: {_fmt_pct(dado.get('concentracao_top20_pct'))} do impacto absoluto.",
        tamanho=10, cor=CINZA_TEXTO,
    )
    return slide


def _bucket_farol(buckets: list, *rotulos: str):
    """Acha o bucket do Farol de Shelf pelo texto do título da tabela de
    origem (ex.: "Top 10 — Urgente (0-30 dias)") - mais robusto que confiar
    na ORDEM em que dashboards_externos_extrator.extrair_farol_shelf varreu
    as tabelas do HTML."""
    for b in buckets:
        titulo = (b.get("titulo") or "").lower()
        if any(r.lower() in titulo for r in rotulos):
            return b
    return None


# Ordem/cor fixa dos status do Farol de Shelf-Life (usada nos dois gráficos
# de detalhamento por almoxarifado/grupo, Fase 2 22/08/2026) - a MESMA ordem
# de severidade dos cartões deste mesmo slide, pra quem olhar os cartões e os
# gráficos reconhecer a cor sem precisar reler a legenda.
_CORES_STATUS_FAROL = {
    "Vencido": COR_FAROL_VENCIDO,
    "Urgente": COR_FAROL_URGENTE,
    "Perigo": COR_FAROL_PERIGO,
    "Atenção": COR_FAROL_ATENCAO,
}
_ORDEM_STATUS_FAROL = ["Vencido", "Urgente", "Perigo", "Atenção"]


def _slide_farol_shelf_externo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Farol de Shelf-Life (20/08/2026, reestruturado em 22/08/2026 - pedido
    do usuário: "mudar a estrutura pra caber tudo no indicador... cards com
    a quantidade de lotes em aberto e uma representação do farol em grana:
    0-30 dias e valor total em aberto, 31-60 e valor total, 61-90") - dashboard
    é uma FOTO do estoque no momento da exportação (sem dimensão de mês:
    lotes com saldo agora e validade em até 90 dias), não um recorte do mês
    deste relatório - entra como retrato datado, rotulado com a data real da
    exportação (decisão do usuário), não filtrado pelo mês (ver
    dashboards_externos_extrator.extrair_farol_shelf).

    NOTA (21/08/2026, Fase 3 - pedido da usuária com o arquivo
    "MBR_Atlas_202607_15.pptx" em anexo, "deixe os slides ... exatamente
    assim conforme anexo. Exclua o slide com indicador individual anexado"):
    as duas visões "Risco por Almoxarifado" e "Custo por Grupo e Status", que
    na Fase 2 tinham ido pro slide companheiro _slide_farol_shelf_risco_
    almoxarifado (por não caberem aqui junto com os cards + 3 tabelas),
    voltaram pra ESTE slide, medidas e posicionadas por geometria exata do
    arquivo que a usuária anexou - o slide companheiro foi removido. O
    gráfico "Custo por Grupo e Status" também passou de barra empilhada em
    R$ (COLUMN_STACKED) pra barra empilhada em 100% (BAR_STACKED_100),
    exatamente como no anexo - lê-se como proporção de cada status dentro do
    grupo, não como custo absoluto (esse já está nos cartões e nas 3
    tabelas)."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Farol de Shelf-Life", mes_label, pagina,
               "Retrato do estoque em risco de validade — dados reais do dashboard de Farol de Shelf-Life")

    dado = d["farol_shelf_externo"]
    if _slide_externo_indisponivel(slide, dado, "Farol de Shelf-Life"):
        return slide

    qtd = dado.get("qtd_lotes") or {}
    buckets = dado.get("buckets") or []
    bucket_urgente = _bucket_farol(buckets, "urgente")
    bucket_perigo = _bucket_farol(buckets, "perigo")
    bucket_atencao = _bucket_farol(buckets, "atenção", "atencao")

    def _valor_bucket(bucket):
        return bucket["total"]["custo"] if bucket and bucket.get("total") else None

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(qtd.get("vencidos")), "rotulo": "Lotes Já Vencidos",
         "cor": COR_FAROL_VENCIDO, "contexto": _fmt_moeda(dado.get("perda_ja_vencida")), "cor_contexto": COR_FAROL_VENCIDO},
        {"valor": _fmt_num(qtd.get("0_30")), "rotulo": "0-30 Dias (Urgente)",
         "cor": COR_FAROL_URGENTE, "contexto": _fmt_moeda(_valor_bucket(bucket_urgente)), "cor_contexto": COR_FAROL_URGENTE},
        {"valor": _fmt_num(qtd.get("31_60")), "rotulo": "31-60 Dias (Perigo)",
         "cor": COR_FAROL_PERIGO, "contexto": _fmt_moeda(_valor_bucket(bucket_perigo)), "cor_contexto": COR_FAROL_PERIGO},
        {"valor": _fmt_num(qtd.get("61_90")), "rotulo": "61-90 Dias (Atenção)",
         "cor": COR_FAROL_ATENCAO, "contexto": _fmt_moeda(_valor_bucket(bucket_atencao)), "cor_contexto": COR_FAROL_ATENCAO},
    ], altura=0.85)

    # Risco por Almoxarifado (barra empilhada em R$) e Custo por Grupo e
    # Status (barra empilhada em 100%), lado a lado - geometria (posição e
    # tamanho) baseada no arquivo "MBR_Atlas_202607_15.pptx" que a usuária
    # anexou como referência (21/08/2026, Fase 3). Sem título de texto acima
    # de cada gráfico - o anexo não tem, a legenda + eixo já identificam
    # cada um.
    #
    # O anexo tem os cartões de KPI mais alto na página (y=1,27in) que este
    # gerador (y=1,55in, ver _linha_kpis abaixo - posição já usada por todo
    # o resto do relatório, não é algo que esta rodada pediu pra mudar) -
    # y_grafico e altura_grafico foram ajustados (não copiados 1:1 do anexo)
    # pra abrir espaço sem sobrepor os cartões, mas mantendo o rodapé do
    # gráfico exatamente onde o anexo tem (y=4,90in, colado no título da
    # 1ª linha de tabelas Top 10 abaixo).
    largura_grafico = (LARGURA_IN - 2 * MARGEM_IN - 0.4) / 2
    x2_grafico = MARGEM_IN + largura_grafico + 0.4
    y_grafico, altura_grafico = 2.50, 2.40

    risco_almox = dado.get("risco_por_almoxarifado")
    custo_grupo = dado.get("custo_por_grupo_status")

    if risco_almox:
        categorias = sorted(risco_almox["totais"], key=lambda c: -risco_almox["totais"][c])
        series = []
        for status in _ORDEM_STATUS_FAROL:
            valores_status = risco_almox["series"].get(status)
            if not valores_status:
                continue
            series.append((status, [valores_status.get(c, 0.0) for c in categorias],
                           _CORES_STATUS_FAROL.get(status, COR_SEM_DADO)))
        _grafico_categoria_multi(slide, MARGEM_IN, y_grafico, largura_grafico, altura_grafico, categorias, series,
                                  tipo=XL_CHART_TYPE.COLUMN_STACKED, formato_numero='R$ #,##0', mostrar_rotulos=False)
    else:
        _caixa_leitura(slide, MARGEM_IN, y_grafico, largura_grafico, altura_grafico, "Sem dados",
                       "Risco por Almoxarifado não disponível neste retrato.", tamanho_texto=11)

    if custo_grupo:
        grupos_ordenados = sorted(custo_grupo, key=lambda g: -g["total"])
        categorias_grupo = [g["grupo"] for g in grupos_ordenados]
        series_grupo = []
        for status in _ORDEM_STATUS_FAROL:
            valores_status = [g["por_status"].get(status, {}).get("valor", 0.0) for g in grupos_ordenados]
            if not any(valores_status):
                continue
            series_grupo.append((status, valores_status, _CORES_STATUS_FAROL.get(status, COR_SEM_DADO)))
        _grafico_categoria_multi(slide, x2_grafico, y_grafico, largura_grafico, altura_grafico, categorias_grupo,
                                  series_grupo, tipo=XL_CHART_TYPE.BAR_STACKED_100, mostrar_rotulos=False)
    else:
        _caixa_leitura(slide, x2_grafico, y_grafico, largura_grafico, altura_grafico, "Sem dados",
                       "Custo por Grupo e Status não disponível neste retrato.", tamanho_texto=11)

    # As 3 listas Top 10 (uma por faixa), lado a lado - "abaixo as tabelas
    # como estão" (pedido do usuário) - encolhidas de 8 pra 5 linhas cada pra
    # caber tudo (cards + 2 gráficos + 3 tabelas) numa página só, com
    # posição/altura também medidas do anexo (21/08/2026, Fase 3).
    colunas = [
        ("0-30 DIAS — URGENTE", bucket_urgente, COR_FAROL_URGENTE),
        ("31-60 DIAS — PERIGO", bucket_perigo, COR_FAROL_PERIGO),
        ("61-90 DIAS — ATENÇÃO", bucket_atencao, COR_FAROL_ATENCAO),
    ]
    gap = 0.3
    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 2 * gap) / 3
    y_titulo, y_tabela, altura_tabela = 4.902, 5.142, 1.891
    algum_bucket_com_itens = False
    for i, (titulo, bucket, cor) in enumerate(colunas):
        x = MARGEM_IN + i * (largura_col + gap)
        _texto(slide, x, y_titulo, largura_col, 0.24, titulo, tamanho=10, negrito=True, cor=cor)
        if bucket and bucket.get("itens"):
            algum_bucket_com_itens = True
            linhas = [[it["descricao"][:26], _fmt_moeda(it["custo"])] for it in bucket["itens"][:5]]
            _tabela(slide, x, y_tabela, largura_col, altura_tabela, ["Descrição", "Custo"], linhas,
                    larguras_relativas=[2.2, 1.0], tamanho_fonte=9)
            if bucket.get("total"):
                _texto(slide, x, y_tabela + altura_tabela + 0.06, largura_col, 0.22,
                       f"Total: {_fmt_moeda(bucket['total']['custo'])}", tamanho=9, negrito=True, cor=CINZA_TEXTO)
        else:
            _caixa_leitura(slide, x, y_tabela, largura_col, altura_tabela, "Sem itens",
                            "Nenhum lote nesta faixa.", tamanho_texto=10)
    if not algum_bucket_com_itens:
        _caixa_leitura(slide, MARGEM_IN, y_tabela, LARGURA_IN - 2 * MARGEM_IN, altura_tabela, "Lotes em risco",
                        "Nenhum lote em risco de validade neste retrato.")
    return slide


# Ordem/cor fixa das 3 séries do gráfico de evolução mensal da Recuperação
# de Shelf (Fase 2, 22/08/2026) - mesma paleta de sucesso/erro já usada nos
# cartões de KPI deste mesmo slide (Perda Real em vermelho, Receita
# Recuperada e Saving Recuperado em verde).
_CORES_SERIE_RECUPERACAO = {
    "Perda": COR_ERRO,
    "Receita Recuperada": COR_SUCESSO,
    "Saving Recuperado": VERDE_AMAZONIA,
}


def _slide_recuperacao_shelf_externo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Recuperação de Shelf (20/08/2026) - dashboard cobre um período agregado
    (ex.: jan-ago/26 no arquivo de exemplo), não um mês calendário isolado -
    entra como retrato datado do período real do arquivo (decisão do usuário),
    com os KPIs de metodologia de recuperação financeira e os 2 rankings Top 10
    reais (ver dashboards_externos_extrator.extrair_recuperacao_shelf).

    NOTA (21/08/2026, Fase 3 - mesmo pedido da usuária descrito em
    _slide_farol_shelf_externo, com o arquivo "MBR_Atlas_202607_15.pptx" em
    anexo): o gráfico "Evolução Mensal" (Perda × Receita Recuperada × Saving
    Recuperado por mês), que na Fase 2 tinha ido pro slide companheiro
    _slide_recuperacao_shelf_evolucao, voltou pra ESTE slide - largura
    cheia, entre os cartões e as 2 tabelas, geometria medida do anexo. O
    slide companheiro foi removido. Os cartões deste slide não têm uma 3ª
    linha de "contexto" (só valor + rótulo), então usam uma altura mais
    baixa que os demais KPIs do relatório (0,65in contra 0,85in) e o rótulo
    cola direto embaixo do valor - ver `deslocamento_rotulo` em
    _cartao_kpi/_linha_kpis, e a nota de decisão registrada ali. A nota
    "ROI operacional" e o texto de retrato datado que existiam no rodapé
    deste slide não estão no anexo (não há mais espaço com o gráfico
    ocupando a largura toda) - removidos; o ROI continua calculado e
    disponível em `kpis["roi_operacional_pct"]" para quem for reutilizar o
    dado em outro lugar."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Recuperação de Shelf", mes_label, pagina,
               "Retrato do período — dados reais do dashboard de Recuperação de Shelf (perda × recuperação por ações de lote)")

    dado = d["recuperacao_shelf_externo"]
    if _slide_externo_indisponivel(slide, dado, "Recuperação de Shelf"):
        return slide

    kpis = dado.get("kpis") or {}
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_moeda(kpis.get("receita_recuperada")), "rotulo": "Receita Recuperada", "cor": COR_SUCESSO},
        {"valor": _fmt_moeda(kpis.get("perda_evitada")), "rotulo": "Perda Evitada", "cor": COR_SUCESSO},
        {"valor": _fmt_moeda(kpis.get("perda_real")), "rotulo": "Perda Real", "cor": COR_ERRO},
        {"valor": _fmt_moeda(kpis.get("saving_recuperado")), "rotulo": "Saving Recuperado"},
    ], altura=0.65, deslocamento_rotulo=0.383, tamanho_valor_base=20)

    # Evolução Mensal (Perda × Receita Recuperada × Saving Recuperado),
    # largura cheia, entre os cartões e as 2 tabelas - geometria medida do
    # arquivo "MBR_Atlas_202607_15.pptx" (21/08/2026, Fase 3).
    evolucao = dado.get("evolucao_mensal")
    if evolucao:
        categorias = evolucao["categorias"]
        series = [
            (nome, [evolucao["series"][nome].get(c, 0.0) for c in categorias], cor)
            for nome, cor in _CORES_SERIE_RECUPERACAO.items()
            if nome in evolucao["series"]
        ]
        _grafico_categoria_multi(slide, MARGEM_IN, 2.185, LARGURA_IN - 2 * MARGEM_IN, 2.231, categorias, series,
                                  tipo=XL_CHART_TYPE.COLUMN_CLUSTERED, formato_numero='R$ #,##0')
    else:
        _caixa_leitura(
            slide, MARGEM_IN, 2.185, LARGURA_IN - 2 * MARGEM_IN, 2.231,
            "Evolução mensal não disponível neste retrato",
            "Não foi possível reconstruir o gráfico de evolução mensal a partir do arquivo exportado desta "
            "vez (o valor recalculado da série \"Perda\" não bateu com o KPI \"Perda Real\" do próprio "
            "arquivo, ou o layout do export mudou) — os KPIs e tabelas abaixo continuam confiáveis.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=12,
        )

    tabelas = dado.get("tabelas") or {}
    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 0.35) / 2
    x2 = MARGEM_IN + largura_col + 0.35
    y_tabelas = 4.403
    for i, (titulo, tabela) in enumerate(list(tabelas.items())[:2]):
        x = MARGEM_IN if i == 0 else x2
        _texto(slide, x, y_tabelas, largura_col, 0.26, titulo.upper(), tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        linhas = [[c[:22] if isinstance(c, str) else c for c in linha] for linha in tabela["linhas"][:6]]
        larguras = [1.6] + [1.0] * (len(tabela["cabecalho"]) - 1) if tabela["cabecalho"] else None
        _tabela(slide, x, y_tabelas + 0.26, largura_col, 2.43, tabela["cabecalho"], linhas,
                larguras_relativas=larguras, tamanho_fonte=9.5)

    return slide


def _slide_baixas_operacionais_externo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Dashboard Baixas Operacionais externo (20/08/2026 - criado como
    controle paralelo; 22/08/2026 - pedido não-negociável da usuária,
    promovido a FONTE OFICIAL do assunto "Passivos/Baixas" no MBR,
    substituindo os slides nativos _slide_mapeamento_passivos e
    _slide_passivos_evolucao, removidos - ver comentário de decisão logo
    antes de _slide_controle_movimentados). Usa categorização de
    motivo própria da equipe (Stock Savvy), diferente da categorização
    nativa do Atlas que os slides removidos usavam. Cobre uma janela móvel
    (ex.: últimos ~60 dias no arquivo de exemplo), não um mês calendário
    isolado - entra como retrato datado (decisão do usuário)
    (ver dashboards_externos_extrator.extrair_baixas_operacionais_externo).

    NOTA (21/08/2026, Fase 3 - pedido da usuária: "só faltou unificar uma
    análise. O pacote de baixas operacionais, por gentileza, replique a
    mesma lógica para o KPI. Slide 14 e 15"): o gráfico "Total de Baixas
    por Mês", que na Fase 2 tinha ido pro slide companheiro
    _slide_baixas_operacionais_evolucao (não cabia junto com os KPIs + 2
    tabelas já existentes aqui), voltou pra ESTE slide - mesma fusão já
    feita com Farol de Shelf-Life e Recuperação de Shelf. O slide
    companheiro foi removido. Os cartões de KPI (que já não tinham 3ª
    linha de contexto) ganharam a mesma "lógica compacta" da Recuperação de
    Shelf (altura 0,65in, número em 20pt - ver `deslocamento_rotulo`/
    `tamanho_valor_base` em _cartao_kpi) - "replique a mesma lógica para o
    KPI", nas palavras da usuária. Cobre o histórico mensal completo do
    export (ex.: janeiro-agosto) - uma janela BEM mais longa que o
    "Prejuízo Total no Período" dos cartões (janela móvel curta, ex.
    últimos ~60 dias); os dois números não batem entre si por desenho
    (janelas diferentes), não é inconsistência - nota no rodapé abaixo."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    # Título/subtítulo atualizados em 22/08/2026 - "(Controle Paralelo)" e
    # "mantido em paralelo ao módulo nativo" não fazem mais sentido depois
    # que este dashboard passou a ser a fonte OFICIAL (nativo removido).
    _cabecalho(slide, "Dashboard Baixas Operacionais", mes_label, pagina,
               "Retrato do período — fonte oficial de Baixas/Passivos no MBR (dashboard externo, modelo aprovado)")

    dado = d["baixas_operacionais_externo"]
    if _slide_externo_indisponivel(slide, dado, "Dashboard Baixas Operacionais"):
        return slide

    resumo = dado.get("resumo") or {}
    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_moeda(resumo.get("prejuizo_total")), "rotulo": "Prejuízo Total no Período", "cor": COR_ERRO},
        {"valor": _fmt_pct(resumo.get("pct_concentrado")), "rotulo": f"Concentrado em {resumo.get('motivo_concentrado', '—')}"},
        {"valor": resumo.get("setor_maior_impacto") or "—", "rotulo": "Setor de Maior Impacto"},
        {"valor": resumo.get("grupo_maior_impacto") or "—", "rotulo": "Grupo de Maior Impacto"},
    ], altura=0.65, deslocamento_rotulo=0.383, tamanho_valor_base=20)

    # Total de Baixas por Mês, largura cheia, entre os cartões e as 2
    # tabelas (21/08/2026, Fase 3).
    evolucao = dado.get("evolucao_mensal")
    y_grafico, altura_grafico = 2.32, 1.95
    if evolucao:
        categorias = evolucao["categorias"]
        nomes_series = list(evolucao["series"].keys())
        empilhado = len(nomes_series) > 1
        paleta = [COR_ERRO, COR_ATENCAO, AZUL_INSTITUCIONAL, VERDE_AMAZONIA, COR_INFO,
                  COR_FAROL_URGENTE, COR_FAROL_PERIGO, AZUL_CLARO, COR_SEM_DADO]
        series = [
            (nome, [evolucao["series"][nome].get(c, 0.0) for c in categorias], paleta[i % len(paleta)])
            for i, nome in enumerate(nomes_series)
        ]
        # Empilhado por motivo tem até 9 séries por coluna - rótulo por
        # segmento fica ilegível/sobreposto (visto na QA visual da Fase 2),
        # então continua desligado nesse caso. Em vez disso (21/08/2026,
        # pedido explícito da usuária: "adicione o rótulo do valor total
        # por gráfico, não por motivo... a somatório dos eventos por mês"),
        # um rótulo de TOTAL por coluna - ver _adicionar_rotulo_total_
        # empilhado. Com série única (coluna simples, quando o export não
        # tem quebra por motivo) o rótulo por coluna já é o total, então
        # continua pelo caminho normal (mostrar_rotulos=True).
        chart = _grafico_categoria_multi(
            slide, MARGEM_IN, y_grafico, LARGURA_IN - 2 * MARGEM_IN, altura_grafico, categorias, series,
            tipo=XL_CHART_TYPE.COLUMN_STACKED if empilhado else XL_CHART_TYPE.COLUMN_CLUSTERED,
            formato_numero='R$ #,##0', mostrar_rotulos=not empilhado)
        if empilhado:
            totais = [sum(evolucao["series"][nome].get(c, 0.0) for nome in nomes_series) for c in categorias]
            _adicionar_rotulo_total_empilhado(chart, categorias, totais,
                                               cor_texto_hex=_hex_cor(CINZA_TEXTO), formato_numero='R$ #,##0')
    else:
        _caixa_leitura(
            slide, MARGEM_IN, y_grafico, LARGURA_IN - 2 * MARGEM_IN, altura_grafico,
            "Evolução mensal não disponível neste retrato",
            "Não foi possível reconstruir o gráfico de evolução mensal a partir do arquivo exportado desta "
            "vez (eixo do gráfico fora do padrão esperado, ou o layout do export mudou) — os KPIs e tabelas "
            "abaixo continuam confiáveis.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=12,
        )

    tabelas = dado.get("tabelas") or {}
    tabela_motivo = tabelas.get("Baixas por Motivo")
    tabela_sku = tabelas.get("Ranking de SKU — Top 10 Baixas")
    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 0.35) / 2
    x2 = MARGEM_IN + largura_col + 0.35
    y_tabelas, altura_tabela = 4.37, 2.05
    if tabela_sku:
        _texto(slide, MARGEM_IN, y_tabelas, largura_col, 0.24, "TOP BAIXAS POR SKU", tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        linhas = [[c[:24] if isinstance(c, str) else c for c in linha] for linha in tabela_sku["linhas"][:6]]
        # Coluna "Ranking" alargada de 0.6 pra 1.0 (21/08/2026, Fase 3):
        # nesta altura de tabela mais compacta, o cabeçalho "Ranking" não
        # cabia mais numa linha só (quebrava em "Rankin"/"g"), o que
        # inflava a linha de cabeçalho e estourava a altura reservada pra
        # tabela, invadindo o rodapé abaixo - visto na QA visual. Largura
        # total das 4 colunas mantida igual (5.4), só redistribuída.
        _tabela(slide, MARGEM_IN, y_tabelas + 0.24, largura_col, altura_tabela, tabela_sku["cabecalho"], linhas,
                larguras_relativas=[1.0, 2.2, 1.1, 1.1], tamanho_fonte=9.5)
    if tabela_motivo:
        _texto(slide, x2, y_tabelas, largura_col, 0.24, "BAIXAS POR MOTIVO", tamanho=10.5, negrito=True, cor=AZUL_INSTITUCIONAL)
        linhas = [[c[:22] if isinstance(c, str) else c for c in linha] for linha in tabela_motivo["linhas"][:6]]
        _tabela(slide, x2, y_tabelas + 0.24, largura_col, altura_tabela, tabela_motivo["cabecalho"], linhas,
                larguras_relativas=[1.8, 1.2, 0.7, 1.2], tamanho_fonte=9.5)

    periodo = (dado.get("filtros") or {}).get("Período", "—")
    exportado = dado.get("exportado_em") or "—"
    _texto(
        slide, MARGEM_IN, 6.70, LARGURA_IN - 2 * MARGEM_IN, 0.35,
        f"Cartões: retrato do período {periodo}, exportado em {exportado}. Gráfico: histórico mensal completo "
        "do export (janela mais longa) — os dois totais não baterem entre si é esperado, não é inconsistência.",
        tamanho=9, cor=CINZA_TEXTO, italico=True,
    )
    return slide


def _slide_dashboard_externo_generico(prs: Presentation, mes_label: str, pagina: int, item: dict):
    """Slide genérico pra indicadores dinâmicos (18/08/2026, pedido do usuário:
    "adicione a opção de adicionar mais indicadores e adicionar automaticamente
    na construção do MBR") - um slide desses por item em d["dashboards_extras"],
    com extração deliberadamente conservadora (só tabelas + metadados de
    exportação, ver dashboards_externos_extrator.extrair_generico) - sem tentar
    adivinhar cartões de KPI num HTML de layout desconhecido."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    nome = item.get("nome_exibicao") or "Indicador"
    _cabecalho(slide, nome, mes_label, pagina,
               "Indicador adicionado pela equipe — retrato do arquivo enviado em Auditoria > Outros Dashboards")

    dado = item.get("dado")
    if not dado:
        _caixa_leitura(slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, "Sem conteúdo extraível",
                       "O arquivo enviado para este indicador não teve nenhuma tabela nem metadado de exportação "
                       "reconhecido — confirme se o .html enviado é o export autocontido correto em "
                       "Auditoria > Outros Dashboards.",
                       cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13)
        return slide

    tabelas = dado.get("tabelas") or []
    if not tabelas:
        _caixa_leitura(slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, "Sem tabelas encontradas",
                       "Este indicador foi reconhecido, mas o arquivo não teve nenhuma tabela — ele aparece com o "
                       "conteúdo completo assim que o export enviado contiver ao menos uma.",
                       cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13)
        return slide

    y_topo = 1.65
    if len(tabelas) == 1:
        tabela = tabelas[0]
        if tabela.get("titulo"):
            _texto(slide, MARGEM_IN, y_topo, LARGURA_IN - 2 * MARGEM_IN, 0.26, tabela["titulo"].upper(),
                   tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
        linhas = [[c[:30] if isinstance(c, str) else c for c in linha] for linha in tabela["linhas"][:8]]
        _tabela(slide, MARGEM_IN, y_topo + 0.30, LARGURA_IN - 2 * MARGEM_IN, 3.9, tabela["cabecalho"], linhas,
                tamanho_fonte=10.5)
        tabelas_exibidas = 1
    else:
        largura_col = (LARGURA_IN - 2 * MARGEM_IN - 0.35) / 2
        x2 = MARGEM_IN + largura_col + 0.35
        duas = tabelas[:2]
        # O título de cada tabela vem do heading/negrito mais próximo ANTES dela
        # no HTML original (ver extrair_generico) - em exports sem um heading
        # próprio por tabela, as duas acabam herdando o mesmo heading do topo
        # da página, o que rende duas colunas com o título idêntico lado a lado.
        # Desambigua nesse caso só na exibição (não nos dados extraídos).
        titulos_iguais = (
            len(duas) == 2 and duas[0].get("titulo") and duas[0]["titulo"] == duas[1]["titulo"]
        )
        for i, tabela in enumerate(duas):
            x = MARGEM_IN if i == 0 else x2
            titulo = tabela.get("titulo")
            if titulo:
                if titulos_iguais:
                    titulo = f"{titulo} — Tabela {i + 1}"
                _texto(slide, x, y_topo, largura_col, 0.26, titulo.upper(), tamanho=10.5,
                       negrito=True, cor=AZUL_INSTITUCIONAL)
            linhas = [[c[:22] if isinstance(c, str) else c for c in linha] for linha in tabela["linhas"][:6]]
            _tabela(slide, x, y_topo + 0.30, largura_col, 3.0, tabela["cabecalho"], linhas, tamanho_fonte=9.5)
        tabelas_exibidas = min(2, len(tabelas))

    exportado = dado.get("exportado_em") or item.get("enviado_em") or "—"
    periodo = (dado.get("filtros") or {}).get("Período")
    extras = len(tabelas) - tabelas_exibidas
    rodape = f"Exportado em {exportado}."
    if periodo:
        rodape += f" Período do arquivo: {periodo}."
    if extras > 0:
        rodape += f" +{extras} tabela(s) adicional(is) no arquivo original, não exibida(s) aqui."
    _texto(slide, MARGEM_IN, 6.55, LARGURA_IN - 2 * MARGEM_IN, 0.6, rodape, tamanho=9.5, cor=CINZA_TEXTO, italico=True)
    return slide


# Frentes que o Atlas cobre nativamente (18/08/2026, usado só no slide de
# Impacto do Atlas) - nomes curtos de propósito (um item por linha, sem
# quebrar em 2 linhas na lista de cobertura, ver _slide_impacto_atlas).
MODULOS_NATIVOS_MBR = [
    "Painel de Inventário — fechamento por almoxarifado",
    "Acurácia Ponderada por Valor (IAP)",
    "Mapeamento de Passivos e Baixas Operacionais",
    "Shelf Life — Farol de Validade",
    "Mapeamento de Risco de Obsolescência",
    "Controle de Movimentados — reconciliação diária",
]


def _slide_impacto_atlas(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Impacto do Atlas (18/08/2026, pedido do usuário: mostrar "o impacto do
    atlas no mapeamento completo dos processos e melhorias previstas com o uso
    contínuo da ferramenta") - fecha o relatório com números reais medidos (não
    projeção genérica): quantas frentes o Atlas cobre nativamente, quantos
    controles paralelos da equipe já foram integrados automaticamente ao MBR, e
    o ganho de acurácia medido desde a implantação do Controle de Movimentados."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Impacto do Atlas", mes_label, pagina,
               "Mapeamento completo dos processos monitorados e ganhos mensuráveis com o uso contínuo da ferramenta")

    dashboards_nativos_enviados = sum(
        1 for chave in ("fefo_externo", "testes_industriais_externo", "farol_shelf_externo",
                         "recuperacao_shelf_externo", "baixas_operacionais_externo",
                         "dispersao_ficha_tecnica_externo")
        if (d.get(chave) or {}).get("enviado")
    )
    dashboards_extras = d.get("dashboards_extras") or []
    total_indicadores_integrados = dashboards_nativos_enviados + len(dashboards_extras)

    resumo_mov = d["resumo_movimentados"]
    evolucao_mov = d["evolucao_movimentados"]
    primeiro_mov = evolucao_mov[0] if evolucao_mov else None
    delta_implantacao = None
    if primeiro_mov and primeiro_mov.get("pct_acuracia") is not None and resumo_mov.get("pct_acuracia") is not None:
        delta_implantacao = round(resumo_mov["pct_acuracia"] - primeiro_mov["pct_acuracia"], 2)

    valor_visibilidade = (
        (d["resumo_passivos"]["passivos"].get("valor") or 0)
        + (d["resumo_shelf_life"].get("valor_total") or 0)
        + (d["mapeamento_risco_obsolescencia"].get("valor_total_risco") or 0)
    )

    _linha_kpis(slide, 1.55, [
        {"valor": _fmt_num(len(MODULOS_NATIVOS_MBR)), "rotulo": "Frentes Nativas Mapeadas pelo Atlas"},
        {"valor": _fmt_num(total_indicadores_integrados), "rotulo": "Controles Paralelos Integrados ao MBR",
         "cor": COR_SUCESSO if total_indicadores_integrados else CINZA_TEXTO},
        {"valor": (f"+{_fmt_pct(delta_implantacao)}" if delta_implantacao >= 0 else _fmt_pct(delta_implantacao))
                  if delta_implantacao is not None else "—",
         "rotulo": "Ganho de Acurácia (Movimentados)",
         "cor": COR_SUCESSO if (delta_implantacao or 0) >= 0 else COR_ERRO},
        {"valor": _fmt_moeda(valor_visibilidade), "rotulo": "Valor sob Visibilidade Ativa", "cor": AZUL_INSTITUCIONAL},
    ], altura=0.85)

    largura_esquerda = 6.3
    _texto(slide, MARGEM_IN, 3.05, largura_esquerda, 0.28, "COBERTURA DE PROCESSOS HOJE", tamanho=11,
           negrito=True, cor=AZUL_INSTITUCIONAL)
    itens_cobertura = list(MODULOS_NATIVOS_MBR)
    itens_cobertura.append(
        f"{total_indicadores_integrados} controle(s) paralelo(s) da equipe já lido(s) automaticamente do arquivo exportado."
        if total_indicadores_integrados else
        "Nenhum controle paralelo enviado ainda — suba os arquivos em Auditoria > Outros Dashboards pra integrá-los aqui."
    )
    _lista_com_marcadores(slide, MARGEM_IN, 3.40, largura_esquerda, 3.0, itens_cobertura, tamanho=12)

    x_direita = MARGEM_IN + largura_esquerda + 0.35
    largura_direita = LARGURA_IN - MARGEM_IN - x_direita
    if delta_implantacao is not None and primeiro_mov:
        texto_melhoria = (
            f"Desde a implantação do Controle de Movimentados, em {_nome_mes(primeiro_mov['mes'])}, a acurácia da "
            f"reconciliação diária avançou {_fmt_pct(delta_implantacao)} — mesma lógica de leitura contínua que hoje "
            "também sustenta Shelf Life, Mapeamento de Risco de Obsolescência e os controles paralelos integrados "
            "neste relatório. Quanto mais meses de uso contínuo, mais preciso o histórico e mais cedo cada "
            "divergência aparece — o próximo ganho esperado é reduzir o tempo entre a divergência ocorrer e ela "
            "ser tratada, hoje limitado pela cadência de conferência manual."
        )
    else:
        texto_melhoria = (
            "Ainda não há histórico suficiente pra medir o ganho de acurácia desde a implantação — a leitura "
            "aparece a partir do segundo mês de dados do Controle de Movimentados."
        )
    _caixa_leitura(slide, x_direita, 3.05, largura_direita, 3.35, "Melhorias Previstas com o Uso Contínuo",
                   texto_melhoria, cor_fundo=OFF_WHITE, tamanho_texto=12)

    _texto(
        slide, MARGEM_IN, 6.55, LARGURA_IN - 2 * MARGEM_IN, 0.5,
        "Cobertura e integrações medidas neste relatório — não é uma projeção genérica: cada número acima vem dos "
        "mesmos dados que alimentam as telas do Atlas no dia a dia.",
        tamanho=9.5, cor=CINZA_TEXTO, italico=True,
    )
    return slide


def _slide_diario_bordo(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Constância e Disciplina — Diário de Bordo (21/08/2026, pedido do
    usuário: ver docstring de _DIARIO_BORDO_POR_MES acima pro contexto
    completo e a ressalva sobre esse indicador vir de coleta manual, não de
    consulta ao banco do Atlas). Indicador PARALELO ao resto do MBR: mede
    disciplina operacional (manter o diário de bordo em dia), não estoque -
    entra na seção Atlas como um segundo ângulo de "impacto do Atlas", já
    que a mesma constância sustenta a qualidade dos dados usados em todo o
    resto deste relatório."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Constância e Disciplina — Diário de Bordo", mes_label, pagina,
               "Indicador paralelo: cumprimento da Rotina Master em dias úteis do mês (fins de semana não têm rotina devida)")

    dado = d.get("diario_bordo") or {"tem_dados": False}
    if not dado.get("tem_dados"):
        _caixa_leitura(
            slide, MARGEM_IN, 1.6, LARGURA_IN - 2 * MARGEM_IN, 1.4, "Ainda sem coleta pra este mês",
            "Este indicador vem de uma coleta manual no Dashboard de Performance da Rotina Master "
            "(rotinabusiness.lovable.app), filtrado pelo mês de fechamento — ainda não foi coletado pra este mês "
            "específico. Repita a coleta (filtro De/Até = 1º ao último dia do mês) pra atualizar este slide.",
            cor_fundo=OFF_WHITE, cor_rotulo=COR_ATENCAO, tamanho_texto=13,
        )
        return slide

    # 02/09/2026: a coleta manual completa (_DIARIO_BORDO_POR_MES) tem
    # sequência/lapsos por dia útil e quebra semanal; a extração automática
    # do export HTML (dashboards_externos_extrator.extrair_diario_bordo) só
    # traz os KPIs de topo (não reconstrói o gráfico diário - ver docstring
    # da função). Detecta qual das duas gerou este `dado` pela presença do
    # campo exclusivo da coleta manual, e mostra uma versão mais enxuta do
    # slide quando é extração automática, em vez de inventar os campos que
    # faltam.
    completo = "maior_sequencia_dias_uteis_100" in dado
    lapsos = dado.get("lapsos_dias_uteis") or []

    if completo:
        kpis = [
            {"valor": _fmt_pct(dado["cumprimento_geral_pct"]), "rotulo": "Cumprimento Geral do Mês"},
            {"valor": f"{_fmt_num(dado['rotinas_cumpridas'])}/{_fmt_num(dado['rotinas_devidas'])}",
             "rotulo": "Rotinas Cumpridas / Devidas"},
            {"valor": _fmt_num(dado["maior_sequencia_dias_uteis_100"]), "rotulo": "Maior Sequência a 100%",
             "cor": COR_SUCESSO},
            {"valor": _fmt_num(len(lapsos)), "rotulo": "Lapsos em Dias Úteis",
             "cor": COR_ATENCAO if lapsos else COR_SUCESSO},
        ]
    else:
        pct_atraso = dado.get("pct_em_atraso")
        kpis = [
            {"valor": _fmt_pct(dado["cumprimento_geral_pct"]), "rotulo": "Cumprimento Geral do Mês"},
            {"valor": f"{_fmt_num(dado['rotinas_cumpridas'])}/{_fmt_num(dado['rotinas_devidas'])}",
             "rotulo": "Rotinas Cumpridas / Devidas"},
            {"valor": _fmt_pct(dado.get("pct_no_prazo")), "rotulo": "Conclusões no Prazo",
             "cor": COR_SUCESSO},
            {"valor": _fmt_pct(pct_atraso), "rotulo": "Conclusões em Atraso",
             "cor": COR_ATENCAO if (pct_atraso or 0) > 0 else COR_SUCESSO},
        ]
    _linha_kpis(slide, 1.55, kpis, altura=0.85)

    largura_esquerda = 7.1
    semanas = dado.get("semanas") or []
    _texto(slide, MARGEM_IN, 3.05, largura_esquerda, 0.28, "EVOLUÇÃO SEMANAL DE CUMPRIMENTO (DIAS ÚTEIS)",
           tamanho=11, negrito=True, cor=AZUL_INSTITUCIONAL)
    if semanas:
        categorias = [s["rotulo"] for s in semanas]
        valores = [s["cumprimento_pct"] for s in semanas]
        _grafico_categoria(slide, MARGEM_IN, 3.40, largura_esquerda, 3.0, categorias, "Cumprimento (%)", valores,
                            tipo=XL_CHART_TYPE.LINE_MARKERS, cor_serie=AZUL_INSTITUCIONAL)
    elif completo:
        _caixa_leitura(slide, MARGEM_IN, 3.40, largura_esquerda, 3.0, "Evolução semanal",
                       "Sem quebra semanal coletada pra este mês.")
    else:
        _caixa_leitura(
            slide, MARGEM_IN, 3.40, largura_esquerda, 3.0, "Evolução semanal",
            "A extração automática deste indicador traz só os KPIs de topo do export — quebra semanal e "
            "sequência/lapsos por dia útil ainda dependem de coleta manual complementar na Rotina Master.",
        )

    x_direita = MARGEM_IN + largura_esquerda + 0.35
    largura_direita = LARGURA_IN - MARGEM_IN - x_direita
    if completo and lapsos:
        texto_leitura = (
            f"Média em dias úteis: {_fmt_pct(dado['media_dias_uteis_pct'])} — {len(lapsos)} lapso(s) pontual(is) "
            f"({', '.join(lapsos)}), sempre recuperados no dia útil seguinte, sem arrastar pra outros dias. Maior "
            f"sequência sem falha: {_fmt_num(dado['maior_sequencia_dias_uteis_100'])} dias úteis consecutivos a "
            "100%. Fins de semana ficam fora dessa leitura — não têm rotina devida no app."
        )
    elif completo:
        texto_leitura = (
            f"Média em dias úteis: {_fmt_pct(dado['media_dias_uteis_pct'])}, sem nenhum lapso no mês — maior "
            f"sequência de {_fmt_num(dado['maior_sequencia_dias_uteis_100'])} dias úteis consecutivos a 100%. "
            "Fins de semana ficam fora dessa leitura — não têm rotina devida no app."
        )
    else:
        texto_leitura = (
            f"Cumprimento geral do mês: {_fmt_pct(dado['cumprimento_geral_pct'])} "
            f"({_fmt_num(dado['rotinas_cumpridas'])} de {_fmt_num(dado['rotinas_devidas'])} rotinas), sendo "
            f"{_fmt_pct(dado.get('pct_no_prazo'))} concluídas no prazo e {_fmt_pct(dado.get('pct_em_atraso'))} em "
            "atraso. Extraído automaticamente do Dashboard de Performance da Rotina Master — ainda não traz "
            "sequência/lapsos por dia útil nem quebra semanal; pra esse detalhe é necessária uma coleta manual "
            "complementar."
        )
    _caixa_leitura(slide, x_direita, 3.05, largura_direita, 3.35, "Leitura Executiva",
                   texto_leitura, cor_fundo=OFF_WHITE, tamanho_texto=12)

    if dado.get("coletado_em"):
        texto_fonte = (
            f"Coletado manualmente em {dado.get('coletado_em', '—')} navegando na Rotina Master "
            "(rotinabusiness.lovable.app), filtrado pelo mês de fechamento — app separado do Atlas, sem integração "
            "automática ainda."
        )
    else:
        texto_fonte = (
            f"Extraído automaticamente do Dashboard de Performance da Rotina Master (rotinabusiness.lovable.app), "
            f"enviado em Auditoria > Outros Dashboards em {dado.get('enviado_em') or '—'}"
            + (f" — período do export: {dado['periodo']}." if dado.get("periodo") else ".")
        )
    _texto(
        slide, MARGEM_IN, 6.55, LARGURA_IN - 2 * MARGEM_IN, 0.5,
        texto_fonte,
        tamanho=9.5, cor=CINZA_TEXTO, italico=True,
    )
    return slide


def _slide_atlas_stock_savvy_visao(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Atlas + Stock Savvy (20/08/2026, pedido do usuário: "Análise também o
    Controle desenvolvido em paralelo no lovable, o aplicativo 'Stock
    Savvy'... venda a ideia de ambos os projetos na construção de um
    controle robusto e eficiente"). Conteúdo baseado em navegação real pelo
    Stock Savvy (stockswift-sync-75.lovable.app, workspace Lovable "Stock
    Savvy" - ver claude/sincronizacao-lovable-baixas.md pro histórico da
    integração já existente entre os dois sistemas), não só na descrição
    do usuário - confirmei estrutura e dado real de cada módulo citado
    (Produção, Shelf Life, Gestão) direto nas telas.

    Posicionamento: os dois sistemas não competem - operam em camadas
    diferentes da mesma operação. Stock Savvy é onde a ação acontece (
    solicitar, escanear, aprovar com assinatura dupla, registrar uma ação de
    lote com custo e recuperação); Atlas é onde as frentes se cruzam num
    relatório executivo único, com histórico mensal e plano de ação — a
    maior parte do que o Atlas já consome de Stock Savvy hoje (Farol de
    Shelf, Dashboard Shelf Life/Recuperação, Dashboard de Baixas, Dispersão
    de Lote) entra via os mesmos exports HTML que o Stock Savvy já gera."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Atlas + Stock Savvy", mes_label, pagina,
               "Dois sistemas, duas camadas de controle — por que manter os dois é o que dá robustez")

    largura_col = (LARGURA_IN - 2 * MARGEM_IN - 0.35) / 2
    _caixa_leitura(
        slide, MARGEM_IN, 1.55, largura_col, 3.15, "Stock Savvy — Camada Operacional",
        "Onde a ação acontece, dia a dia, no chão de fábrica e na loja: solicitar uma baixa escaneando o QR do lote, "
        "aprovar com assinatura dupla (Diretor de Operações + Coordenador Financeiro), registrar uma ação de "
        "recuperação de lote (desconto, doação, anúncio) e acompanhar seu custo e retorno em tempo real. É o sistema "
        "de registro e rastreabilidade — cada baixa, cada ação de lote, cada auditoria de ficha técnica nasce lá.",
        cor_fundo=OFF_WHITE, cor_rotulo=AZUL_INSTITUCIONAL, tamanho_texto=12.5,
    )
    x_direita = MARGEM_IN + largura_col + 0.35
    _caixa_leitura(
        slide, x_direita, 1.55, largura_col, 3.15, "Atlas — Camada de Inteligência Executiva",
        "Onde as frentes se cruzam: Inventário, Movimentados, FEFO, Shelf Life, Passivos e os controles paralelos da "
        "equipe (incluindo os exports do próprio Stock Savvy) chegam num único relatório mensal, com histórico "
        "real, evolução mês a mês e plano de ação por frente e por almoxarifado — não um retrato isolado de um "
        "módulo, e sim a leitura de como a operação inteira está andando.",
        cor_fundo=OFF_WHITE, cor_rotulo=VERDE_AMAZONIA, tamanho_texto=12.5,
    )

    texto_por_que = (
        "Um sem o outro deixa uma lacuna real: só Stock Savvy dá execução granular e rastreável, mas sem visão "
        "cruzada entre frentes nem histórico executivo; só Atlas dá a leitura consolidada, mas precisa de onde "
        "vêm os dados operacionais de origem. Juntos, formam um controle de ponta a ponta — o dado nasce com "
        "rastreabilidade e governança no Stock Savvy (quem pediu, quem aprovou, quando), e chega consolidado, "
        "comparável mês a mês e com plano de ação no Atlas. Essa é a base de controle mais robusta possível hoje "
        "para o estoque da Mágio: operação rastreada + inteligência executiva, sem depender de planilha solta "
        "entre uma ponta e outra."
    )
    # Altura fixa (1.85) cortava esse texto pela rede de segurança de
    # _caber_no_espaco - calcula a altura real necessária (com folga de
    # 0.15 acima do mínimo teórico) em vez de chutar um valor fixo (mesmo
    # padrão já usado em _slide_resumo_executivo, ver _altura_necessaria_
    # caixa_leitura).
    altura_por_que = max(1.85, _altura_necessaria_caixa_leitura(
        texto_por_que, LARGURA_IN - 2 * MARGEM_IN, 12.5) + 0.15)
    _caixa_leitura(
        slide, MARGEM_IN, 4.95, LARGURA_IN - 2 * MARGEM_IN, altura_por_que, "Por que manter os dois",
        texto_por_que,
        cor_fundo=OFF_WHITE, cor_rotulo=COR_SUCESSO, tamanho_texto=12.5,
    )
    return slide


def _slide_atlas_stock_savvy_modulos(prs: Presentation, mes_label: str, pagina: int, d: dict):
    """Módulos recentes do Stock Savvy citados pelo usuário (Produção,
    Shelf Life, Gestão) - conteúdo confirmado navegando nas telas reais
    (20/08/2026), não apenas na descrição recebida."""
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Stock Savvy — Módulos Recentes", mes_label, pagina,
               "O que cada módulo executa no Stock Savvy e como já se conecta ao Atlas hoje")

    linhas = [
        [
            ("Produção", AZUL_INSTITUCIONAL, True),
            "Dispersão de Lote (Ficha Técnica x consumo real por OP, com Ações Corretivas rastreadas por "
            "responsável/status) e Auditoria de Ficha Técnica (cobertura de cadastro de BOM por produto).",
            "Alimenta o slide de Dispersão de Ficha Técnica do MBR "
            "(taxa de furo, materiais crônicos, impacto líquido).",
        ],
        [
            ("Shelf Life", AZUL_INSTITUCIONAL, True),
            "Mapeamento de Risco por almoxarifado, Ações de Lote (desconto, doação, anúncio — com custo e valor "
            "recuperado por ação) e Farol de Shelf, consolidados no Dashboard Shelf Life com ROI operacional e "
            "Saving recuperado.",
            "Os mesmos exports HTML (Farol de Shelf, Dashboard Shelf Life) já alimentam os slides de Farol de "
            "Shelf-Life e Recuperação de Shelf do MBR.",
        ],
        [
            ("Gestão", AZUL_INSTITUCIONAL, True),
            "Baixas Operacionais com fluxo completo de solicitação (scanner QR/EAN), fila de aprovação com "
            "assinatura dupla e histórico, além do Dashboard de Baixas por motivo/setor/grupo.",
            "Já sincronizado ao vivo com o módulo de Baixas Operacionais do Atlas (Mapeamento de Passivos) via "
            "rota própria — integração automática, sem exportação manual de arquivo.",
        ],
    ]
    _tabela(slide, MARGEM_IN, 1.65, LARGURA_IN - 2 * MARGEM_IN, 4.7,
            ["Módulo", "Execução no Stock Savvy", "Conexão com o Atlas"], linhas,
            larguras_relativas=[1.1, 3.6, 3.0], tamanho_fonte=12)

    _texto(
        slide, MARGEM_IN, 6.55, LARGURA_IN - 2 * MARGEM_IN, 0.5,
        "Levantamento feito navegando diretamente nas telas do Stock Savvy (workspace Lovable \"Stock Savvy\") "
        "em 20/08/2026 — não é uma descrição de segunda mão.",
        tamanho=9.5, cor=CINZA_TEXTO, italico=True,
    )
    return slide


def _slide_proximos_passos(prs: Presentation, mes_label: str, pagina: int, d: dict):
    slide = _slide_em_branco(prs)
    _fundo(slide, BRANCO)
    _cabecalho(slide, "Próximos Passos", mes_label, pagina, "Ações priorizadas a partir dos pontos de atenção identificados acima")

    scorecard = _montar_scorecard(d)
    acoes = [item for item in scorecard if item["status_label"] != "Em avanço"]

    y_inicio = 1.65
    y_limite = ALTURA_IN - 1.55  # deixa espaço reservado pra caixa "Mensagem" no rodapé
    if not acoes:
        _caixa_leitura(slide, MARGEM_IN, y_inicio, LARGURA_IN - 2 * MARGEM_IN, 1.2, "Sem ação crítica pendente",
                       "Todas as frentes monitoradas estão dentro da faixa de controle neste recorte — manter a cadência atual de fechamento, mapeamento e monitoramento de validade.",
                       cor_fundo=OFF_WHITE)
    else:
        # Altura/gap calculados a partir de quantas ações existem (no máximo 5,
        # uma por frente do scorecard) pra nunca ultrapassar o espaço disponível
        # acima da caixa "Mensagem" - com poucas ações os blocos ficam maiores
        # (mais folgados), com mais ações eles encolhem proporcionalmente.
        n = len(acoes)
        gap = 0.16
        altura_bloco = min(1.05, ((y_limite - y_inicio) - gap * (n - 1)) / n)
        tamanho_indice = 26 if altura_bloco >= 0.95 else max(16, int(26 * altura_bloco / 1.05))
        tamanho_frente = 14 if altura_bloco >= 0.95 else 12
        tamanho_passo = 12 if altura_bloco >= 0.95 else 10.5
        y = y_inicio
        for i, item in enumerate(acoes, start=1):
            _retangulo(slide, MARGEM_IN, y, LARGURA_IN - 2 * MARGEM_IN, altura_bloco, cor_fill=BRANCO, cor_borda=CINZA_CLARO, raio=0.10)
            _texto(slide, MARGEM_IN + 0.22, y, 0.5, altura_bloco, str(i), tamanho=tamanho_indice, negrito=True,
                   cor=AZUL_CLARO, fonte=FONTE_TITULO, ancora_meio=True)
            _texto(slide, MARGEM_IN + 0.85, y + altura_bloco * 0.12, 6.0, altura_bloco * 0.4, item["frente"],
                   tamanho=tamanho_frente, negrito=True, cor=AZUL_INSTITUCIONAL)
            _texto(slide, MARGEM_IN + 0.85, y + altura_bloco * 0.48, LARGURA_IN - 2 * MARGEM_IN - 3.0, altura_bloco * 0.48,
                   item["proximo_passo"], tamanho=tamanho_passo, cor=CINZA_TEXTO)
            _badge_status(slide, LARGURA_IN - MARGEM_IN - 1.6, y + (altura_bloco - 0.38) / 2, 1.35, 0.38,
                          item["status_label"], item["status_cor"])
            y += altura_bloco + gap

    _caixa_leitura(
        slide, MARGEM_IN, ALTURA_IN - 1.35, LARGURA_IN - 2 * MARGEM_IN, 0.9, "Mensagem",
        "Este relatório é gerado automaticamente pelo Atlas a partir dos mesmos dados que alimentam as telas do dia a dia — "
        "qualquer ajuste em um fechamento, baixa ou lote se reflete aqui no próximo mês, sem retrabalho manual.",
        cor_fundo=OFF_WHITE, tamanho_texto=11,
    )
    return slide


# ---------------------------------------------------------------------------
# Ponto de entrada público
# ---------------------------------------------------------------------------
def montar_pptx_mbr(db: Session, usuario: models.Usuario, mes: str) -> bytes:
    """Gera o MBR completo (.pptx) para o mês informado ("AAAA-MM"),
    lendo os dados diretamente das funções de negócio do Atlas (sem HTTP,
    sem navegador) e aplicando a identidade visual da Mágio Chocolates.

    Estrutura em 4 seções temáticas, cada uma com capa própria (18/08/2026,
    pedido do usuário: "traga uma visão detalhada modulando os grupos de
    relatório" — respostas às perguntas de esclarecimento: "Reordenar + capa
    de seção"; reduzida de 7 para 5 seções em 22/08/2026 (FEFO e Testes
    Industriais deixaram de ter seção/capa própria - virou slide de
    detalhe dentro de "Mapeamento de Riscos e Passivos"); e de 5 para 4
    seções ainda em 22/08/2026, pedido não-negociável da usuária de fundir
    a antiga seção "Outros" (dashboards externos) com "Mapeamento de Riscos
    e Passivos" (nativa), substituindo os indicadores nativos que já têm
    dashboard externo aprovado cobrindo o mesmo assunto — ver comentário de
    decisão dentro da seção 3, abaixo). A contagem de página é dinâmica
    (closure `_pag`) porque o número de slides varia com a quantidade de
    indicadores dinâmicos cadastrados em Outros Dashboards (ver
    dados["dashboards_extras"]).

    22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o
    modelo final"): "Painel de Inventário — Detalhamento Financeiro",
    "Mapeamento de Risco — Obsolescência" e "Scorecard de Mapeamento de
    Riscos" (que fazia a síntese de FEFO/Testes Industriais citada no
    parágrafo acima) saíram do relatório - ver os comentários de decisão
    antes de _slide_acuracia_ponderada_iap e antes de
    _slide_controle_movimentados."""
    dados = _coletar_dados_mbr(db, usuario, mes)
    mes_label = _nome_mes(mes)

    prs = _nova_apresentacao()
    _slide_capa(prs, mes_label)

    pagina = [1]  # a capa já é a página 1

    def _pag():
        pagina[0] += 1
        return pagina[0]

    def _secao(numero, titulo, descricao, itens):
        _slide_abertura_secao(prs, mes_label, _pag(), numero, titulo, descricao, itens)

    _secao(1, "Resumo Executivo Geral",
           "Leitura consolidada do mês e status por frente monitorada.",
           ["Resumo Executivo", "Scorecard do Mês"])
    _slide_resumo_executivo(prs, mes_label, _pag(), dados)
    _slide_scorecard(prs, mes_label, _pag(), dados)

    # 22/08/2026 (mockups aprovados v4/v5/v8): Painel de Inventário e Acurácia
    # Ponderada ganharam slide de detalhamento financeiro/tendência cada, e
    # Acurácia Ponderada foi de 1 slide combinado (IAQ+IAP) pra 2 dedicados
    # (IAP, IAQ) + o novo "Detalhamento por Faixa de Magnitude" - a seção foi
    # de 5 pra 8 slides de indicador.
    #
    # 22/08/2026, pedido não-negociável da usuária ("dar ênfase ao inventário
    # ponderado IAP e IAQ. Deixe a visão item a item, porém não dê destaque
    # à mesma, apenas para conhecimento"): IAP/IAQ (ponderados) passam a
    # abrir a seção; a leitura item a item (Inventário Item a Item) desce
    # pro final da seção, antes de Movimentados/Scorecard — o CONTEÚDO de
    # cada slide não muda (ver docstring de _slide_painel_inventario), só a
    # ORDEM.
    #
    # 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o
    # modelo final"): o slide companheiro "Inventário — SKUs Recorrentes e
    # Cobertura" (_slide_painel_inventario_detalhe) foi removido - ver
    # comentário de decisão antes de _slide_acuracia_ponderada_iap.
    _secao(2, "Inventários e Movimentados",
           "Acurácia de fechamento, ponderação por valor e reconciliação diária sistema x físico.",
           ["Acurácia Ponderada (IAP)", "Acurácia Ponderada (IAQ)",
            "Acurácia Ponderada — Concentração de Risco", "Acurácia Ponderada — Detalhamento por Faixa",
            "Inventário Item a Item",
            "Controle de Movimentados", "Scorecard de Inventário por Almoxarifado"])
    _slide_acuracia_ponderada_iap(prs, mes_label, _pag(), dados)
    _slide_acuracia_ponderada_iaq(prs, mes_label, _pag(), dados)
    _slide_acuracia_ponderada_detalhe(prs, mes_label, _pag(), dados)
    _slide_acuracia_ponderada_faixas(prs, mes_label, _pag(), dados)
    _slide_painel_inventario(prs, mes_label, _pag(), dados)
    _slide_controle_movimentados(prs, mes_label, _pag(), dados)
    _slide_scorecard_inventario_almoxarifado(prs, mes_label, _pag(), dados)

    # 22/08/2026, pedido não-negociável da usuária: Seção 3 (nativa,
    # "Mapeamento de Riscos e Passivos") e a antiga Seção 4 ("Outros",
    # dashboards externos) foram FUNDIDAS numa seção só - a usuária pediu
    # explicitamente pra substituir os indicadores nativos que têm um
    # dashboard externo aprovado cobrindo o mesmo assunto, em vez de manter
    # as duas fontes lado a lado. Mapeamento de Passivos/Passivos-Evolução e
    # Shelf Life (nativos) foram REMOVIDOS (ver comentário de decisão antes
    # de _slide_controle_movimentados, mais acima no arquivo) - Dashboard
    # Baixas Operacionais e Farol de Shelf-Life (externos, modelo já
    # aprovado) passam a ser a única fonte pra esses 2 assuntos. Testes
    # Industriais e FEFO ficam (sem equivalente externo aprovado);
    # Recuperação de Shelf e Dispersão de Ficha Técnica entram como
    # conteúdo novo (não substituem nada).
    #
    # FEFO e Testes Industriais entraram nesta seção em 22/08/2026 (pedido
    # anterior do usuário) - tinham cada um sua própria seção de 2 slides
    # (capa + indicador); as duas capas de seção ("slides de apresentação")
    # foram removidas e os dois slides de indicador viraram detalhe desta
    # seção.
    #
    # 22/08/2026, pedido da usuária ("fiz uma versão manual [...] esse é o
    # modelo final"): "Mapeamento de Risco — Obsolescência" e "Scorecard de
    # Mapeamento de Riscos" (capítulo-síntese que fechava a seção) foram
    # REMOVIDOS - ver comentário de decisão antes de _slide_controle_
    # movimentados, mais acima no arquivo.
    #
    # Cap defensivo na lista "NESTA SEÇÃO" da capa (não no relatório em si -
    # todo indicador dinâmico enviado ainda gera seu slide completo abaixo)
    # pra não estourar a capa de seção se a equipe cadastrar muitos
    # indicadores.
    #
    # 21/08/2026, Fase 3 (pedido da usuária, ver nota de decisão em
    # _slide_farol_shelf_externo e _slide_recuperacao_shelf_externo): os
    # slides companheiros "Farol de Shelf-Life — Risco por Almoxarifado" e
    # "Recuperação de Shelf — Evolução Mensal" foram removidos (conteúdo
    # fundido de volta nos slides principais) - saíram também desta lista.
    #
    # 21/08/2026, mesma Fase 3, pedido seguinte da usuária ("replique a
    # mesma lógica para o KPI. Slide 14 e 15" - Dashboard Baixas
    # Operacionais): "Dashboard Baixas Operacionais — Evolução Mensal"
    # também foi removido (mesma fusão) - saiu desta lista também.
    nomes_extras = [item["nome_exibicao"] for item in dados["dashboards_extras"]]
    limite_itens_capa = 4
    itens_riscos_passivos = [
        "Dashboard Baixas Operacionais",
        "Farol de Shelf-Life", "Recuperação de Shelf",
        "Dispersão de Ficha Técnica", "Testes Industriais", "FEFO",
    ]
    if len(nomes_extras) > limite_itens_capa:
        itens_riscos_passivos += nomes_extras[:limite_itens_capa]
        itens_riscos_passivos.append(f"+ {len(nomes_extras) - limite_itens_capa} indicador(es) adicional(is)")
    else:
        itens_riscos_passivos += nomes_extras
    _secao(3, "Mapeamento de Riscos e Passivos",
           "Passivos e baixas (fonte oficial: dashboards externos aprovados), validade de lotes, "
           "recuperação de shelf, dispersão de ficha técnica, FEFO e Testes Industriais.",
           itens_riscos_passivos)
    _slide_baixas_operacionais_externo(prs, mes_label, _pag(), dados)
    _slide_farol_shelf_externo(prs, mes_label, _pag(), dados)
    _slide_recuperacao_shelf_externo(prs, mes_label, _pag(), dados)
    _slide_dispersao_ficha_tecnica(prs, mes_label, _pag(), dados)
    _slide_testes_industriais(prs, mes_label, _pag(), dados)
    _slide_fefo(prs, mes_label, _pag(), dados)
    for item in dados["dashboards_extras"]:
        _slide_dashboard_externo_generico(prs, mes_label, _pag(), item)

    _secao(4, "Atlas",
           "Cobertura de processos hoje e melhorias esperadas com o uso contínuo da ferramenta.",
           ["Impacto do Atlas", "Constância e Disciplina — Diário de Bordo", "Atlas + Stock Savvy",
            "Módulos Recentes do Stock Savvy", "Próximos Passos"])
    _slide_impacto_atlas(prs, mes_label, _pag(), dados)
    _slide_diario_bordo(prs, mes_label, _pag(), dados)
    _slide_atlas_stock_savvy_visao(prs, mes_label, _pag(), dados)
    _slide_atlas_stock_savvy_modulos(prs, mes_label, _pag(), dados)
    _slide_proximos_passos(prs, mes_label, _pag(), dados)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
